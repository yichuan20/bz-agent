#!/usr/bin/env python3
from __future__ import annotations

"""
read-pptx.py — inspect a PPTX file and report its structure.

Outputs JSON describing slide layouts (name + placeholder types) and
existing slide content so the agent can build an outline that correctly
maps to the template's layouts before calling create-pptx.py.

Usage:
  python3 read-pptx.py --file /path/to/deck.pptx
  python3 read-pptx.py --file deck.pptx --slides          # include existing slide content
  python3 read-pptx.py --file deck.pptx --layouts-only    # layouts only, no slide content

Output JSON schema:
  {
    "path": "/abs/path/to/deck.pptx",
    "slide_width_inches": 13.33,
    "slide_height_inches": 7.5,
    "slide_count": 8,
    "layouts": [
      {
        "index": 0,
        "name": "Title Slide",
        "placeholders": [
          {"idx": 0, "type": "CENTER_TITLE", "name": "Title 1"},
          {"idx": 1, "type": "SUBTITLE",     "name": "Subtitle 2"}
        ],
        "suggested_use": "title"
      },
      ...
    ],
    "slides": [           # only when --slides flag given
      {
        "index": 0,
        "layout": "Title Slide",
        "title": "Slide Title Text",
        "body": "Body text or bullets joined by newline"
      },
      ...
    ]
  }
"""

import argparse
import json
import sys
from pathlib import Path

# Maps pptx placeholder type enum values to readable strings
_PP_PLACEHOLDER = {
    1: "BODY",
    2: "CENTER_TITLE",
    3: "DATE",
    4: "FOOTER",
    5: "SLIDE_NUMBER",
    6: "SUBTITLE",
    7: "TITLE",
    8: "BITMAP",
    9: "CHART",
    10: "CLIP_ART",
    11: "FLOWCHART",
    12: "MEDIA_CLIP",
    13: "OBJECT",
    14: "ORG_CHART",
    15: "PICTURE",
    16: "TABLE",
    17: "VERTICAL_BODY",
    18: "VERTICAL_TITLE",
}

# Heuristic mapping from layout name keywords to suggested use
_LAYOUT_USE_HINTS = [
    (["title slide", "cover"], "title"),
    (["title and content", "title, content"], "content"),
    (["section header", "section"], "section"),
    (["two content", "two column", "2 col"], "two-column"),
    (["comparison"], "comparison"),
    (["blank"], "blank"),
    (["picture", "image", "photo"], "image"),
    (["title only"], "title-only"),
    (["caption"], "caption"),
]


def _guess_use(layout_name: str) -> str:
    lower = layout_name.lower()
    for keywords, use in _LAYOUT_USE_HINTS:
        if any(k in lower for k in keywords):
            return use
    return "content"  # safe default


def _placeholder_type_name(ph) -> str:
    try:
        return _PP_PLACEHOLDER.get(ph.placeholder_format.type, str(ph.placeholder_format.type))
    except Exception:
        return "UNKNOWN"


def read_layouts(prs) -> list[dict]:
    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for ph in layout.placeholders:
            placeholders.append(
                {
                    "idx": ph.placeholder_format.idx,
                    "type": _placeholder_type_name(ph),
                    "name": ph.name,
                }
            )
        layouts.append(
            {
                "index": i,
                "name": layout.name,
                "placeholders": placeholders,
                "suggested_use": _guess_use(layout.name),
            }
        )
    return layouts


def read_slides(prs) -> list[dict]:
    from pptx.util import Inches

    slides = []
    for i, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name if slide.slide_layout else ""
        title_text = ""
        body_parts: list[str] = []

        for ph in slide.placeholders:
            try:
                text = ph.text_frame.text.strip() if ph.has_text_frame else ""
            except Exception:
                text = ""
            if not text:
                continue
            idx = ph.placeholder_format.idx
            ph_type = _placeholder_type_name(ph)
            if idx == 0 or ph_type in ("TITLE", "CENTER_TITLE"):
                title_text = text
            else:
                body_parts.append(text)

        slides.append(
            {
                "index": i,
                "layout": layout_name,
                "title": title_text,
                "body": "\n".join(body_parts),
            }
        )
    return slides


def main() -> None:
    parser = argparse.ArgumentParser(description="Inspect a PPTX file and report its structure as JSON.")
    parser.add_argument("--file", required=True, help="Path to the .pptx file.")
    parser.add_argument("--slides", action="store_true", help="Include existing slide content in output.")
    parser.add_argument("--layouts-only", action="store_true", help="Output layouts only (no slide content).")
    args = parser.parse_args()

    path = Path(args.file)
    if not path.exists():
        print(json.dumps({"error": f"file not found: {path}"}))
        sys.exit(1)
    if path.suffix.lower() not in (".pptx",):
        print(json.dumps({"error": "only .pptx files are supported"}))
        sys.exit(1)

    try:
        from pptx import Presentation
        from pptx.util import Inches, Emu

        prs = Presentation(str(path))
    except Exception as exc:
        print(json.dumps({"error": f"could not open presentation: {exc}"}))
        sys.exit(1)

    EMU_PER_INCH = 914400
    result: dict = {
        "path": str(path.resolve()),
        "slide_width_inches": round(prs.slide_width / EMU_PER_INCH, 2),
        "slide_height_inches": round(prs.slide_height / EMU_PER_INCH, 2),
        "slide_count": len(prs.slides),
        "layouts": read_layouts(prs),
    }

    if args.slides and not args.layouts_only:
        result["slides"] = read_slides(prs)

    print(json.dumps(result, indent=2))


if __name__ == "__main__":
    main()
