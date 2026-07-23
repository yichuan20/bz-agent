#!/usr/bin/env python3
from __future__ import annotations

"""
create-pptx-sidecar.py — build a new PPTX by cloning slides from a template
sidecar JSON and applying content text updates.

Workflow:
  1. Agent calls read-pptx-sidecar.py to learn which template slides exist
     and which box ids hold editable text.
  2. Agent writes a content-spec JSON (see schema below).
  3. Agent calls this script — which clones the chosen template slides,
     replaces text in the specified boxes (preserving all fonts/colors/layout),
     writes a new sidecar JSON, and exports a .pptx.

Usage:
  python3 create-pptx-sidecar.py \\
      --template-sidecar /path/to/.template.pptx.json \\
      --spec-file /tmp/deck_spec.json \\
      --out /tmp/new_deck.pptx

Content-spec JSON schema:
  {
    "slides": [
      {
        "template_slide_index": 0,          // which template slide to clone
        "text_updates": {                   // box_id → new text (preserves styling)
          "abc123": "Boltzbit",
          "def456": "The AI platform for knowledge workers"
        }
      },
      {
        "template_slide_index": 2,
        "text_updates": {
          "ghi789": "The Problem",
          "jkl012": "80% of enterprise data is unstructured\\nKnowledge workers lose 2.5 hrs/day to manual tasks\\nExisting AI tools add chat on top of chaos"
        }
      }
    ]
  }

Rules for text_updates:
  - Only list box ids you want to change.  Unmentioned boxes (shapes, logos,
    decorative dots, etc.) are cloned unchanged.
  - For multi-line content use \\n — each line becomes a separate paragraph
    with the same styling as the original run.
  - To blank a box, pass "" as the value.

Output:
  Prints JSON: { "ok": true, "path": "/abs/path.pptx", "slides": N,
                 "sidecar": "/abs/path/.new_deck.pptx.json" }
"""

import argparse
import base64 as _b64
import copy
import io
import json
import os
import shutil
import sys
import tempfile
from pathlib import Path


# ── Text update helpers ───────────────────────────────────────────────────────


def _first_run_style(box: dict) -> dict:
    """Extract styling from the first run of the first paragraph."""
    for para in box.get("paragraphs", []):
        for run in para.get("runs", []):
            return {k: v for k, v in run.items() if k != "text"}
    style = box.get("boxStyle", {})
    return {
        "fontSize": style.get("fontSize", 16),
        "fontFamily": style.get("fontFamily", "Montserrat"),
        "color": style.get("color", "#000000"),
    }


def _apply_text_update(box: dict, new_text: str) -> dict:
    """Clone box, replacing its text content while preserving all styling."""
    box = copy.deepcopy(box)
    run_style = _first_run_style(box)
    lines = new_text.split("\n")

    new_paragraphs = []
    for line in lines:
        run = dict(run_style)
        run["text"] = line
        existing_para = box.get("paragraphs", [{}])[0] if box.get("paragraphs") else {}
        new_paragraphs.append(
            {
                "text": line,
                "align": existing_para.get("align", "left"),
                "spaceBefore": existing_para.get("spaceBefore", 0),
                "runs": [run],
            }
        )

    box["text"] = lines[0] if lines else ""
    box["paragraphs"] = new_paragraphs
    return box


# ── Sidecar → PPTX export (mirrors app.py _pptx_export) ─────────────────────


def _export_pptx(out_path: Path, slides: list[dict]) -> None:
    """Convert sidecar slide list to a .pptx file."""
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

    # Determine coordinate system from sidecar.
    # If slideWidthPt is present, x/y/w/h are in points (pt).
    # Otherwise assume the 896×504 canvas coordinate system.
    first_slide = slides[0] if slides else {}
    slide_width_pt = first_slide.get("slideWidthPt")

    def hex_rgb(h):
        h = str(h).lstrip("#")
        if len(h) == 6:
            try:
                return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
            except Exception:
                pass
        return None

    def align(s):
        if s == "center":
            return PP_ALIGN.CENTER
        if s == "right":
            return PP_ALIGN.RIGHT
        return PP_ALIGN.LEFT

    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(5143500)
    sw, sh = prs.slide_width, prs.slide_height

    if slide_width_pt:
        # pt-based sidecar: 1 pt = 12700 EMU
        sx = 12700.0
        sy = 12700.0
    else:
        # canvas-based sidecar: 896×504 canvas units → EMU
        sx = sw / 896.0
        sy = sh / 504.0

    blank = prs.slide_layouts[6]

    PRST = {
        "ellipse": "ellipse",
        "roundRect": "roundRect",
        "rect": "rect",
        "triangle": "triangle",
        "isoscelesTri": "triangle",
        "rtTriangle": "rtTriangle",
        "diamond": "diamond",
        "parallelogram": "parallelogram",
        "trapezoid": "trapezoid",
        "pentagon": "pentagon",
        "hexagon": "hexagon",
        "plus": "plus",
        "rightArrow": "rightArrow",
        "leftArrow": "leftArrow",
        "upArrow": "upArrow",
        "downArrow": "downArrow",
        "chevron": "chevron",
        "star4": "star4",
        "star5": "star5",
        "snip1Rect": "snip1Rect",
    }

    def apply_prst(sp, st, cr=0):
        sp_el = sp._element
        sp_pr = sp_el.find(qn("p:spPr"))
        if sp_pr is None:
            return
        pg = sp_pr.find(qn("a:prstGeom"))
        if pg is None:
            return
        pg.set("prst", PRST.get(st, st))
        av = pg.find(qn("a:avLst"))
        if av is not None:
            for g in av.findall(qn("a:gd")):
                av.remove(g)
        if st == "roundRect" and cr > 0:
            adj = max(0, min(50000, int(cr * 500)))
            if av is None:
                av = etree.SubElement(pg, qn("a:avLst"))
            g2 = etree.SubElement(av, qn("a:gd"))
            g2.set("name", "adj")
            g2.set("fmla", f"val {adj}")

    def apply_fill(sp, fill, bs):
        ft = fill.get("type") if isinstance(fill, dict) else None
        fc = fill.get("color") if isinstance(fill, dict) else None
        op = float(fill.get("opacity", 1.0)) if isinstance(fill, dict) else 1.0
        grad = fill.get("gradient") if isinstance(fill, dict) else None
        bg_grad = (bs or {}).get("bgGradient")
        if bg_grad:
            grad = bg_grad
            ft = "gradient"

        if ft == "gradient" and grad and grad.get("stops"):
            try:
                sp_pr = sp._element.find(qn("p:spPr"))
                if sp_pr is not None:
                    for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:blipFill"):
                        for el in sp_pr.findall(qn(tag)):
                            sp_pr.remove(el)
                    gf = etree.SubElement(sp_pr, qn("a:gradFill"))
                    gsl = etree.SubElement(gf, qn("a:gsLst"))
                    for stop in grad["stops"]:
                        gs = etree.SubElement(gsl, qn("a:gs"))
                        gs.set("pos", str(int(stop["pos"] * 100000)))
                        sr = etree.SubElement(gs, qn("a:srgbClr"))
                        sr.set("val", stop["color"].lstrip("#"))
                    ln = etree.SubElement(gf, qn("a:lin"))
                    ln.set("ang", str(int(grad.get("angle", 0) * 60000)))
                    ln.set("scaled", "0")
            except Exception:
                pass
        elif fc:
            rgb = hex_rgb(fc)
            if rgb:
                sp.fill.solid()
                sp.fill.fore_color.rgb = rgb
                if op < 1.0:
                    try:
                        sp_pr2 = sp._element.find(qn("p:spPr"))
                        if sp_pr2 is not None:
                            clr = sp_pr2.find(".//" + qn("a:srgbClr"))
                            if clr is not None:
                                ae = etree.SubElement(clr, qn("a:alpha"))
                                ae.set("val", str(max(0, min(100000, int(op * 100000)))))
                    except Exception:
                        pass
        else:
            sp.fill.background()

    def apply_line(sp, bs):
        bc = (bs or {}).get("borderColor", "")
        bw = float((bs or {}).get("borderWidth", 0) or 0)
        if bc and bc != "transparent" and bw > 0:
            try:
                rgb = hex_rgb(bc)
                if rgb:
                    sp.line.color.rgb = rgb
                    sp.line.width = Pt(bw)
            except Exception:
                pass
        else:
            try:
                sp.line.fill.background()
            except Exception:
                pass

    def apply_text(tf, paragraphs, plain, bs, font):
        if paragraphs:
            for pi, pd in enumerate(paragraphs):
                para = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
                try:
                    para.alignment = align(pd.get("align") or (bs or {}).get("textAlign", "left"))
                except Exception:
                    pass
                try:
                    sb = pd.get("spaceBefore", 0) or 0
                    if sb:
                        para.space_before = Pt(sb)
                except Exception:
                    pass
                for rd in pd.get("runs") or []:
                    rt = rd.get("text", "")
                    if not rt:
                        continue
                    run = para.add_run()
                    run.text = rt
                    run.font.size = Pt(rd.get("fontSize") or (bs or {}).get("fontSize", 16))
                    run.font.name = rd.get("fontFamily") or font
                    if rd.get("bold"):
                        run.font.bold = True
                    if rd.get("italic"):
                        run.font.italic = True
                    if rd.get("underline"):
                        run.font.underline = True
                    rgb = hex_rgb(rd.get("color") or (bs or {}).get("color", "#000000"))
                    if rgb:
                        run.font.color.rgb = rgb
        else:
            ba = (bs or {}).get("textAlign", "left")
            for li, line in enumerate((plain or "").split("\n")):
                para = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                try:
                    para.alignment = align(ba)
                except Exception:
                    pass
                if not line:
                    continue
                run = para.add_run()
                run.text = line
                run.font.size = Pt((bs or {}).get("fontSize", 16))
                run.font.name = font
                if (bs or {}).get("fontWeight") == "bold":
                    run.font.bold = True
                rgb = hex_rgb((bs or {}).get("color", "#000000"))
                if rgb:
                    run.font.color.rgb = rgb

    def apply_anchor(tf, bs):
        anchor = (bs or {}).get("textAnchor", "")
        if anchor:
            try:
                bp = tf._txBody.find(qn("a:bodyPr"))
                if bp is not None:
                    bp.set("anchor", anchor)
            except Exception:
                pass

    for sd in slides:
        slide = prs.slides.add_slide(blank)

        # Background
        try:
            bg_grad = sd.get("bgGradient")
            if bg_grad and bg_grad.get("stops"):
                bg_el = slide.background._element
                bg_pr = bg_el.find(qn("p:bg"))
                if bg_pr is None:
                    bg_pr = etree.SubElement(bg_el, qn("p:bg"))
                bg_prop = bg_pr.find(qn("p:bgPr"))
                if bg_prop is None:
                    bg_prop = etree.SubElement(bg_pr, qn("p:bgPr"))
                for tag in ("a:solidFill", "a:gradFill", "a:noFill"):
                    for el in bg_prop.findall(qn(tag)):
                        bg_prop.remove(el)
                gf = etree.SubElement(bg_prop, qn("a:gradFill"))
                gsl = etree.SubElement(gf, qn("a:gsLst"))
                for stop in bg_grad["stops"]:
                    gs = etree.SubElement(gsl, qn("a:gs"))
                    gs.set("pos", str(int(stop["pos"] * 100000)))
                    sr = etree.SubElement(gs, qn("a:srgbClr"))
                    sr.set("val", stop["color"].lstrip("#"))
                ln = etree.SubElement(gf, qn("a:lin"))
                ln.set("ang", str(int(bg_grad.get("angle", 0) * 60000)))
                ln.set("scaled", "0")
            else:
                bg = slide.background.fill
                bg.solid()
                rgb = hex_rgb(sd.get("bgColor", "#ffffff"))
                if rgb:
                    bg.fore_color.rgb = rgb
        except Exception:
            pass

        for box in sd.get("boxes", []):
            try:
                x = Emu(int(box.get("x", 0) * sx))
                y = Emu(int(box.get("y", 0) * sy))
                w = Emu(max(1, int(box.get("w", 100) * sx)))
                h = Emu(max(1, int(box.get("h", 50) * sy)))
                rot = float(box.get("rotation", 0) or 0)
                bs = box.get("boxStyle", {})
                font = bs.get("fontFamily") or "Montserrat"
                st = box.get("shapeType", "rect") or "rect"
                paras = box.get("paragraphs", [])
                has_text = bool(box.get("text") or paras)
                fill = box.get("fill")
                img = box.get("imageData", "")

                if img and "," in img:
                    try:
                        _, b64 = img.split(",", 1)
                        pic = slide.shapes.add_picture(io.BytesIO(_b64.b64decode(b64)), x, y, w, h)
                        if rot:
                            pic.rotation = rot
                    except Exception:
                        pass
                    continue

                if has_text and st not in ("", "textbox", "rect"):
                    sp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
                    apply_prst(sp, st, box.get("cornerRadius", 0))
                    apply_fill(sp, fill or {"type": "none"}, bs)
                    apply_line(sp, bs)
                    if rot:
                        sp.rotation = rot
                    tf = sp.text_frame
                    tf.word_wrap = True
                    apply_text(tf, paras, box.get("text", ""), bs, font)
                    apply_anchor(tf, bs)
                    continue

                if fill and not has_text and st not in ("", "textbox"):
                    sp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
                    apply_prst(sp, st, box.get("cornerRadius", 0))
                    apply_fill(sp, fill, bs)
                    apply_line(sp, bs)
                    if rot:
                        sp.rotation = rot
                    if sp.has_text_frame:
                        sp.text_frame.text = ""
                    continue

                tb = slide.shapes.add_textbox(x, y, w, h)
                if rot:
                    tb.rotation = rot
                tf = tb.text_frame
                tf.word_wrap = True
                apply_text(tf, paras, box.get("text", ""), bs, font)
                apply_anchor(tf, bs)
                apply_line(tb, bs)
                bgc = bs.get("bgColor")
                if bgc and bgc != "transparent":
                    rgb = hex_rgb(bgc)
                    if rgb:
                        tb.fill.solid()
                        tb.fill.fore_color.rgb = rgb

            except Exception:
                pass

    out_path.parent.mkdir(parents=True, exist_ok=True)
    fd, tmp = tempfile.mkstemp(suffix=".pptx", dir=out_path.parent)
    try:
        os.close(fd)
        prs.save(tmp)
        shutil.move(tmp, str(out_path))
    except Exception:
        try:
            os.unlink(tmp)
        except Exception:
            pass
        raise


# ── Main ──────────────────────────────────────────────────────────────────────


def main() -> None:
    parser = argparse.ArgumentParser(description="Build a PPTX from template sidecar JSON + content spec.")
    parser.add_argument("--template-sidecar", required=True, help="Path to template sidecar JSON (.<name>.pptx.json).")
    g = parser.add_mutually_exclusive_group(required=True)
    g.add_argument("--spec-file", help="Path to content-spec JSON file.")
    g.add_argument("--spec", help="Inline content-spec JSON string.")
    parser.add_argument("--out", required=True, help="Output .pptx path.")
    args = parser.parse_args()

    # Load template sidecar
    tpl_path = Path(args.template_sidecar)
    if not tpl_path.exists():
        print(json.dumps({"error": f"template sidecar not found: {tpl_path}"}))
        sys.exit(1)
    try:
        template = json.loads(tpl_path.read_text(encoding="utf-8"))
    except Exception as e:
        print(json.dumps({"error": f"could not read template sidecar: {e}"}))
        sys.exit(1)

    tpl_slides = template.get("slides", [])
    if not tpl_slides:
        print(json.dumps({"error": "template sidecar has no slides"}))
        sys.exit(1)

    # Load spec
    if args.spec_file:
        sp = Path(args.spec_file)
        if not sp.exists():
            print(json.dumps({"error": f"spec file not found: {sp}"}))
            sys.exit(1)
        try:
            spec = json.loads(sp.read_text(encoding="utf-8"))
        except Exception as e:
            print(json.dumps({"error": f"invalid spec JSON: {e}"}))
            sys.exit(1)
    else:
        try:
            spec = json.loads(args.spec)
        except Exception as e:
            print(json.dumps({"error": f"invalid spec JSON: {e}"}))
            sys.exit(1)

    out_path = Path(args.out)
    if out_path.suffix.lower() != ".pptx":
        out_path = out_path.with_suffix(".pptx")

    # Build new slide list
    new_slides: list[dict] = []
    for entry in spec.get("slides", []):
        tpl_idx = int(entry.get("template_slide_index", 0))
        tpl_idx = max(0, min(tpl_idx, len(tpl_slides) - 1))
        slide = copy.deepcopy(tpl_slides[tpl_idx])

        updates: dict[str, str] = entry.get("text_updates", {})
        if updates:
            new_boxes = []
            for box in slide.get("boxes", []):
                bid = box.get("id", "")
                if bid in updates:
                    box = _apply_text_update(box, updates[bid])
                new_boxes.append(box)
            slide["boxes"] = new_boxes

        new_slides.append(slide)

    # Export PPTX
    try:
        _export_pptx(out_path, new_slides)
    except Exception as exc:
        print(json.dumps({"error": f"export failed: {exc}"}))
        sys.exit(1)

    # Write companion sidecar JSON
    sidecar_path = out_path.parent / f".{out_path.name}.json"
    sidecar_path.write_text(
        json.dumps({"slides": new_slides}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    print(
        json.dumps(
            {
                "ok": True,
                "path": str(out_path.resolve()),
                "sidecar": str(sidecar_path.resolve()),
                "slides": len(new_slides),
            }
        )
    )


if __name__ == "__main__":
    main()
