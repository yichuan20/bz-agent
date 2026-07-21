#!/usr/bin/env python3
from __future__ import annotations

"""
diff-docs.py — compare two documents and return only the changed sections.

Instead of reading both full documents and asking the LLM to spot differences,
this script diffs the extracted text and returns a concise change summary.

Usage:
  python3 bzcode/scripts/diff-docs.py --before old.docx --after new.docx
  python3 bzcode/scripts/diff-docs.py --before v1.pdf --after v2.pdf --max-chars 10000
  python3 bzcode/scripts/diff-docs.py --before budget_q1.xlsx --after budget_q2.xlsx

Output:
  Prints JSON:
  {
    "before": "old.docx",
    "after":  "new.docx",
    "added_lines":   N,
    "removed_lines": N,
    "changed_sections": N,
    "diff": "<unified diff text>",
    "summary": "Added N sections, removed M sections, N lines changed.",
    "truncated": false
  }
"""
import argparse
import difflib
import json
import sys
from pathlib import Path

_MAX_CHARS = 16_000
_MAX_BYTES = 50 * 1024 * 1024


def _extract(path: Path) -> str:
    ext = path.suffix.lower()
    data = path.read_bytes()

    if ext == ".pdf":
        import pypdf, io

        reader = pypdf.PdfReader(io.BytesIO(data))
        parts = []
        for i, page in enumerate(reader.pages, 1):
            t = (page.extract_text() or "").strip()
            if t:
                parts.append(f"[Page {i}]\n{t}")
        return "\n\n".join(parts)

    if ext in (".docx", ".doc"):
        import docx, io

        doc = docx.Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    if ext in (".xlsx", ".xls"):
        import openpyxl, io

        wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
        parts = []
        for name in wb.sheetnames:
            ws = wb[name]
            rows = [r for r in ws.iter_rows(values_only=True) if any(c is not None for c in r)][:500]
            lines = ["  ".join(str(c) if c is not None else "" for c in row) for row in rows]
            parts.append(f"[Sheet: {name}]\n" + "\n".join(lines))
        return "\n\n".join(parts)

    if ext in (".pptx", ".ppt"):
        from pptx import Presentation
        import io

        prs = Presentation(io.BytesIO(data))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [s.text_frame.text.strip() for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
            parts.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(parts)

    if ext in (".txt", ".md", ".csv"):
        return data.decode("utf-8", errors="replace")

    raise ValueError(f"Unsupported format: {ext}")


def main() -> None:
    parser = argparse.ArgumentParser(description="Diff two documents and return changed sections.")
    parser.add_argument("--before", required=True, help="Original document path.")
    parser.add_argument("--after", required=True, help="New document path.")
    parser.add_argument("--max-chars", type=int, default=_MAX_CHARS, help="Character budget for diff output.")
    parser.add_argument("--context", type=int, default=3, help="Context lines around each change (default 3).")
    args = parser.parse_args()

    for label, path_str in [("before", args.before), ("after", args.after)]:
        p = Path(path_str)
        if not p.exists():
            print(json.dumps({"error": f"{label} file not found: {p}"}), flush=True)
            sys.exit(1)
        if p.stat().st_size > _MAX_BYTES:
            print(json.dumps({"error": f"{label} file too large (max 50 MB)"}), flush=True)
            sys.exit(1)

    p_before = Path(args.before)
    p_after = Path(args.after)

    try:
        text_before = _extract(p_before)
        text_after = _extract(p_after)
    except Exception as exc:
        print(json.dumps({"error": f"could not parse: {exc}"}), flush=True)
        sys.exit(1)

    lines_before = text_before.splitlines(keepends=True)
    lines_after = text_after.splitlines(keepends=True)

    diff_lines = list(
        difflib.unified_diff(
            lines_before,
            lines_after,
            fromfile=p_before.name,
            tofile=p_after.name,
            n=args.context,
        )
    )

    added = sum(1 for l in diff_lines if l.startswith("+") and not l.startswith("+++"))
    removed = sum(1 for l in diff_lines if l.startswith("-") and not l.startswith("---"))

    # Count changed sections (hunk headers)
    changed_sections = sum(1 for l in diff_lines if l.startswith("@@"))

    diff_text = "".join(diff_lines)
    truncated = len(diff_text) > args.max_chars
    if truncated:
        diff_text = diff_text[: args.max_chars]

    if not diff_lines:
        summary = "No differences found — the documents are identical."
    else:
        parts = []
        if added:
            parts.append(f"{added} line{'s' if added != 1 else ''} added")
        if removed:
            parts.append(f"{removed} line{'s' if removed != 1 else ''} removed")
        summary = f"{changed_sections} section{'s' if changed_sections != 1 else ''} changed: {', '.join(parts)}."

    result = {
        "before": p_before.name,
        "after": p_after.name,
        "added_lines": added,
        "removed_lines": removed,
        "changed_sections": changed_sections,
        "summary": summary,
        "diff": diff_text,
        "truncated": truncated,
    }
    print(json.dumps(result, ensure_ascii=False), flush=True)


if __name__ == "__main__":
    main()
