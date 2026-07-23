"""Document service — parse, save, and download .docx/.pdf/.xlsx/.pptx files.

All parsing/conversion functions are copied verbatim from the old server.py
(§ 7 DOCUMENT PARSING).  Only the sidecar helpers (from app.py) and the
``DocService`` wrapper are new; the logic itself is unchanged.

Sidecar convention (from app.py):
  ``<file>.json``       — DOCX block sidecar (same location, name + ".json")
  Sidecar is authoritative; DOCX binary regenerated from it on save.
  force_refresh discard policy: only discard when the DOCX is newer than sidecar.

Cursor store: in-memory, resets on restart (UI state only).
"""

from __future__ import annotations

import asyncio
import copy
import hashlib
import json
from pathlib import Path
from typing import Any

# ── Constants (from server.py) ────────────────────────────────────────────────

_MAX_DOC_CHARS = 80_000
_MAX_DOC_BYTES = 50 * 1024 * 1024  # 50 MB

_DOCX_EXTS = {".docx", ".doc"}

# ── In-memory cursor store ────────────────────────────────────────────────────

_CURSOR_STORE: dict[str, dict[str, int]] = {}

# ── Sidecar helpers (from app.py) ─────────────────────────────────────────────


def _doc_sidecar_path(p: Path) -> Path:
    return Path(str(p) + ".json")


def _extract_images(blocks: list) -> tuple[list, dict]:
    """Move imageUrl out of block styles into a shared images dict keyed by stable ID."""
    images: dict = {}
    for block in blocks:
        for style in block.get("styles") or []:
            url = style.get("imageUrl")
            if not url:
                continue
            prefix = url[:200] if isinstance(url, str) else ""
            key = "img-" + hashlib.md5(prefix.encode()).hexdigest()[:8]
            images[key] = {
                "url": url,
                **({"width": style["imageWidth"]} if style.get("imageWidth") else {}),
                **({"height": style["imageHeight"]} if style.get("imageHeight") else {}),
            }
            style.pop("imageUrl", None)
            style.pop("imageWidth", None)
            style.pop("imageHeight", None)
            style["imageId"] = key
    return blocks, images


def _expand_images(blocks: list, images: dict) -> list:
    """Inline imageUrl/Width/Height back into block styles from the images dict."""
    for block in blocks:
        for style in block.get("styles") or []:
            img_id = style.get("imageId")
            if img_id and img_id in images:
                img = images[img_id]
                style["imageUrl"] = img["url"]
                if img.get("width"):
                    style["imageWidth"] = img["width"]
                if img.get("height"):
                    style["imageHeight"] = img["height"]
    return blocks


# ── DOCX ↔ Block JSON conversion — copied verbatim from server.py § 7 ─────────


def _docx_to_blocks(data: bytes) -> list:
    """Convert DOCX binary → Block[] in bz-office JSON format."""
    import base64
    import io
    import secrets

    import docx as _docx
    from docx.oxml.ns import qn

    doc = _docx.Document(io.BytesIO(data))
    blocks = []

    # Resolve major/minor theme fonts once for the whole document
    _major_font = _minor_font = None
    try:
        import xml.etree.ElementTree as _ET
        import zipfile as _zf

        with _zf.ZipFile(io.BytesIO(data)) as _z:
            _theme_names = [n for n in _z.namelist() if n.lower().endswith("theme1.xml") and "theme" in n.lower()]
            if _theme_names:
                _theme_xml = _z.read(_theme_names[0])
                _theme_root = _ET.fromstring(_theme_xml)
                _ans = "http://schemas.openxmlformats.org/drawingml/2006/main"
                _fs = _theme_root.find(f".//{{{_ans}}}fontScheme")
                if _fs is not None:
                    _mj = _fs.find(f"{{{_ans}}}majorFont/{{{_ans}}}latin")
                    _mn = _fs.find(f"{{{_ans}}}minorFont/{{{_ans}}}latin")
                    if _mj is not None:
                        _major_font = _mj.get("typeface")
                    if _mn is not None:
                        _minor_font = _mn.get("typeface")
    except Exception:
        pass

    def _resolve_theme(name):
        if not name:
            return None
        if name in ("+mn-lt", "+Body"):
            return _minor_font
        if name in ("+mj-lt", "+Heading"):
            return _major_font
        if name.startswith("+"):
            return None
        return name

    _default_font = None
    try:
        v = doc.styles["Normal"].font.name
        _default_font = _resolve_theme(v) or v or None
    except Exception:
        pass

    _default_font_size_pt = None
    try:
        sz = doc.styles["Normal"].font.size
        if sz:
            _default_font_size_pt = int(sz.pt)
    except Exception:
        pass
    if not _default_font_size_pt:
        try:
            _wns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
            styles_el = doc.part.styles._element
            defs = styles_el.find(f"{{{_wns}}}docDefaults")
            if defs is not None:
                rPr = defs.find(f"{{{_wns}}}rPrDefault/{{{_wns}}}rPr")
                if rPr is not None:
                    sz_el = rPr.find(f"{{{_wns}}}sz")
                    if sz_el is not None:
                        val = sz_el.get(qn("w:val"))
                        if val:
                            _default_font_size_pt = int(val) // 2
        except Exception:
            pass
    if not _default_font:
        try:
            _wns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
            defs = doc.element.find(f".//{{{_wns}}}docDefaults")
            if defs is not None:
                rFonts = defs.find(f".//{{{_wns}}}rFonts")
                if rFonts is not None:
                    raw = rFonts.get(qn("w:ascii")) or rFonts.get(qn("w:hAnsi"))
                    _default_font = _resolve_theme(raw)
        except Exception:
            pass

    def _resolve_font(name):
        if not name:
            return None
        if name in ("+mj-lt", "+Heading"):
            return _major_font
        if name in ("+mn-lt", "+Body"):
            return _minor_font
        nl = name.lower()
        if nl.startswith("major"):
            return _major_font
        if nl.startswith("minor"):
            return _minor_font
        if name.startswith("+"):
            return None
        return name

    def _rFonts_font(rFonts):
        for attr in (qn("w:ascii"), qn("w:hAnsi"), qn("w:asciiTheme"), qn("w:hAnsiTheme"), qn("w:cs")):
            v = _resolve_font(rFonts.get(attr))
            if v:
                return v
        return None

    def _get_run_font(run, para=None):
        try:
            rPr = run._r.rPr
            if rPr is not None:
                rFonts = rPr.find(qn("w:rFonts"))
                if rFonts is not None:
                    v = _rFonts_font(rFonts)
                    if v:
                        return v
        except Exception:
            pass
        if para is not None:
            try:
                _wns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
                style_el = para.style.element if para.style else None
                if style_el is not None:
                    style_rPr = style_el.find(f".//{{{_wns}}}rPr")
                    if style_rPr is not None:
                        rFonts = style_rPr.find(qn("w:rFonts"))
                        if rFonts is not None:
                            v = _rFonts_font(rFonts)
                            if v:
                                return v
            except Exception:
                pass
        try:
            v = _resolve_font(run.font.name)
            if v:
                return v
        except Exception:
            pass
        return _minor_font or _default_font or None

    def _effective_font_size_pt(run, para):
        if run.font.size:
            return int(run.font.size.pt)
        try:
            if run.style and run.style.font.size:
                return int(run.style.font.size.pt)
        except Exception:
            pass
        try:
            if para.style and para.style.font.size:
                return int(para.style.font.size.pt)
        except Exception:
            pass
        return _default_font_size_pt

    def _run_styles(para) -> list:
        styles, pos = [], 0
        for run in para.runs:
            n = len(run.text)
            if not n:
                pos += n
                continue
            sr: dict[str, Any] = {"start": pos, "end": pos + n}
            if run.bold:
                sr["isBold"] = True
            if run.italic:
                sr["isItalic"] = True
            if run.underline:
                sr["isUnderlined"] = True
            if getattr(run.font, "strike", None):
                sr["isStrikethrough"] = True
            eff_size = _effective_font_size_pt(run, para)
            if eff_size:
                sr["fontSize"] = eff_size
            if run.font.color and run.font.color.type is not None:
                try:
                    sr["textColor"] = f"#{run.font.color.rgb}"
                except Exception:
                    pass
            fname = _get_run_font(run, para)
            if fname:
                sr["fontFamily"] = fname
            if len(sr) > 2:
                styles.append(sr)
            pos += n
        return styles

    _H_WNS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    _FALLBACK_HEADING_SIZES = {"1": 24, "2": 20, "3": 18, "4": 16}

    def _resolve_style_size_pt(style_name: str):
        visited: set = set()
        try:
            s = doc.styles[style_name]
        except Exception:
            return None
        while s is not None and s.name not in visited:
            visited.add(s.name)
            if s.font.size:
                return int(s.font.size.pt)
            s = s.base_style
        return _default_font_size_pt

    _heading_props: dict = {}
    for _level in ("1", "2", "3", "4"):
        _props: dict = {}
        _sz = _resolve_style_size_pt(f"Heading {_level}")
        if _sz:
            _props["fontSize"] = _sz
        try:
            _hs = doc.styles[f"Heading {_level}"]
            if _hs.font.bold:
                _props["isBold"] = True
            if _hs.font.italic:
                _props["isItalic"] = True
            _rPr = _hs._element.find(f"{{{_H_WNS}}}rPr")
            _col = _rPr.find(f"{{{_H_WNS}}}color") if _rPr is not None else None
            if _col is not None:
                _cv = _col.get(qn("w:val"))
                if _cv and _cv.lower() != "auto":
                    _props["textColor"] = f"#{_cv.upper()}"
        except Exception:
            pass
        _heading_props[_level] = _props

    def _heading_size(style_name: str):
        for level in ("1", "2", "3", "4"):
            if style_name == f"Heading {level}":
                return _heading_props[level].get("fontSize") or _FALLBACK_HEADING_SIZES[level]
        return None

    _ALIGN_MAP: dict = {}
    try:
        from docx.enum.text import WD_ALIGN_PARAGRAPH as _WAP

        _ALIGN_MAP = {
            _WAP.CENTER: "center",
            _WAP.RIGHT: "right",
            _WAP.JUSTIFY: "justify",
        }
    except Exception:
        pass

    def _extract_drawing_style(drawing, para):
        try:
            _wp = "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
            _a = "http://schemas.openxmlformats.org/drawingml/2006/main"
            _r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
            container = drawing.find(f"{{{_wp}}}inline") or drawing.find(f"{{{_wp}}}anchor")
            if container is None:
                return None
            extent = container.find(f"{{{_wp}}}extent")
            cx = int(extent.get("cx", 0)) if extent is not None else 0
            cy = int(extent.get("cy", 0)) if extent is not None else 0
            width_px = max(1, round(cx / 9525))
            height_px = max(1, round(cy / 9525))
            blip = container.find(f".//{{{_a}}}blip")
            if blip is None:
                return None
            r_id = blip.get(f"{{{_r}}}embed")
            if not r_id:
                return None
            image_part = para.part.related_parts[r_id]
            b64 = base64.b64encode(image_part.blob).decode("ascii")
            data_url = f"data:{image_part.content_type};base64,{b64}"
            return {"imageUrl": data_url, "imageWidth": width_px, "imageHeight": height_px}
        except Exception:
            return None

    def _para_to_block(para) -> dict:
        text = para.text
        styles = _run_styles(para)
        block: dict[str, Any] = {"text": text, "styles": styles}

        pos = 0
        insertions = []
        for run_el in para._p.findall(qn("w:r")):
            drawing = run_el.find(qn("w:drawing"))
            if drawing is not None:
                img = _extract_drawing_style(drawing, para)
                if img:
                    insertions.append((pos, img))
            t_el = run_el.find(qn("w:t"))
            pos += len(t_el.text if t_el is not None and t_el.text else "")

        if insertions:
            chars = list(block["text"])
            cur_styles = list(block["styles"])
            for ins_pos, img_style in sorted(insertions, key=lambda x: -x[0]):
                chars.insert(ins_pos, " ")
                shifted = []
                for sr in cur_styles:
                    nr = dict(sr)
                    if nr["start"] >= ins_pos:
                        nr["start"] += 1
                        nr["end"] += 1
                    elif nr["end"] > ins_pos:
                        nr["end"] += 1
                    shifted.append(nr)
                shifted.append({"start": ins_pos, "end": ins_pos + 1, **img_style})
                cur_styles = sorted(shifted, key=lambda s: s["start"])
            block["text"] = "".join(chars)
            block["styles"] = cur_styles

        sname = para.style.name if para.style else ""
        size = _heading_size(sname)
        if size:
            level_key = sname.split()[-1] if sname.startswith("Heading ") else None
            props = _heading_props.get(level_key, {}) if level_key else {}
            heading_font = _major_font or _minor_font or None
            heading_style: dict = {"start": 0, "end": len(block["text"]), "fontSize": size}
            if props.get("isBold", True):
                heading_style["isBold"] = True
            if props.get("isItalic"):
                heading_style["isItalic"] = True
            if props.get("textColor"):
                heading_style["textColor"] = props["textColor"]
            if heading_font:
                heading_style["fontFamily"] = heading_font
            block["styles"] = [heading_style]
            block["headingLevel"] = int(level_key) if level_key else None

        if para.alignment in _ALIGN_MAP:
            block["alignment"] = _ALIGN_MAP[para.alignment]

        try:
            numPr = para._p.pPr.numPr if para._p.pPr is not None else None
            if numPr is not None:
                block["prefix"] = "•"
                ilvl = numPr.ilvl
                block["indent"] = int(ilvl.val) + 1 if ilvl is not None else 1
        except Exception:
            pass

        return block

    for child in doc.element.body.iterchildren():
        tag = child.tag.split("}")[-1]

        if tag == "p":
            try:
                from docx.text.paragraph import Paragraph

                para = Paragraph(child, doc)
                block = _para_to_block(para)
                if block["text"].strip() or not blocks:
                    blocks.append(block)
            except Exception:
                pass

        elif tag == "tbl":
            try:
                from docx.table import Table

                table = Table(child, doc)
                tid = secrets.token_hex(8)
                n_rows = len(table.rows)
                n_cols = max((len(r.cells) for r in table.rows), default=0)
                for r_idx, row in enumerate(table.rows):
                    for c_idx, cell in enumerate(row.cells):
                        cell_text = ""
                        cell_styles: list = []
                        for p_idx, para in enumerate(cell.paragraphs):
                            if p_idx > 0:
                                cell_text += "\n"
                            offset = len(cell_text)
                            for sr in _run_styles(para):
                                cell_styles.append(
                                    {
                                        **sr,
                                        "start": sr["start"] + offset,
                                        "end": sr["end"] + offset,
                                    }
                                )
                            cell_text += para.text
                        blocks.append(
                            {
                                "text": cell_text,
                                "styles": cell_styles,
                                "isTableCell": True,
                                "tableId": tid,
                                "rowIndex": r_idx,
                                "columnIndex": c_idx,
                                "numberOfRows": n_rows,
                                "numberOfColumns": n_cols,
                            }
                        )
            except Exception:
                pass

    effective_default = _minor_font or _default_font or None
    return {"blocks": blocks, "defaultFont": effective_default}


def _blocks_to_docx(blocks: list) -> bytes:
    """Convert Block[] (bz-office format) → DOCX binary."""
    import base64
    import io

    import docx as _docx
    from docx.oxml.ns import qn
    from docx.shared import Emu, Pt, RGBColor
    from lxml import etree

    PX_TO_EMU = 9525
    WP = "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"

    def _make_float(run, wrap, placed_x_css):
        drawing = run._r.find(qn("w:drawing"))
        if drawing is None:
            return
        inline = drawing.find(f"{{{WP}}}inline")
        if inline is None:
            return
        extent = inline.find(f"{{{WP}}}extent")
        cx = extent.get("cx", "0") if extent is not None else "0"
        cy = extent.get("cy", "0") if extent is not None else "0"
        behind = "1" if wrap == "behind" else "0"
        CANVAS_LEFT_CSS = 100
        COLUMN_W_CSS = 621
        col_x = max(0, (placed_x_css or 0) - CANVAS_LEFT_CSS)
        pos_x = int(col_x * PX_TO_EMU)
        wrap_side = "right" if col_x < COLUMN_W_CSS / 2 else "left"
        wrap_el = "<wp:wrapBehindDoc/>" if wrap == "behind" else f'<wp:wrapSquare wrapText="{wrap_side}"/>'
        anchor = etree.fromstring(
            f'<wp:anchor xmlns:wp="{WP}"'
            f' distT="0" distB="114300" distL="114300" distR="114300"'
            f' simplePos="0" relativeHeight="251658240" behindDoc="{behind}"'
            f' locked="0" layoutInCell="1" allowOverlap="0">'
            f'<wp:simplePos x="0" y="0"/>'
            f'<wp:positionH relativeFrom="column"><wp:posOffset>{pos_x}</wp:posOffset></wp:positionH>'
            f'<wp:positionV relativeFrom="paragraph"><wp:posOffset>0</wp:posOffset></wp:positionV>'
            f'<wp:extent cx="{cx}" cy="{cy}"/>'
            f'<wp:effectExtent l="0" t="0" r="0" b="0"/>'
            f"{wrap_el}"
            f"</wp:anchor>"
        )
        for tag in (f"{{{WP}}}docPr", f"{{{WP}}}cNvGraphicFramePr", f"{{{A}}}graphic"):
            child = inline.find(tag)
            if child is not None:
                anchor.append(child)
        drawing.remove(inline)
        drawing.append(anchor)

    doc = _docx.Document()
    i = 0
    while i < len(blocks):
        b = blocks[i]
        if b.get("isTableCell"):
            tid = b.get("tableId")
            cells = [c for c in blocks if c.get("tableId") == tid]
            n_rows = b.get("numberOfRows", 1)
            n_cols = b.get("numberOfColumns", 1)
            tbl = doc.add_table(rows=n_rows, cols=n_cols)
            tbl.style = "Table Grid"
            for cell in cells:
                r, c = cell.get("rowIndex", 0), cell.get("columnIndex", 0)
                try:
                    tbl.rows[r].cells[c].text = cell.get("text", "")
                    if r == 0:
                        for run in tbl.rows[r].cells[c].paragraphs[0].runs:
                            run.bold = True
                except Exception:
                    pass
            while i < len(blocks) and blocks[i].get("tableId") == tid:
                i += 1
            continue

        text = b.get("text", "")
        styles = b.get("styles", [])
        prefix = b.get("prefix", "")
        indent = b.get("indent", 0)
        heading_size = b.get("headingLevel") or None

        if heading_size:
            para = doc.add_heading(text, level=heading_size)
            heading_font = next((sr.get("fontFamily") for sr in styles if sr.get("fontFamily")), None)
            if heading_font:
                for run in para.runs:
                    run.font.name = heading_font
        elif prefix == "•":
            para = doc.add_paragraph(style="List Bullet")
            para.add_run(text)
        else:
            para = doc.add_paragraph()
            if not styles:
                para.add_run(text)
            else:
                cursor = 0
                for sr in sorted(styles, key=lambda s: s.get("start", 0)):
                    s, e = sr.get("start", 0), sr.get("end", len(text))
                    if cursor < s:
                        para.add_run(text[cursor:s])
                    if sr.get("imageUrl"):
                        try:
                            data_url = sr["imageUrl"]
                            if data_url.startswith("data:"):
                                _, b64_data = data_url.split(",", 1)
                                img_bytes = base64.b64decode(b64_data)
                            else:
                                import urllib.request as _req

                                with _req.urlopen(data_url, timeout=10) as resp:  # noqa: ASYNC210
                                    img_bytes = resp.read()
                            wrap = sr.get("imageWrap", "inline")
                            img_run = para.add_run()
                            w_emu = sr.get("imageWidth", 64) * PX_TO_EMU
                            h_emu = sr.get("imageHeight", 64) * PX_TO_EMU
                            img_run.add_picture(io.BytesIO(img_bytes), width=Emu(w_emu), height=Emu(h_emu))
                            if wrap in ("square", "behind"):
                                _make_float(img_run, wrap, sr.get("imagePlacedX"))
                        except Exception:
                            para.add_run(text[s:e])
                    else:
                        run = para.add_run(text[s:e])
                        run.bold = sr.get("isBold", False)
                        run.italic = sr.get("isItalic", False)
                        run.underline = sr.get("isUnderlined", False)
                        run.font.strike = sr.get("isStrikethrough", False) or None
                        if sr.get("fontFamily"):
                            run.font.name = sr["fontFamily"]
                        if sr.get("fontSize"):
                            run.font.size = Pt(sr["fontSize"])
                        if sr.get("textColor"):
                            try:
                                hex_c = sr["textColor"].lstrip("#")
                                run.font.color.rgb = RGBColor(
                                    int(hex_c[0:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16)
                                )
                            except Exception:
                                pass
                    cursor = e
                if cursor < len(text):
                    para.add_run(text[cursor:])

        if indent and not heading_size:
            from docx.shared import Inches

            para.paragraph_format.left_indent = Inches(indent * 0.25)

        i += 1

    buf = __import__("io").BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _parse_pdf(data: bytes) -> tuple[int, str]:
    import io

    import pypdf

    reader = pypdf.PdfReader(io.BytesIO(data))
    pages = len(reader.pages)
    parts = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        if text.strip():
            parts.append(f"# Page {i}\n\n{text.strip()}")
    return pages, "\n\n".join(parts)


def _parse_docx_text(data: bytes) -> tuple[int, str]:
    import io

    import docx

    doc = docx.Document(io.BytesIO(data))
    parts = []
    heading_map = {1: "#", 2: "##", 3: "###", 4: "####"}
    for para in doc.paragraphs:
        style = para.style.name if para.style else ""
        text = para.text.strip()
        if not text:
            continue
        level = next((int(s) for s in ("1", "2", "3", "4") if style == f"Heading {s}"), None)
        if level:
            parts.append(f"{heading_map[level]} {text}")
        else:
            parts.append(text)
    for table in doc.tables:
        rows = []
        for i, row in enumerate(table.rows):
            cells = [c.text.strip() for c in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
            if i == 0:
                rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
        parts.append("\n".join(rows))
    page_count = max(1, len(parts) // 10)
    return page_count, "\n\n".join(parts)


def _parse_xlsx(data: bytes) -> tuple[int, str]:
    import io

    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        rows = [r for r in rows if any(c is not None for c in r)]
        if not rows:
            continue
        rows = rows[:1000]
        parts.append(f"## Sheet: {sheet_name}")
        header = rows[0]
        parts.append("| " + " | ".join(str(c) if c is not None else "" for c in header) + " |")
        parts.append("| " + " | ".join(["---"] * len(header)) + " |")
        for row in rows[1:]:
            parts.append("| " + " | ".join(str(c) if c is not None else "" for c in row) + " |")
    return len(wb.sheetnames), "\n\n".join(parts)


def _parse_pptx(data: bytes) -> tuple[int, str]:
    import io

    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        title_text = ""
        body_lines = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if shape.shape_type == 13:
                continue
            if not title_text and hasattr(shape, "placeholder_format") and shape.placeholder_format:
                title_text = text
            else:
                for para in shape.text_frame.paragraphs:
                    pt = para.text.strip()
                    if pt:
                        body_lines.append(f"- {pt}")
        header = f"## Slide {i}: {title_text}" if title_text else f"## Slide {i}"
        parts.append(header + ("\n" + "\n".join(body_lines) if body_lines else ""))
    return len(prs.slides), "\n\n".join(parts)


def _detect_and_parse(filename: str, data: bytes) -> dict:
    ext = Path(filename).suffix.lower()
    fmt = ext.lstrip(".")
    if fmt in ("doc", "xls", "ppt"):
        fmt = {"doc": "docx", "xls": "xlsx", "ppt": "pptx"}[fmt]

    if ext in _DOCX_EXTS:
        result = _docx_to_blocks(data)
        blocks = result["blocks"]
        word_count = sum(len(b.get("text", "").split()) for b in blocks)
        return {
            "filename": filename,
            "type": fmt,
            "pages": max(1, len([b for b in blocks if not b.get("isTableCell")]) // 30),
            "wordCount": word_count,
            "truncated": False,
            "blocks": blocks,
            "defaultFont": result.get("defaultFont"),
        }

    parsers = {
        ".pdf": _parse_pdf,
        ".xlsx": _parse_xlsx,
        ".xls": _parse_xlsx,
        ".pptx": _parse_pptx,
        ".ppt": _parse_pptx,
    }
    if ext not in parsers:
        raise ValueError(f"unsupported format: {ext or '(no extension)'}")
    pages, content = parsers[ext](data)
    truncated = len(content) > _MAX_DOC_CHARS
    if truncated:
        content = content[:_MAX_DOC_CHARS]
    return {
        "filename": filename,
        "type": fmt,
        "pages": pages,
        "wordCount": len(content.split()),
        "truncated": truncated,
        "content": content,
    }


# ── DocService — thin wrapper exposed to routes ────────────────────────────────


class DocService:
    """Thin wrapper around the parsing/save functions for use as a FastAPI dependency."""

    async def parse(self, path: str, *, force: bool = False) -> dict[str, Any]:
        """Parse a document at *path*.  Uses sidecar cache for DOCX unless *force*."""
        p = Path(path)

        def _do() -> dict[str, Any]:
            if not p.exists():
                raise FileNotFoundError(path)
            if p.stat().st_size > _MAX_DOC_BYTES:
                raise ValueError(f"file too large (max 50 MB): {path}")
            if p.suffix.lower() in _DOCX_EXTS:
                sc_path = _doc_sidecar_path(p)
                if force and sc_path.exists():
                    try:
                        if p.stat().st_mtime > sc_path.stat().st_mtime:
                            sc_path.unlink()
                    except Exception:
                        pass
                if sc_path.exists():
                    try:
                        sidecar = json.loads(sc_path.read_text(encoding="utf-8"))
                        has_valid = any(
                            s.get("fontFamily") or s.get("imageUrl") or s.get("imageId")
                            for b in (sidecar.get("blocks") or [])
                            for s in (b.get("styles") or [])
                        )
                        if has_valid or not sidecar.get("blocks"):
                            if sidecar.get("images"):
                                _expand_images(sidecar["blocks"], sidecar["images"])
                            return sidecar
                    except Exception:
                        pass

            data = p.read_bytes()
            result = _detect_and_parse(p.name, data)
            if p.suffix.lower() in _DOCX_EXTS:
                try:
                    _doc_sidecar_path(p).write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
                except Exception:
                    pass
            return result

        return await asyncio.to_thread(_do)

    async def parse_upload(self, data: bytes, filename: str) -> dict[str, Any]:
        """Parse an uploaded file (bytes). No sidecar written."""
        return await asyncio.to_thread(_detect_and_parse, filename, data)

    def set_cursor(self, path: str, sel_start: int, sel_end: int) -> None:
        _CURSOR_STORE[path] = {"selStart": sel_start, "selEnd": sel_end}

    def get_cursor(self, path: str) -> dict[str, int]:
        return _CURSOR_STORE.get(path, {"selStart": 0, "selEnd": 0})

    async def save(self, path: str, blocks: list[dict[str, Any]]) -> dict[str, Any]:
        p = Path(path)
        if p.suffix.lower() not in _DOCX_EXTS:
            raise ValueError("only DOCX files can be saved")

        def _do() -> dict[str, Any]:
            sc_path = _doc_sidecar_path(p)
            word_count = sum(len(b.get("text", "").split()) for b in blocks)
            if sc_path.exists():
                try:
                    sidecar = json.loads(sc_path.read_text(encoding="utf-8"))
                except Exception:
                    sidecar = {"filename": p.name, "type": "docx", "truncated": False}
            else:
                sidecar = {"filename": p.name, "type": "docx", "truncated": False}
            save_blocks = copy.deepcopy(blocks)
            save_blocks, new_images = _extract_images(save_blocks)
            existing_images = sidecar.get("images") or {}
            existing_images.update(new_images)
            if existing_images:
                sidecar["images"] = existing_images
            sidecar["blocks"] = save_blocks
            sidecar["wordCount"] = word_count
            try:
                docx_bytes = _blocks_to_docx(blocks)
                p.write_bytes(docx_bytes)
            except Exception:
                pass
            try:
                sc_path.write_text(json.dumps(sidecar, ensure_ascii=False, indent=2), encoding="utf-8")
            except Exception:
                pass
            return {"ok": True, "path": path, "wordCount": word_count}

        return await asyncio.to_thread(_do)

    async def download(self, path: str) -> bytes:
        p = Path(path)
        sc_path = _doc_sidecar_path(p)

        def _do() -> bytes:
            if sc_path.exists():
                try:
                    sidecar = json.loads(sc_path.read_text(encoding="utf-8"))
                    blocks = sidecar.get("blocks") or []
                    if sidecar.get("images"):
                        _expand_images(blocks, sidecar["images"])
                    return _blocks_to_docx(blocks)
                except Exception:
                    pass
            if p.exists():
                return p.read_bytes()
            raise FileNotFoundError(path)

        return await asyncio.to_thread(_do)
