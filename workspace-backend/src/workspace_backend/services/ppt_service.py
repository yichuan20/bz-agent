"""PPT service — load and save .pptx presentations.

``_pptx_to_slides`` is copied verbatim from the old app.py ppt_load handler.
The ``PptService`` wrapper is the only new code.

Slide format: { bgColor, bgGradient?, slideWidthPt, boxes: [...] }
Box format:   { id, x, y, w, h, rotation, shapeType, text, paragraphs, styles, boxStyle, ... }
"""

from __future__ import annotations

import asyncio
import json
from pathlib import Path
from typing import Any


def _sidecar(path: Path) -> Path:
    return path.parent / f".{path.name}.json"


def _read_sidecar(path: Path) -> dict[str, Any] | None:
    sc = _sidecar(path)
    if not sc.exists():
        return None
    try:
        return json.loads(sc.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return None


def _write_sidecar(path: Path, data: dict[str, Any]) -> None:
    sc = _sidecar(path)
    sc.write_text(json.dumps(data, ensure_ascii=False), encoding="utf-8")


def _pptx_to_slides(p: Path) -> list[dict[str, Any]]:
    """Convert a .pptx file → list of slide dicts. Copied verbatim from old app.py."""
    import base64 as _b64
    import secrets as _sec2
    import zipfile as _zf

    from lxml import etree as _etree
    from pptx import Presentation
    from pptx.oxml.ns import qn as _qn

    prs = Presentation(str(p))
    sw = prs.slide_width
    sh = prs.slide_height
    CW, CH = 896, 504
    sx = CW / sw if sw else 1
    sy = CH / sh if sh else 1

    # ── Build theme color lookup ──────────────────────────────────────────────
    theme_colors: dict = {}
    try:
        _NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
        with _zf.ZipFile(str(p)) as _z:
            _theme_files = sorted([n for n in _z.namelist() if "/theme/" in n and n.endswith(".xml")])
            if _theme_files:
                _tel = _etree.fromstring(_z.read(_theme_files[0]))
                _clr = _tel.find(f".//{{{_NS_A}}}clrScheme")
                if _clr is not None:
                    for _child in _clr:
                        _tag = _child.tag.split("}")[1]
                        _srgb = _child.find(f"{{{_NS_A}}}srgbClr")
                        _sysc = _child.find(f"{{{_NS_A}}}sysClr")
                        if _srgb is not None:
                            _v = _srgb.get("val", "")
                            if len(_v) == 6:
                                theme_colors[_tag] = f"#{_v.lower()}"
                        elif _sysc is not None:
                            _last = _sysc.get("lastClr", "")
                            if len(_last) == 6:
                                theme_colors[_tag] = f"#{_last.lower()}"
    except Exception:
        pass

    _THEME_SLOT = {
        1: "dk1",
        2: "lt1",
        3: "dk2",
        4: "lt2",
        5: "accent1",
        6: "accent2",
        7: "accent3",
        8: "accent4",
        9: "accent5",
        10: "accent6",
        11: "hlink",
        12: "folHlink",
    }
    _NS_A_FULL = "http://schemas.openxmlformats.org/drawingml/2006/main"

    def _resolve_color(color_obj):
        try:
            rgb = color_obj.rgb
            hex6 = str(rgb)
            if len(hex6) == 6:
                return f"#{hex6.lower()}"
            return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
        except Exception:
            pass
        try:
            tc = color_obj.theme_color
            tc_int = int(getattr(tc, "real", tc))
            slot = _THEME_SLOT.get(tc_int)
            if slot and slot in theme_colors:
                raw = theme_colors[slot]
                try:
                    lum_mod = color_obj._element.find(f".//{{{_NS_A_FULL}}}lumMod")
                    if lum_mod is not None:
                        factor = int(lum_mod.get("val", "100000")) / 100000
                        h = raw.lstrip("#")
                        r2, g2, b2 = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
                        r2 = int(min(255, r2 * factor))
                        g2 = int(min(255, g2 * factor))
                        b2 = int(min(255, b2 * factor))
                        return f"#{r2:02x}{g2:02x}{b2:02x}"
                except Exception:
                    pass
                return raw
        except Exception:
            pass
        return None

    def _color_from_solid_elem(solid_elem):
        _SCHEME_MAP = {"bg1": "lt1", "tx1": "dk1", "bg2": "lt2", "tx2": "dk2"}
        for _c in solid_elem:
            _ctag = _c.tag.split("}")[-1]
            if _ctag == "srgbClr":
                _v = _c.get("val", "")
                return f"#{_v.lower()}" if len(_v) == 6 else None
            elif _ctag == "schemeClr":
                _key = _SCHEME_MAP.get(_c.get("val", ""), _c.get("val", ""))
                return theme_colors.get(_key)
            elif _ctag == "sysClr":
                _v = _c.get("lastClr", "")
                return f"#{_v.lower()}" if len(_v) == 6 else None
        return None

    def _get_fill(fill_obj, shape_element=None):
        try:
            if shape_element is not None:
                _sp = shape_element.find(_qn("p:spPr"))
                if _sp is not None:
                    _gf = _sp.find(".//" + _qn("a:gradFill"))
                    if _gf is not None:
                        _lin = _gf.find(_qn("a:lin"))
                        _ang = int(_lin.get("ang", "0")) / 60000.0 if _lin is not None else 0.0
                        _stops = []
                        for _gs in sorted(_gf.findall(".//" + _qn("a:gs")), key=lambda g: int(g.get("pos", "0"))):
                            _pos = int(_gs.get("pos", "0")) / 100000.0
                            _c = _color_from_solid_elem(_gs)
                            if _c:
                                _stops.append({"pos": _pos, "color": _c})
                        if _stops:
                            return {
                                "type": "gradient",
                                "color": _stops[0]["color"],
                                "gradient": {"angle": _ang, "stops": _stops},
                            }
            ft = str(fill_obj.type)
            if ft == "None" or "BACKGROUND" in ft:
                return {"type": "none"}
            if "SOLID" in ft or ft == "1":
                c = _resolve_color(fill_obj.fore_color)
                opacity = 1.0
                try:
                    if shape_element is not None:
                        _sp = shape_element.find(_qn("p:spPr"))
                        _solid = _sp.find(".//" + _qn("a:solidFill")) if _sp is not None else None
                        if _solid is not None:
                            for _child in _solid:
                                _alpha = _child.find(f"{{{_NS_A_FULL}}}alpha")
                                if _alpha is not None:
                                    opacity = int(_alpha.get("val", "100000")) / 100000.0
                                    break
                except Exception:
                    pass
                result = {"type": "solid", "color": c}
                if opacity < 0.999:
                    result["opacity"] = round(opacity, 4)
                return result
        except Exception:
            pass
        return {"type": "none"}

    def _get_bg(slide):
        _SCHEME_MAP = {"bg1": "lt1", "tx1": "dk1", "bg2": "lt2", "tx2": "dk2"}

        def _resolve_scheme(val):
            return theme_colors.get(_SCHEME_MAP.get(val, val))

        def _color_from_node(node):
            tag = node.tag.split("}")[-1]
            if tag == "srgbClr":
                v = node.get("val", "")
                return f"#{v.lower()}" if len(v) == 6 else None
            if tag == "schemeClr":
                return _resolve_scheme(node.get("val", ""))
            if tag == "sysClr":
                v = node.get("lastClr", "")
                return f"#{v.lower()}" if len(v) == 6 else None
            return None

        def _fill_from_csld(cSld_elem):
            bg = cSld_elem.find(_qn("p:bg"))
            if bg is None:
                return None, None
            bgPr = bg.find(_qn("p:bgPr"))
            if bgPr is not None:
                solid = bgPr.find(_qn("a:solidFill"))
                if solid is not None:
                    for child in solid:
                        c = _color_from_node(child)
                        if c:
                            return c, None
                grad = bgPr.find(_qn("a:gradFill"))
                if grad is not None:
                    lin = grad.find(_qn("a:lin"))
                    angle_deg = int(lin.get("ang", "0")) / 60000.0 if lin is not None else 0.0
                    gs_list = sorted(grad.findall(".//" + _qn("a:gs")), key=lambda g: int(g.get("pos", "0")))
                    stops = []
                    for gs in gs_list:
                        pos = int(gs.get("pos", "0")) / 100000.0
                        for child in gs:
                            c = _color_from_node(child)
                            if c:
                                stops.append({"pos": pos, "color": c})
                                break
                    if stops:
                        fallback = stops[0]["color"]
                        return fallback, {"angle": angle_deg, "stops": stops}
                return None, None
            bgRef = bg.find(_qn("p:bgRef"))
            if bgRef is not None:
                for child in bgRef:
                    c = _color_from_node(child)
                    if c:
                        return c, None
            return None, None

        try:
            fd = _get_fill(slide.background.fill)
            if fd["type"] == "solid" and fd.get("color"):
                return fd["color"], None
        except Exception:
            pass

        for src in [slide, slide.slide_layout, slide.slide_layout.slide_master]:
            try:
                c, grad = _fill_from_csld(src.background._element)
                if c:
                    return c, grad
            except Exception:
                pass

        return "#ffffff", None

    def _extract_geometry(shape):
        st, cr = "rect", 0
        try:
            sp_pr = shape._element.find(_qn("p:spPr"))
            if sp_pr is not None:
                pg = sp_pr.find(".//" + _qn("a:prstGeom"))
                if pg is not None:
                    prst = pg.get("prst", "rect")
                    st = prst
                    if prst == "roundRect":
                        cr = 33.334
                        av = pg.find(_qn("a:avLst"))
                        if av is not None:
                            for gd in av:
                                fmla = gd.get("fmla", "")
                                if fmla.startswith("val "):
                                    try:
                                        cr = int(fmla.split()[1]) / 500
                                    except Exception:
                                        pass
        except Exception:
            pass
        return st, cr

    def _get_layout_ph_defaults(shape):
        defaults: dict = {}
        try:
            ph_fmt = shape.placeholder_format
            if ph_fmt is None:
                return defaults
            layout = shape.part.slide.slide_layout
            for lshp in layout.shapes:
                try:
                    lpf = lshp.placeholder_format
                    if lpf is None or lpf.idx != ph_fmt.idx:
                        continue
                    sp = lshp._element
                    tx = sp.find(_qn("p:txBody"))
                    if tx is None:
                        break
                    bpr = tx.find(_qn("a:bodyPr"))
                    if bpr is not None:
                        anchor = bpr.get("anchor", "")
                        if anchor:
                            defaults["textAnchor"] = anchor
                    lst = tx.find(_qn("a:lstStyle"))
                    if lst is not None:
                        lvl1 = lst.find(_qn("a:lvl1pPr"))
                        if lvl1 is not None:
                            drpr = lvl1.find(_qn("a:defRPr"))
                            if drpr is not None:
                                sz = drpr.get("sz")
                                if sz:
                                    defaults["fontSize"] = int(sz) // 100
                                if drpr.get("b") == "1":
                                    defaults["bold"] = True
                                if drpr.get("cap", "") == "all":
                                    defaults["allCaps"] = True
                                solid = drpr.find(_qn("a:solidFill"))
                                if solid is not None:
                                    c = _color_from_solid_elem(solid)
                                    if c:
                                        defaults["color"] = c
                    break
                except Exception:
                    continue
        except Exception:
            pass
        return defaults

    def _iter_shapes_flat(shapes):
        for s in shapes:
            if s.shape_type == 6:  # GROUP
                try:
                    yield from _iter_shapes_flat(s.shapes)
                except Exception:
                    pass
            else:
                yield s

    slides_out = []
    for slide in prs.slides:
        bg_color = "#ffffff"
        bg_gradient = None
        try:
            bg_color, bg_gradient = _get_bg(slide)
        except Exception:
            pass

        boxes: list[dict] = []

        # Layout non-placeholder shapes
        try:
            for lshp in slide.slide_layout.shapes:
                try:
                    if lshp.shape_type == 14:
                        continue
                    lx = int((lshp.left or 0) * sx)
                    ly = int((lshp.top or 0) * sy)
                    sw_emu = lshp.width or 0
                    sh_emu = lshp.height or 0
                    if lshp.shape_type == 9:  # connector
                        ln_w_px = max(1, int(25400 * sx))
                        ln_color = "#888888"
                        spPr2 = lshp._element.find(_qn("p:spPr"))
                        if spPr2 is not None:
                            ln_el = spPr2.find(_qn("a:ln"))
                            if ln_el is not None:
                                ln_w_px = max(1, int(int(ln_el.get("w", "25400")) * sx))
                                sol2 = ln_el.find(_qn("a:solidFill"))
                                if sol2 is not None:
                                    c2 = _color_from_solid_elem(sol2)
                                    if c2:
                                        ln_color = c2
                        if sw_emu == 0:
                            bw, bh = ln_w_px, max(1, int(sh_emu * sy))
                            bx = lx - ln_w_px // 2
                            by_ = ly
                        elif sh_emu == 0:
                            bw, bh = max(1, int(sw_emu * sx)), ln_w_px
                            bx = lx
                            by_ = ly - ln_w_px // 2
                        else:
                            bx, by_ = lx, ly
                            bw = max(1, int(sw_emu * sx))
                            bh = max(1, int(sh_emu * sy))
                        boxes.append(
                            {
                                "id": _sec2.token_hex(4),
                                "x": bx,
                                "y": by_,
                                "w": bw,
                                "h": bh,
                                "rotation": 0,
                                "shapeType": "rect",
                                "text": "",
                                "styles": [],
                                "boxStyle": {"bgColor": ln_color},
                            }
                        )
                    elif lshp.shape_type == 1:  # rectangle
                        _fd2: dict = {"type": "none"}
                        _grad2 = None
                        _spPr2 = lshp._element.find(_qn("p:spPr"))
                        if _spPr2 is not None:
                            _sol2 = _spPr2.find(".//" + _qn("a:solidFill"))
                            if _sol2 is not None:
                                _c2 = _color_from_solid_elem(_sol2)
                                if _c2:
                                    _fd2 = {"type": "solid", "color": _c2}
                            _gf2 = _spPr2.find(".//" + _qn("a:gradFill"))
                            if _gf2 is not None and _fd2["type"] == "none":
                                _lin2 = _gf2.find(_qn("a:lin"))
                                _ang2 = int(_lin2.get("ang", "0")) / 60000.0 if _lin2 is not None else 0.0
                                _stops2: list = []
                                for _gs2 in sorted(
                                    _gf2.findall(".//" + _qn("a:gs")), key=lambda g: int(g.get("pos", "0"))
                                ):
                                    _pos2 = int(_gs2.get("pos", "0")) / 100000.0
                                    _c2 = _color_from_solid_elem(_gs2)
                                    if _c2:
                                        _stops2.append({"pos": _pos2, "color": _c2})
                                if _stops2:
                                    _fd2 = {"type": "solid", "color": _stops2[0]["color"]}
                                    _grad2 = {"angle": _ang2, "stops": _stops2}
                        if _fd2.get("type") == "solid" or _grad2:
                            _bs2: dict = {"bgColor": _fd2.get("color", "transparent")}
                            if _grad2:
                                _bs2["bgGradient"] = _grad2
                            boxes.append(
                                {
                                    "id": _sec2.token_hex(4),
                                    "x": lx,
                                    "y": ly,
                                    "w": max(1, int(sw_emu * sx)),
                                    "h": max(1, int(sh_emu * sy)),
                                    "rotation": 0,
                                    "shapeType": "rect",
                                    "text": "",
                                    "styles": [],
                                    "boxStyle": _bs2,
                                }
                            )
                except Exception:
                    pass
        except Exception:
            pass

        for shape in _iter_shapes_flat(slide.shapes):
            try:
                box_id = _sec2.token_hex(4)
                x = int((shape.left or 0) * sx)
                y = int((shape.top or 0) * sy)
                w = max(1, int((shape.width or 100) * sx))
                h = max(1, int((shape.height or 50) * sy))
                rotation = float(getattr(shape, "rotation", 0) or 0)

                _is_picture = shape.shape_type == 13
                if not _is_picture and not shape.has_text_frame:
                    try:
                        _blob_test = shape.image.blob  # noqa: F841
                        _is_picture = True
                    except Exception:
                        pass
                if _is_picture:
                    try:
                        img_bytes = shape.image.blob
                        ct = shape.image.content_type or "image/png"
                        img_b64 = _b64.b64encode(img_bytes).decode()
                        boxes.append(
                            {
                                "id": box_id,
                                "x": x,
                                "y": y,
                                "w": w,
                                "h": h,
                                "rotation": rotation,
                                "imageData": f"data:{ct};base64,{img_b64}",
                                "shapeType": "rect",
                                "text": "",
                                "styles": [],
                                "boxStyle": {},
                            }
                        )
                    except Exception:
                        pass
                    continue

                if not shape.has_text_frame:
                    _st, _cr = _extract_geometry(shape)
                    _fd: dict = {"type": "none"}
                    try:
                        _fd = _get_fill(shape.fill, shape._element)
                    except Exception:
                        pass
                    if _fd.get("type") != "none" or _st != "rect":
                        _bs: dict = {"bgColor": "transparent"}
                        if _fd.get("type") in ("solid", "gradient") and _fd.get("color"):
                            _bs["bgColor"] = _fd["color"]
                        if _fd.get("type") == "gradient" and _fd.get("gradient"):
                            _bs["bgGradient"] = _fd["gradient"]
                        try:
                            _spPr = shape._element.find(_qn("p:spPr"))
                            if _spPr is not None:
                                _ln = _spPr.find(_qn("a:ln"))
                                if _ln is not None and _ln.find(_qn("a:noFill")) is None:
                                    _ln_w_emu = int(_ln.get("w", "0") or "0")
                                    if _ln_w_emu > 0:
                                        _ln_sol = _ln.find(_qn("a:solidFill"))
                                        _ln_color = "#000000"
                                        if _ln_sol is not None:
                                            _lc = _color_from_solid_elem(_ln_sol)
                                            if _lc:
                                                _ln_color = _lc
                                        _bs["borderColor"] = _ln_color
                                        _bs["borderWidth"] = round(_ln_w_emu / 12700, 1)
                        except Exception:
                            pass
                        boxes.append(
                            {
                                "id": box_id,
                                "x": x,
                                "y": y,
                                "w": w,
                                "h": h,
                                "rotation": rotation,
                                "shapeType": _st,
                                "cornerRadius": _cr,
                                "fill": _fd,
                                "text": "",
                                "styles": [],
                                "boxStyle": _bs,
                            }
                        )
                    continue

                shape_type, corner_radius = _extract_geometry(shape)
                fill_dict: dict = {"type": "none"}
                try:
                    fill_dict = _get_fill(shape.fill, shape._element)
                except Exception:
                    pass

                tf = shape.text_frame
                paras_out = []
                box_style: dict = {"bgColor": "transparent"}
                first_done = False

                for para in tf.paragraphs:
                    align_str = "left"
                    try:
                        from pptx.enum.text import PP_ALIGN  # type: ignore[import-untyped]

                        al = para.alignment
                        if al == PP_ALIGN.CENTER:
                            align_str = "center"
                        elif al == PP_ALIGN.RIGHT:
                            align_str = "right"
                    except Exception:
                        pass

                    space_before = 0
                    try:
                        sb = para.space_before
                        if sb is not None:
                            space_before = round(float(sb) / 12700, 1)
                    except Exception:
                        pass

                    runs_out = []
                    para_text_parts = []
                    for run in para.runs:
                        rt = run.text or ""
                        if not rt:
                            continue
                        para_text_parts.append(rt)
                        rs: dict = {"text": rt}
                        try:
                            if run.font.size:
                                rs["fontSize"] = int(run.font.size / 12700)
                        except Exception:
                            pass
                        try:
                            if run.font.bold:
                                rs["bold"] = True
                        except Exception:
                            pass
                        try:
                            if run.font.italic:
                                rs["italic"] = True
                        except Exception:
                            pass
                        try:
                            if run.font.underline:
                                rs["underline"] = True
                        except Exception:
                            pass
                        try:
                            fn = run.font.name
                            if not fn:
                                _rpr = run._r.find(_qn("a:rPr"))
                                if _rpr is not None:
                                    _lat = _rpr.find(_qn("a:latin"))
                                    if _lat is not None:
                                        _tf = _lat.get("typeface", "")
                                        if _tf and not _tf.startswith("+"):
                                            fn = _tf
                            if fn:
                                rs["fontFamily"] = fn
                        except Exception:
                            pass
                        try:
                            _rpr2 = run._r.find(_qn("a:rPr"))
                            if _rpr2 is not None:
                                _spc = _rpr2.get("spc")
                                if _spc is not None:
                                    rs["letterSpacing"] = int(_spc) / 100
                        except Exception:
                            pass
                        try:
                            c = _resolve_color(run.font.color)
                            if c:
                                rs["color"] = c
                        except Exception:
                            pass
                        runs_out.append(rs)

                        if not first_done:
                            if "fontSize" in rs:
                                box_style["fontSize"] = rs["fontSize"]
                            if rs.get("bold"):
                                box_style["fontWeight"] = "bold"
                            if rs.get("italic"):
                                box_style["fontStyle"] = "italic"
                            if "color" in rs:
                                box_style["color"] = rs["color"]
                            if "fontFamily" in rs:
                                box_style["fontFamily"] = rs["fontFamily"]
                            box_style["textAlign"] = align_str
                            first_done = True

                    para_text = "".join(para_text_parts)
                    if para_text or runs_out:
                        paras_out.append(
                            {"text": para_text, "align": align_str, "spaceBefore": space_before, "runs": runs_out}
                        )

                if fill_dict.get("type") in ("solid", "gradient") and fill_dict.get("color"):
                    box_style["bgColor"] = fill_dict["color"]
                if fill_dict.get("type") == "gradient" and fill_dict.get("gradient"):
                    box_style["bgGradient"] = fill_dict["gradient"]

                try:
                    _spPr_t = shape._element.find(_qn("p:spPr"))
                    if _spPr_t is not None:
                        _ln_t = _spPr_t.find(_qn("a:ln"))
                        if _ln_t is not None and _ln_t.find(_qn("a:noFill")) is None:
                            _ln_w_emu_t = int(_ln_t.get("w", "0") or "0")
                            if _ln_w_emu_t > 0:
                                _ln_sol_t = _ln_t.find(_qn("a:solidFill"))
                                _ln_color_t = "#000000"
                                if _ln_sol_t is not None:
                                    _lct = _color_from_solid_elem(_ln_sol_t)
                                    if _lct:
                                        _ln_color_t = _lct
                                box_style["borderColor"] = _ln_color_t
                                box_style["borderWidth"] = round(_ln_w_emu_t / 12700, 1)
                except Exception:
                    pass

                try:
                    _txBody_t = shape._element.find(_qn("p:txBody"))
                    if _txBody_t is not None:
                        _bpr_t = _txBody_t.find(_qn("a:bodyPr"))
                        if _bpr_t is not None:
                            _anchor = _bpr_t.get("anchor", "")
                            if _anchor:
                                box_style["textAnchor"] = _anchor
                            _l = int(_bpr_t.get("lIns", "91440") or "91440")
                            _r = int(_bpr_t.get("rIns", "91440") or "91440")
                            _t = int(_bpr_t.get("tIns", "45720") or "45720")
                            _b = int(_bpr_t.get("bIns", "45720") or "45720")
                            box_style["padL"] = round(_l * sx, 2)
                            box_style["padR"] = round(_r * sx, 2)
                            box_style["padT"] = round(_t * sy, 2)
                            box_style["padB"] = round(_b * sy, 2)
                            if _bpr_t.find(_qn("a:normAutofit")) is not None:
                                box_style["normAutofit"] = True
                except Exception:
                    pass

                try:
                    if shape.shape_type == 14:
                        ld = _get_layout_ph_defaults(shape)
                        if ld:
                            if "textAnchor" in ld:
                                box_style["textAnchor"] = ld["textAnchor"]
                            if "fontSize" in ld and "fontSize" not in box_style:
                                box_style["fontSize"] = ld["fontSize"]
                            if ld.get("bold") and "fontWeight" not in box_style:
                                box_style["fontWeight"] = "bold"
                            if ld.get("allCaps"):
                                box_style["allCaps"] = True
                            if "color" in ld and "color" not in box_style:
                                box_style["color"] = ld["color"]
                except Exception:
                    pass

                full_text = "\n".join(pa["text"] for pa in paras_out)
                boxes.append(
                    {
                        "id": box_id,
                        "x": x,
                        "y": y,
                        "w": w,
                        "h": h,
                        "rotation": rotation,
                        "shapeType": shape_type,
                        "cornerRadius": corner_radius,
                        "fill": fill_dict,
                        "text": full_text,
                        "paragraphs": paras_out,
                        "styles": [],
                        "boxStyle": box_style,
                    }
                )
            except Exception:
                pass

        slide_width_pt = int(sw / 12700)
        slide_out: dict = {"bgColor": bg_color, "boxes": boxes, "slideWidthPt": slide_width_pt}
        if bg_gradient:
            slide_out["bgGradient"] = bg_gradient
        slides_out.append(slide_out)

    return slides_out


def _slides_to_pptx(slides: list[dict[str, Any]], source_path: Path | None = None) -> bytes:
    """Rebuild .pptx from slide JSON. Copied verbatim from old app.py _pptx_export."""
    import base64 as _b64mod
    import io as _iomod

    from lxml import etree as _etree_save  # type: ignore[import-untyped]
    from pptx import Presentation  # type: ignore[import-untyped]
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # type: ignore[import-untyped]
    from pptx.enum.text import PP_ALIGN as _PP_ALIGN  # type: ignore[import-untyped]
    from pptx.oxml.ns import qn as _qn_save  # type: ignore[import-untyped]
    from pptx.util import Emu, Pt  # type: ignore[import-untyped]

    CW, CH = 896, 504

    def hex_to_rgb(hex_str: str) -> RGBColor | None:
        h = str(hex_str).lstrip("#")
        if len(h) == 6:
            try:
                return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
            except Exception:
                pass
        return None

    if source_path and source_path.exists():  # noqa: ASYNC240
        _orig = Presentation(str(source_path))
        sw, sh = _orig.slide_width, _orig.slide_height
    else:
        _orig = Presentation()
        sw = _orig.slide_width or Emu(9144000)
        sh = _orig.slide_height or Emu(5143500)
    sx = sw / CW if CW else 1
    sy = sh / CH if CH else 1
    prs = Presentation()
    prs.slide_width = sw
    prs.slide_height = sh
    blank_layout = prs.slide_layouts[6]

    _PRST_MAP = {
        "ellipse": "ellipse",
        "roundRect": "roundRect",
        "rect": "rect",
        "triangle": "triangle",
        "isoscelesTri": "triangle",
        "rtTriangle": "rtTriangle",
        "diamond": "diamond",
        "parallelogram": "parallelogram",
        "trapezoid": "trapezoid",
        "pentagon": "pentagon",
        "hexagon": "hexagon",
        "plus": "plus",
        "mathPlus": "mathPlus",
        "mathMinus": "mathMinus",
        "mathMultiply": "mathMultiply",
        "rightArrow": "rightArrow",
        "leftArrow": "leftArrow",
        "upArrow": "upArrow",
        "downArrow": "downArrow",
        "chevron": "chevron",
        "bentArrow": "bentArrow",
        "curvedRightArrow": "curvedRightArrow",
        "can": "can",
        "cube": "cube",
        "star4": "star4",
        "star5": "star5",
        "snip1Rect": "snip1Rect",
        "round2SameRect": "round2SameRect",
        "round2DiagRect": "round2DiagRect",
    }

    def _align(s: str):
        if s == "center":
            return _PP_ALIGN.CENTER
        if s == "right":
            return _PP_ALIGN.RIGHT
        return _PP_ALIGN.LEFT

    def _apply_prst_geom(sp: Any, st_str: str, corner_radius: float = 0) -> None:
        sp_el = sp._element
        sp_pr = sp_el.find(_qn_save("p:spPr"))
        if sp_pr is None:
            return
        pg = sp_pr.find(_qn_save("a:prstGeom"))
        if pg is None:
            return
        prst = _PRST_MAP.get(st_str, st_str)
        pg.set("prst", prst)
        av = pg.find(_qn_save("a:avLst"))
        if av is not None:
            for gd_el in av.findall(_qn_save("a:gd")):
                av.remove(gd_el)
        if st_str == "roundRect" and corner_radius > 0:
            adj_val = max(0, min(50000, int(corner_radius * 500)))
            if av is None:
                av = _etree_save.SubElement(pg, _qn_save("a:avLst"))
            gd_new = _etree_save.SubElement(av, _qn_save("a:gd"))
            gd_new.set("name", "adj")
            gd_new.set("fmla", f"val {adj_val}")

    def _apply_fill(sp: Any, fill_data: Any, box_style: dict) -> None:
        fill_type = fill_data.get("type") if isinstance(fill_data, dict) else None
        fill_color = fill_data.get("color") if isinstance(fill_data, dict) else None
        opacity = float(fill_data.get("opacity", 1.0)) if isinstance(fill_data, dict) else 1.0
        gradient = fill_data.get("gradient") if isinstance(fill_data, dict) else None
        bg_grad = box_style.get("bgGradient") if box_style else None
        if bg_grad:
            gradient = bg_grad
            fill_type = "gradient"
        if fill_type == "gradient" and gradient and gradient.get("stops"):
            try:
                sp_pr_g = sp._element.find(_qn_save("p:spPr"))
                if sp_pr_g is not None:
                    for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:blipFill", "a:pattFill"):
                        for el in sp_pr_g.findall(_qn_save(tag)):
                            sp_pr_g.remove(el)
                    gf = _etree_save.SubElement(sp_pr_g, _qn_save("a:gradFill"))
                    gsLst = _etree_save.SubElement(gf, _qn_save("a:gsLst"))
                    for stop in gradient["stops"]:
                        gs = _etree_save.SubElement(gsLst, _qn_save("a:gs"))
                        gs.set("pos", str(int(stop["pos"] * 100000)))
                        srgb = _etree_save.SubElement(gs, _qn_save("a:srgbClr"))
                        srgb.set("val", stop["color"].lstrip("#"))
                    lin = _etree_save.SubElement(gf, _qn_save("a:lin"))
                    lin.set("ang", str(int(gradient.get("angle", 0) * 60000)))
                    lin.set("scaled", "0")
            except Exception:
                pass
        elif fill_color:
            rgb3 = hex_to_rgb(fill_color)
            if rgb3:
                sp.fill.solid()
                sp.fill.fore_color.rgb = rgb3
                if opacity < 1.0:
                    alpha_val = max(0, min(100000, int(opacity * 100000)))
                    try:
                        sp_pr2 = sp._element.find(_qn_save("p:spPr"))
                        if sp_pr2 is not None:
                            clr_el = sp_pr2.find(".//" + _qn_save("a:srgbClr"))
                            if clr_el is not None:
                                alpha_el = _etree_save.SubElement(clr_el, _qn_save("a:alpha"))
                                alpha_el.set("val", str(alpha_val))
                    except Exception:
                        pass
        else:
            sp.fill.background()

    def _apply_line(sp: Any, box_style: dict) -> None:
        border_color = (box_style or {}).get("borderColor", "")
        border_width = float((box_style or {}).get("borderWidth", 0) or 0)
        if border_color and border_color != "transparent" and border_width > 0:
            try:
                rgb_ln = hex_to_rgb(border_color)
                if rgb_ln:
                    sp.line.color.rgb = rgb_ln
                    sp.line.width = Pt(border_width)
            except Exception:
                pass
        else:
            try:
                sp.line.fill.background()
            except Exception:
                pass

    def _apply_text(tf: Any, paragraphs: list, plain_text: str, box_style: dict, box_font: str) -> None:
        if paragraphs:
            for pi, para_data in enumerate(paragraphs):
                para = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
                try:
                    para.alignment = _align(para_data.get("align") or box_style.get("textAlign", "left"))
                except Exception:
                    pass
                try:
                    sb = para_data.get("spaceBefore", 0) or 0
                    if sb:
                        para.space_before = Pt(sb)
                except Exception:
                    pass
                for run_data in para_data.get("runs") or []:
                    rt = run_data.get("text", "")
                    if not rt:
                        continue
                    run = para.add_run()
                    run.text = rt
                    run.font.size = Pt(run_data.get("fontSize") or box_style.get("fontSize", 16))
                    run.font.name = run_data.get("fontFamily") or box_font
                    if run_data.get("bold"):
                        run.font.bold = True
                    if run_data.get("italic"):
                        run.font.italic = True
                    if run_data.get("underline"):
                        run.font.underline = True
                    c_hex = run_data.get("color") or box_style.get("color", "#000000")
                    rgb_r = hex_to_rgb(c_hex)
                    if rgb_r:
                        run.font.color.rgb = rgb_r
        else:
            box_align = box_style.get("textAlign", "left")
            for li, line in enumerate((plain_text or "").split("\n")):
                para = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                try:
                    para.alignment = _align(box_align)
                except Exception:
                    pass
                if not line:
                    continue
                run = para.add_run()
                run.text = line
                run.font.size = Pt(box_style.get("fontSize", 16))
                run.font.name = box_font
                if box_style.get("fontWeight") == "bold":
                    run.font.bold = True
                if box_style.get("fontStyle") == "italic":
                    run.font.italic = True
                if box_style.get("textDecoration") == "underline":
                    run.font.underline = True
                rgb = hex_to_rgb(box_style.get("color", "#000000"))
                if rgb:
                    run.font.color.rgb = rgb

    def _apply_text_anchor(tf: Any, box_style: dict) -> None:
        anchor = (box_style or {}).get("textAnchor", "")
        if anchor:
            try:
                body_pr = tf._txBody.find(_qn_save("a:bodyPr"))
                if body_pr is not None:
                    body_pr.set("anchor", anchor)
            except Exception:
                pass

    for slide_data in slides:
        slide = prs.slides.add_slide(blank_layout)

        # ── Slide background ──────────────────────────────────────────────────
        try:
            bg_grad = slide_data.get("bgGradient")
            if bg_grad and bg_grad.get("stops"):
                bg_pr = slide.background._element.find(_qn_save("p:bg"))
                if bg_pr is None:
                    bg_pr = _etree_save.SubElement(slide.background._element, _qn_save("p:bg"))
                bg_prop = bg_pr.find(_qn_save("p:bgPr"))
                if bg_prop is None:
                    bg_prop = _etree_save.SubElement(bg_pr, _qn_save("p:bgPr"))
                for tag in ("a:solidFill", "a:gradFill", "a:noFill"):
                    for el in bg_prop.findall(_qn_save(tag)):
                        bg_prop.remove(el)
                gf_bg = _etree_save.SubElement(bg_prop, _qn_save("a:gradFill"))
                gsLst_bg = _etree_save.SubElement(gf_bg, _qn_save("a:gsLst"))
                for stop in bg_grad["stops"]:
                    gs_bg = _etree_save.SubElement(gsLst_bg, _qn_save("a:gs"))
                    gs_bg.set("pos", str(int(stop["pos"] * 100000)))
                    srgb_bg = _etree_save.SubElement(gs_bg, _qn_save("a:srgbClr"))
                    srgb_bg.set("val", stop["color"].lstrip("#"))
                lin_bg = _etree_save.SubElement(gf_bg, _qn_save("a:lin"))
                lin_bg.set("ang", str(int(bg_grad.get("angle", 0) * 60000)))
                lin_bg.set("scaled", "0")
            else:
                bg = slide.background.fill
                bg.solid()
                rgb_bg = hex_to_rgb(slide_data.get("bgColor", "#ffffff"))
                if rgb_bg:
                    bg.fore_color.rgb = rgb_bg
        except Exception:
            pass

        for box in slide_data.get("boxes", []):
            try:
                x = Emu(int(box.get("x", 0) * sx))
                y = Emu(int(box.get("y", 0) * sy))
                w = Emu(max(1, int(box.get("w", 100) * sx)))
                h = Emu(max(1, int(box.get("h", 50) * sy)))
                rotation = float(box.get("rotation", 0) or 0)
                box_style = box.get("boxStyle", {})

                # ── Image box ────────────────────────────────────────────────
                img_data = box.get("imageData", "")
                if img_data and "," in img_data:
                    try:
                        _, b64part = img_data.split(",", 1)
                        img_bytes = _b64mod.b64decode(b64part)
                        pic = slide.shapes.add_picture(_iomod.BytesIO(img_bytes), x, y, w, h)
                        if rotation:
                            pic.rotation = rotation
                    except Exception:
                        pass
                    continue

                fill_data = box.get("fill")
                st_str = box.get("shapeType", "rect") or "rect"
                has_text = bool(box.get("text") or box.get("paragraphs"))
                paragraphs = box.get("paragraphs") or []
                box_font = box_style.get("fontFamily") or "Montserrat"

                # ── Geometric shape with text ─────────────────────────────────
                if has_text and st_str not in ("", "textbox", "rect"):
                    sp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
                    _apply_prst_geom(sp, st_str, box.get("cornerRadius", 0))
                    _apply_fill(sp, fill_data or {"type": "none"}, box_style)
                    _apply_line(sp, box_style)
                    if rotation:
                        sp.rotation = rotation
                    tf = sp.text_frame
                    tf.word_wrap = True
                    _apply_text(tf, paragraphs, box.get("text", ""), box_style, box_font)
                    _apply_text_anchor(tf, box_style)
                    continue

                # ── Pure geometric shape (no text) ────────────────────────────
                if fill_data and not has_text and st_str not in ("", "textbox"):
                    sp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
                    _apply_prst_geom(sp, st_str, box.get("cornerRadius", 0))
                    _apply_fill(sp, fill_data, box_style)
                    _apply_line(sp, box_style)
                    if rotation:
                        sp.rotation = rotation
                    if sp.has_text_frame:
                        sp.text_frame.text = ""
                    continue

                # ── Text box ──────────────────────────────────────────────────
                txBox = slide.shapes.add_textbox(x, y, w, h)
                if rotation:
                    txBox.rotation = rotation
                tf = txBox.text_frame
                tf.word_wrap = True
                _apply_text(tf, paragraphs, box.get("text", ""), box_style, box_font)
                _apply_text_anchor(tf, box_style)
                _apply_line(txBox, box_style)
                bg_hex2 = box_style.get("bgColor")
                if bg_hex2 and bg_hex2 != "transparent":
                    rgb2 = hex_to_rgb(bg_hex2)
                    if rgb2:
                        txBox.fill.solid()
                        txBox.fill.fore_color.rgb = rgb2
                elif box_style.get("bgGradient"):
                    _apply_fill(
                        txBox,
                        {
                            "type": "gradient",
                            "gradient": box_style["bgGradient"],
                            "color": box_style.get("bgColor", "#ffffff"),
                        },
                        box_style,
                    )
            except Exception:
                pass

    buf = _iomod.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class PptService:
    async def load(self, path: str) -> list[dict[str, Any]]:
        p = Path(path)

        def _do() -> list[dict[str, Any]]:
            sc = _read_sidecar(p)
            if sc is not None:
                return sc.get("slides", [])
            slides = _pptx_to_slides(p)
            _write_sidecar(p, {"slides": slides})
            return slides

        return await asyncio.to_thread(_do)

    async def save(self, path: str, slides: list[dict[str, Any]]) -> None:
        p = Path(path)

        def _do() -> None:
            _write_sidecar(p, {"slides": slides})
            try:
                binary = _slides_to_pptx(slides, source_path=p)
                p.write_bytes(binary)
            except Exception:
                pass

        await asyncio.to_thread(_do)

    def has_sidecar(self, path: str) -> bool:
        return _sidecar(Path(path)).exists()  # noqa: ASYNC240
