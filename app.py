#!/usr/bin/env python3
"""
BoltzAgent — FastAPI server.

Run (dev):
    uvicorn app:app --host localhost --port 18789 --reload

Run (production with built frontend):
    python app.py --bzcode ./bzcode --dist ./dist

All business logic is imported from server.py.  This file only defines the
FastAPI routes, Pydantic request models, and the WebSocket endpoint.
"""

from __future__ import annotations

import argparse
import asyncio
import json
import logging
import os
import sys
from contextlib import asynccontextmanager
from pathlib import Path
from typing import Any, Dict, List, Optional

import uvicorn
from fastapi import (
    APIRouter,
    BackgroundTasks,
    FastAPI,
    HTTPException,
    Path as FPath,
    Query,
    Request,
    WebSocket,
    WebSocketDisconnect,
)
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import (
    FileResponse,
    JSONResponse,
    StreamingResponse,
)
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel

# ── Import all business logic from server.py ──────────────────────────────────
# We deliberately import the module-level globals and helpers; only the
# aiohttp HTTP/WebSocket layer is replaced here.
from server import (
    BACKEND_VERSION,
    BOLTZHUB_API,
    BOLTZHUB_AUTH,
    SESSIONS_DIR,
    SERVER_DATA_DIR,
    _active_sessions,
    _batch_store,
    _boltzhub_token,
    _credentials_valid,
    _load_defaults,
    _load_index,
    _load_mode_config,
    _load_titles,
    _now,
    _read_app_config,
    _read_session_file,
    _running_sessions,
    _save_code,
    _save_default,
    _save_index,
    _save_title,
    _clear_default,
    _token_stats,
    _add_tokens,
    _widget_load,
    _widget_save,
    _widget_lock,
    _widget_path,
    _wj,
    _write_session_config,
    _CANVAS_ID_RE,
    _wj,
    _BatchItem,
    _WASess,
    _whatsapp_sessions,
    _whatsapp_lock,
    _bz_headers,
    _bz_auth,
    _write_app_config,
    _sync_env_oauth_client_id,
    _load_code,
    _code_path,
    handle_ws_client,
    _write_bzcode_credentials,
    _read_api_keys,
    agent_pool,
    # Canvas / widget helpers
    WIDGETS_DIR,
    CUSTOM_WIDGETS_DIR,
    _custom_widgets_dir,
    _canvas_file,
    # Document parsing
    _detect_and_parse,
    _blocks_to_docx,
    # Dev-server & cursor
    _cursor_store,
    _dev_servers,
    _find_free_port,
    # _add_frontend kept separate — called after app is built
)




# ── Pydantic request models ───────────────────────────────────────────────────

class AuthBody(BaseModel):
    accessToken: str
    refreshToken: Optional[str] = None
    expiresAt: Optional[int] = None
    authUrl: str = "https://boltzhub.com"

class ProxyBody(BaseModel):
    url: str
    method: str = "GET"
    headers: Dict[str, str] = {}
    body: Optional[str] = None

class CredentialBody(BaseModel):
    key: str
    value: str

class CanvasBody(BaseModel):
    widgets: Optional[List[Any]] = None

class WidgetBody(BaseModel):
    id: str
    code: Optional[str] = None
    class Config:
        extra = "allow"

class SeedWidgetsBody(BaseModel):
    widgets: Optional[List[Dict[str, Any]]] = None

class SessionTitleBody(BaseModel):
    title: str

class SetDefaultBody(BaseModel):
    cwd: str
    sessionId: str = ""

class CreateSessionBody(BaseModel):
    cwd: str = ""
    mode: str = "general"

class SearchBody(BaseModel):
    pass

class MkdirBody(BaseModel):
    parent: str
    name: str

class CreateAppBody(BaseModel):
    cwd: str = ""
    name: str
    description: Optional[str] = None
    visibility: str = "private"
    priceMonthly: Optional[float] = None
    buildCommand: Optional[str] = None

class PushBody(BaseModel):
    cwd: str = ""
    releaseNotes: Optional[str] = None
    versionNumber: str = "1.0.0"

class SyncBody(BaseModel):
    cwd: str = ""
    appId: Optional[str] = None

class CreateVersionBody(BaseModel):
    appId: str
    releaseNotes: str = ""
    versionNumber: str = "1.0.0"

class PublishBody(BaseModel):
    appId: str

class BatchRunBody(BaseModel):
    cwds: List[str]
    message: str
    sessions: Dict[str, str] = {}

class WidgetRowBody(BaseModel):
    row: Optional[Dict[str, Any]] = None
    rows: Optional[List[Dict[str, Any]]] = None

class WidgetUpdateBody(BaseModel):
    data: Dict[str, Any]

class WidgetExecBody(BaseModel):
    code: str

class WidgetSchemaBody(BaseModel):
    columns: List[Dict[str, Any]] = []

class WriteFileBody(BaseModel):
    path: str
    content: str

class WhatsAppStatus(BaseModel):
    pass

class FileRenameBody(BaseModel):
    path: str
    newName: str

class FileDuplicateBody(BaseModel):
    path: str

class WriteFileBody2(BaseModel):
    path: str
    content: str = ""

class CursorBody(BaseModel):
    path: str
    selStart: int = 0
    selEnd: int = 0

class DocPathBody(BaseModel):
    path: str

class DocSaveBody(BaseModel):
    path: str
    blocks: list

class ExcelSaveBody(BaseModel):
    path: str
    sheets: list = []

class ExcelPatchBody(BaseModel):
    path: str
    sheet: str = ""       # sheet name; empty = first sheet
    cells: dict = {}      # { ref: { f?, v?, s? } } in sidecar schema

class ExcelGridBody(BaseModel):
    path: str
    sheet: str = ""
    columnIndexToWidth: dict = {}
    rowIndexToHeight: dict = {}

class ExcelAddSheetBody(BaseModel):
    path: str
    sheetName: str = "Sheet2"

class ExcelRenameSheetBody(BaseModel):
    path: str
    oldName: str
    newName: str

class ExcelMergeBody(BaseModel):
    path: str
    sheet: str = ""
    mergedCells: list = []   # list of range strings e.g. ["A1:C2"]

class PptSaveBody(BaseModel):
    path: str
    slides: list = []

class DeployWidgetBody(BaseModel):
    sessionId: str = ""
    cwd: str = ""
    title: str = "Widget"
    code: str = ""
    w: int = 380
    h: int = 280
    x: Optional[int] = None
    y: Optional[int] = None
    initialData: list = []

class CustomWidgetCodeBody(BaseModel):
    code: str = ""

class LogoutBody(BaseModel):
    authUrl: str = "https://boltzhub.com"

class DevServerBody(BaseModel):
    cwd: str = ""

class BzHubSyncBody(BaseModel):
    cwd: str = ""
    appId: str = ""

class BzHubVersionBody(BaseModel):
    appId: str
    releaseNotes: str = ""
    versionNumber: str = "1.0.0"

class BzHubPublishBody(BaseModel):
    appId: str

# ── Lifespan (startup / shutdown) ────────────────────────────────────────────

@asynccontextmanager
async def lifespan(app: FastAPI):
    # ── Startup ───────────────────────────────────────────────────────────────
    app.state.db = None
    await agent_pool.start()

    # ── bzcode version detection (run once at startup) ────────────────────────
    app.state.bzcode_version = None
    app.state.bzcode_latest  = None
    try:
        import shutil as _shv, re as _rev, aiohttp as _aio_v
        _bzc = app.state.bzcode_path
        _bzc_resolved = _shv.which(_bzc) or _bzc
        if os.path.isfile(_bzc_resolved):
            _vp = await asyncio.create_subprocess_exec(
                _bzc_resolved, "--version",
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
            )
            _vout, _ = await asyncio.wait_for(_vp.communicate(), timeout=10)
            _vm = _rev.search(r'(\d+\.\d+[\.\d]*)', _vout.decode())
            if _vm:
                app.state.bzcode_version = _vm.group(1)
                print(f"[startup] bzcode version: {app.state.bzcode_version}", file=sys.stderr)
    except Exception as _ve:
        print(f"[startup] bzcode version check failed: {_ve}", file=sys.stderr)

    try:
        import aiohttp as _aio_l
        async with _aio_l.ClientSession() as _ls:
            async with _ls.get(
                "https://boltzhub.com/bz-appstore-api/v1/bzcode/latest",
                timeout=_aio_l.ClientTimeout(total=8),
                ssl=False,
            ) as _lr:
                if _lr.ok:
                    import re as _rel2
                    _ld = await _lr.json(content_type=None)
                    _lv = (_ld.get("version") or _ld.get("latestVersion")
                           or _ld.get("tag") or "")
                    if not _lv and isinstance(_ld, str):
                        _lv = _ld
                    _lm = _rel2.search(r'(\d+\.\d+[\.\d]*)', str(_lv))
                    if _lm:
                        app.state.bzcode_latest = _lm.group(1)
                        print(f"[startup] bzcode latest: {app.state.bzcode_latest}", file=sys.stderr)
    except Exception as _le:
        print(f"[startup] bzcode latest check failed: {_le}", file=sys.stderr)

    yield

    # ── Shutdown ──────────────────────────────────────────────────────────────
    await agent_pool.stop()


# ── PPTX binary export (shared by load + save) ───────────────────────────────

def _pptx_export(p, slides):
    """Rebuild the .pptx binary from slide JSON. Called on save and on load-from-sidecar."""
    import json as _json, tempfile, shutil as _shutil, os as _os, base64 as _b64mod, io as _iomod
    from pathlib import Path as _Path
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.oxml.ns import qn as _qn_save
    from lxml import etree as _etree_save

    CW, CH = 896, 504
    p = _Path(p)

    def hex_to_rgb(hex_str):
        h = str(hex_str).lstrip("#")
        if len(h) == 6:
            try:
                return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
            except Exception:
                pass
        return None

    if p.exists():
        _orig = Presentation(str(p))
        sw, sh = _orig.slide_width, _orig.slide_height
    else:
        _orig = Presentation()
        sw = _orig.slide_width or Emu(9144000)
        sh = _orig.slide_height or Emu(5143500)
    sx = sw / CW if CW else 1
    sy = sh / CH if CH else 1
    prs = Presentation()
    prs.slide_width = sw
    prs.slide_height = sh
    blank_layout = prs.slide_layouts[6]

    _PRST_MAP = {
        'ellipse': 'ellipse', 'roundRect': 'roundRect', 'rect': 'rect',
        'triangle': 'triangle', 'isoscelesTri': 'triangle',
        'rtTriangle': 'rtTriangle', 'diamond': 'diamond',
        'parallelogram': 'parallelogram', 'trapezoid': 'trapezoid',
        'pentagon': 'pentagon', 'hexagon': 'hexagon',
        'plus': 'plus', 'mathPlus': 'mathPlus',
        'mathMinus': 'mathMinus', 'mathMultiply': 'mathMultiply',
        'rightArrow': 'rightArrow', 'leftArrow': 'leftArrow',
        'upArrow': 'upArrow', 'downArrow': 'downArrow',
        'chevron': 'chevron', 'bentArrow': 'bentArrow',
        'curvedRightArrow': 'curvedRightArrow',
        'can': 'can', 'cube': 'cube',
        'star4': 'star4', 'star5': 'star5',
        'snip1Rect': 'snip1Rect',
        'round2SameRect': 'round2SameRect',
        'round2DiagRect': 'round2DiagRect',
    }

    from pptx.enum.text import PP_ALIGN as _PP_ALIGN
    def _align(s):
        if s == 'center': return _PP_ALIGN.CENTER
        if s == 'right':  return _PP_ALIGN.RIGHT
        return _PP_ALIGN.LEFT

    def _apply_prst_geom(sp, st_str, corner_radius=0):
        """Patch prstGeom prst attribute and re-apply corner radius adjust."""
        sp_el = sp._element
        sp_pr = sp_el.find(_qn_save('p:spPr'))
        if sp_pr is None:
            return
        pg = sp_pr.find(_qn_save('a:prstGeom'))
        if pg is None:
            return
        prst = _PRST_MAP.get(st_str, st_str)
        pg.set('prst', prst)
        av = pg.find(_qn_save('a:avLst'))
        if av is not None:
            for gd_el in av.findall(_qn_save('a:gd')):
                av.remove(gd_el)
        if st_str == 'roundRect' and corner_radius > 0:
            adj_val = max(0, min(50000, int(corner_radius * 500)))
            if av is None:
                av = _etree_save.SubElement(pg, _qn_save('a:avLst'))
            gd_new = _etree_save.SubElement(av, _qn_save('a:gd'))
            gd_new.set('name', 'adj')
            gd_new.set('fmla', f'val {adj_val}')

    def _apply_fill(sp, fill_data, box_style):
        """Apply solid or gradient fill to a shape."""
        fill_type = fill_data.get('type') if isinstance(fill_data, dict) else None
        fill_color = fill_data.get('color') if isinstance(fill_data, dict) else None
        opacity = float(fill_data.get('opacity', 1.0)) if isinstance(fill_data, dict) else 1.0
        gradient = fill_data.get('gradient') if isinstance(fill_data, dict) else None
        # Prefer boxStyle bgGradient if present
        bg_grad = box_style.get('bgGradient') if box_style else None
        if bg_grad:
            gradient = bg_grad
            fill_type = 'gradient'

        if fill_type == 'gradient' and gradient and gradient.get('stops'):
            try:
                sp_pr_g = sp._element.find(_qn_save('p:spPr'))
                if sp_pr_g is not None:
                    # Remove existing fill elements
                    for tag in ('a:solidFill', 'a:gradFill', 'a:noFill', 'a:blipFill', 'a:pattFill'):
                        for el in sp_pr_g.findall(_qn_save(tag)):
                            sp_pr_g.remove(el)
                    gf = _etree_save.SubElement(sp_pr_g, _qn_save('a:gradFill'))
                    gsLst = _etree_save.SubElement(gf, _qn_save('a:gsLst'))
                    for stop in gradient['stops']:
                        gs = _etree_save.SubElement(gsLst, _qn_save('a:gs'))
                        gs.set('pos', str(int(stop['pos'] * 100000)))
                        srgb = _etree_save.SubElement(gs, _qn_save('a:srgbClr'))
                        srgb.set('val', stop['color'].lstrip('#'))
                    lin = _etree_save.SubElement(gf, _qn_save('a:lin'))
                    ang_60k = int(gradient.get('angle', 0) * 60000)
                    lin.set('ang', str(ang_60k))
                    lin.set('scaled', '0')
            except Exception:
                pass
        elif fill_color:
            rgb3 = hex_to_rgb(fill_color)
            if rgb3:
                sp.fill.solid()
                sp.fill.fore_color.rgb = rgb3
                if opacity < 1.0:
                    alpha_val = max(0, min(100000, int(opacity * 100000)))
                    try:
                        sp_pr2 = sp._element.find(_qn_save('p:spPr'))
                        if sp_pr2 is not None:
                            clr_el = sp_pr2.find('.//' + _qn_save('a:srgbClr'))
                            if clr_el is not None:
                                alpha_el = _etree_save.SubElement(clr_el, _qn_save('a:alpha'))
                                alpha_el.set('val', str(alpha_val))
                    except Exception:
                        pass
        else:
            sp.fill.background()

    def _apply_line(sp, box_style):
        """Apply border/line from boxStyle, or remove it."""
        border_color = (box_style or {}).get('borderColor', '')
        border_width = float((box_style or {}).get('borderWidth', 0) or 0)
        if border_color and border_color != 'transparent' and border_width > 0:
            try:
                rgb_ln = hex_to_rgb(border_color)
                if rgb_ln:
                    sp.line.color.rgb = rgb_ln
                    sp.line.width = Pt(border_width)
            except Exception:
                pass
        else:
            try:
                sp.line.fill.background()
            except Exception:
                pass

    def _apply_text(tf, paragraphs, plain_text, box_style, box_font):
        """Populate a text frame from paragraphs or plain text."""
        if paragraphs:
            for pi, para_data in enumerate(paragraphs):
                para = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
                try:
                    para.alignment = _align(para_data.get("align") or box_style.get("textAlign", "left"))
                except Exception:
                    pass
                try:
                    sb = para_data.get("spaceBefore", 0) or 0
                    if sb:
                        para.space_before = Pt(sb)
                except Exception:
                    pass
                for run_data in (para_data.get("runs") or []):
                    rt = run_data.get("text", "")
                    if not rt:
                        continue
                    run = para.add_run()
                    run.text = rt
                    run.font.size = Pt(run_data.get("fontSize") or box_style.get("fontSize", 16))
                    run.font.name = run_data.get("fontFamily") or box_font
                    if run_data.get("bold"):
                        run.font.bold = True
                    if run_data.get("italic"):
                        run.font.italic = True
                    if run_data.get("underline"):
                        run.font.underline = True
                    c_hex = run_data.get("color") or box_style.get("color", "#000000")
                    rgb_r = hex_to_rgb(c_hex)
                    if rgb_r:
                        run.font.color.rgb = rgb_r
        else:
            box_align = box_style.get("textAlign", "left")
            for li, line in enumerate((plain_text or "").split("\n")):
                para = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                try:
                    para.alignment = _align(box_align)
                except Exception:
                    pass
                if not line:
                    continue
                run = para.add_run()
                run.text = line
                run.font.size = Pt(box_style.get("fontSize", 16))
                run.font.name = box_font
                if box_style.get("fontWeight") == "bold":
                    run.font.bold = True
                if box_style.get("fontStyle") == "italic":
                    run.font.italic = True
                if box_style.get("textDecoration") == "underline":
                    run.font.underline = True
                rgb = hex_to_rgb(box_style.get("color", "#000000"))
                if rgb:
                    run.font.color.rgb = rgb

    def _apply_text_anchor(tf, box_style):
        """Set vertical text anchor from boxStyle.textAnchor (OOXML: t/ctr/b)."""
        anchor = (box_style or {}).get('textAnchor', '')
        if anchor:
            try:
                body_pr = tf._txBody.find(_qn_save('a:bodyPr'))
                if body_pr is not None:
                    body_pr.set('anchor', anchor)
            except Exception:
                pass

    for slide_data in slides:
        slide = prs.slides.add_slide(blank_layout)

        # ── Slide background ──────────────────────────────────────────────────
        try:
            bg_grad = slide_data.get("bgGradient")
            if bg_grad and bg_grad.get('stops'):
                bg_pr = slide.background._element.find(_qn_save('p:bg'))
                if bg_pr is None:
                    bg_pr = _etree_save.SubElement(slide.background._element, _qn_save('p:bg'))
                bg_prop = bg_pr.find(_qn_save('p:bgPr'))
                if bg_prop is None:
                    bg_prop = _etree_save.SubElement(bg_pr, _qn_save('p:bgPr'))
                for tag in ('a:solidFill', 'a:gradFill', 'a:noFill'):
                    for el in bg_prop.findall(_qn_save(tag)):
                        bg_prop.remove(el)
                gf_bg = _etree_save.SubElement(bg_prop, _qn_save('a:gradFill'))
                gsLst_bg = _etree_save.SubElement(gf_bg, _qn_save('a:gsLst'))
                for stop in bg_grad['stops']:
                    gs_bg = _etree_save.SubElement(gsLst_bg, _qn_save('a:gs'))
                    gs_bg.set('pos', str(int(stop['pos'] * 100000)))
                    srgb_bg = _etree_save.SubElement(gs_bg, _qn_save('a:srgbClr'))
                    srgb_bg.set('val', stop['color'].lstrip('#'))
                lin_bg = _etree_save.SubElement(gf_bg, _qn_save('a:lin'))
                lin_bg.set('ang', str(int(bg_grad.get('angle', 0) * 60000)))
                lin_bg.set('scaled', '0')
            else:
                bg = slide.background.fill
                bg.solid()
                rgb = hex_to_rgb(slide_data.get("bgColor", "#ffffff"))
                if rgb:
                    bg.fore_color.rgb = rgb
        except Exception:
            pass

        for box in slide_data.get("boxes", []):
            try:
                x = Emu(int(box.get("x", 0) * sx))
                y = Emu(int(box.get("y", 0) * sy))
                w = Emu(max(1, int(box.get("w", 100) * sx)))
                h = Emu(max(1, int(box.get("h", 50) * sy)))
                rotation = float(box.get("rotation", 0) or 0)
                box_style = box.get("boxStyle", {})

                # ── Image box ────────────────────────────────────────────────
                img_data = box.get("imageData", "")
                if img_data and "," in img_data:
                    try:
                        _, b64part = img_data.split(",", 1)
                        img_bytes = _b64mod.b64decode(b64part)
                        pic = slide.shapes.add_picture(_iomod.BytesIO(img_bytes), x, y, w, h)
                        if rotation:
                            pic.rotation = rotation
                    except Exception:
                        pass
                    continue

                fill_data = box.get("fill")
                st_str = box.get("shapeType", "rect") or "rect"
                has_text = bool(box.get("text") or box.get("paragraphs"))
                paragraphs = box.get("paragraphs", [])
                box_font = box_style.get("fontFamily") or "Montserrat"

                # ── Geometric shape with text (shape + text frame) ────────────
                if has_text and st_str not in ("", "textbox", "rect"):
                    sp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
                    _apply_prst_geom(sp, st_str, box.get("cornerRadius", 0))
                    _apply_fill(sp, fill_data or {'type': 'none'}, box_style)
                    _apply_line(sp, box_style)
                    if rotation:
                        sp.rotation = rotation
                    tf = sp.text_frame
                    tf.word_wrap = True
                    _apply_text(tf, paragraphs, box.get("text", ""), box_style, box_font)
                    _apply_text_anchor(tf, box_style)
                    continue

                # ── Pure geometric shape (no text) ────────────────────────────
                if fill_data and not has_text and st_str not in ("", "textbox"):
                    sp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
                    _apply_prst_geom(sp, st_str, box.get("cornerRadius", 0))
                    _apply_fill(sp, fill_data, box_style)
                    _apply_line(sp, box_style)
                    if rotation:
                        sp.rotation = rotation
                    if sp.has_text_frame:
                        sp.text_frame.text = ""
                    continue

                # ── Text box ──────────────────────────────────────────────────
                txBox = slide.shapes.add_textbox(x, y, w, h)
                if rotation:
                    txBox.rotation = rotation
                tf = txBox.text_frame
                tf.word_wrap = True
                _apply_text(tf, paragraphs, box.get("text", ""), box_style, box_font)
                _apply_text_anchor(tf, box_style)
                _apply_line(txBox, box_style)
                bg_hex2 = box_style.get("bgColor")
                if bg_hex2 and bg_hex2 != "transparent":
                    rgb2 = hex_to_rgb(bg_hex2)
                    if rgb2:
                        txBox.fill.solid()
                        txBox.fill.fore_color.rgb = rgb2
                elif box_style.get('bgGradient'):
                    _apply_fill(txBox, {'type': 'gradient', 'gradient': box_style['bgGradient'],
                                        'color': box_style.get('bgColor', '#ffffff')}, box_style)
            except Exception:
                pass

    tmp_fd, tmp_name = tempfile.mkstemp(suffix=".pptx", dir=p.parent)
    try:
        _os.close(tmp_fd)
        prs.save(tmp_name)
        _shutil.move(tmp_name, str(p))
    except Exception:
        try:
            _os.unlink(tmp_name)
        except OSError:
            pass
        raise


# ── App factory ───────────────────────────────────────────────────────────────

_DEFAULT_BZ_HOME = "/usr/local/boltzbit"


class _TeeWriter:
    """Write to both the original stream and a log file simultaneously."""
    def __init__(self, stream, file_path: Path):
        self._stream = stream
        try:
            self._file = open(file_path, 'a', buffering=1, encoding='utf-8', errors='replace')
        except Exception:
            self._file = None

    def write(self, data: str) -> int:
        if self._file:
            try:
                self._file.write(data)
            except Exception:
                pass
        return self._stream.write(data)

    def flush(self):
        self._stream.flush()
        if self._file:
            try:
                self._file.flush()
            except Exception:
                pass

    def fileno(self):
        return self._stream.fileno()

    def __getattr__(self, name):
        return getattr(self._stream, name)


def create_app(bzcode_path: str = "", default_cwd: str = "",
               bz_home: str = "", port: int = 18789) -> FastAPI:

    app = FastAPI(
        title="BoltzAgent API",
        version="1.0.0",
        description="bzcode bridge + widget canvas + session management",
        lifespan=lifespan,
    )

    bz_home = str(Path(bz_home or _DEFAULT_BZ_HOME).expanduser())
    os.makedirs(bz_home, exist_ok=True)

    # Store config accessible to route handlers
    app.state.bzcode_path = bzcode_path
    app.state.default_cwd = default_cwd
    app.state.bz_home     = bz_home
    app.state.port        = port

    # Tee all stderr to BZ_HOME/server.log so remote deployments have accessible logs.
    _log_path = Path(bz_home) / "server.log"
    if not isinstance(sys.stderr, _TeeWriter):
        sys.stderr = _TeeWriter(sys.stderr, _log_path)

    # Propagate BZ_HOME into the process environment so that _BatchItem and
    # _WASess (which inherit os.environ) also pick it up without needing refactoring.
    os.environ["BZ_HOME"] = bz_home

    # CORS — allow all origins (local-first tool)
    app.add_middleware(
        CORSMiddleware,
        allow_origins=["*"],
        allow_methods=["*"],
        allow_headers=["*"],
    )

    # ── Routers (one per domain — tags drive OpenAPI grouping) ────────────────
    ws_router        = APIRouter(tags=["WebSocket"])
    auth_router      = APIRouter(tags=["Auth"])
    files_router     = APIRouter(tags=["Files"])
    sessions_router  = APIRouter(tags=["Sessions"])
    canvas_router    = APIRouter(tags=["Canvas & Widgets"])
    db_router        = APIRouter(prefix="/db",        tags=["Database"])
    boltzhub_router  = APIRouter(prefix="/boltzhub",  tags=["BoltzHub"])
    batch_router     = APIRouter(tags=["Batch"])
    whatsapp_router  = APIRouter(prefix="/whatsapp",  tags=["WhatsApp"])
    misc_router      = APIRouter(tags=["Misc"])

    # ── § 1 · WebSocket bridge ────────────────────────────────────────────────

    @ws_router.websocket("/ws")
    async def ws_endpoint(websocket: WebSocket,
                          cwd: str = Query(""),
                          sessionId: str = Query(""),
                          mode: str = Query("")):
        """bzcode stdio bridge — one WebSocket per agent session.

        Uses AgentPool so the bzcode process outlives individual WebSocket
        connections and can be reconnected to with the same sessionId.
        """
        class _WsShim:
            """Thin shim: FastAPI WebSocket → aiohttp WebSocketResponse-like API."""
            async def send_str(self, text: str):
                await websocket.send_text(text)

            def __aiter__(self):
                return self

            async def __anext__(self):
                import aiohttp
                try:
                    text = await websocket.receive_text()
                    return type("M", (), {
                        "type": aiohttp.WSMsgType.TEXT,
                        "data": text,
                    })()
                except WebSocketDisconnect:
                    return type("M", (), {
                        "type": aiohttp.WSMsgType.CLOSE,
                        "data": "",
                    })()

        await websocket.accept()

        _bzcode  = app.state.bzcode_path
        _cwd     = app.state.default_cwd
        _bz_home = app.state.bz_home

        # Resolve cwd / sessionId / mode from query params
        effective_cwd = cwd if (cwd and os.path.isdir(cwd)) else _cwd
        req_mode = mode or _load_mode_config().get("default", "general")

        # Validate / generate session ID
        req_session_id = sessionId or None
        if req_session_id:
            # If there's already a live agent in the pool, skip the .jsonl check
            if req_session_id not in agent_pool._entries:
                if not (SESSIONS_DIR / f"{req_session_id}.jsonl").exists():
                    print(f"[ws] session file not found for {req_session_id!r} — starting fresh",
                          file=sys.stderr)
                    req_session_id = None

        if not req_session_id:
            import secrets as _sec
            req_session_id = f"bz-{_sec.token_hex(6)}"
            print(f"[ws] generated new sessionId={req_session_id}", file=sys.stderr)

        # ── Pre-flight checks ────────────────────────────────────────────────
        import shutil as _shutil_ws
        _bzcode_resolved = _shutil_ws.which(_bzcode) or _bzcode
        if not os.path.isfile(_bzcode_resolved):
            await websocket.send_text(json.dumps({
                "type": "result", "status": "error",
                "error": f"bzcode not found: '{_bzcode}' is not on PATH and not a file. Install bzcode or set BZCODE_PATH.",
            }))
            print(f"[ws] bzcode not found: {_bzcode}", file=sys.stderr)
            return
        _bzcode = _bzcode_resolved

        if not os.access(_bzcode, os.X_OK):
            try:
                os.chmod(_bzcode, 0o755)
                print(f"[ws] chmod +x {_bzcode} (was not executable)", file=sys.stderr)
            except OSError as e:
                await websocket.send_text(json.dumps({
                    "type": "result", "status": "error",
                    "error": f"bzcode is not executable: {_bzcode} — run: chmod +x {_bzcode}",
                }))
                print(f"[ws] cannot chmod bzcode: {e}", file=sys.stderr)
                return

        cred_ok, cred_reason = _credentials_valid()
        if not cred_ok:
            print(f"[ws] auth_error: {cred_reason}", file=sys.stderr)
            await websocket.send_text(json.dumps({
                "type": "auth_error",
                "reason": cred_reason,
            }))
            return

        # ── Session config & pool ────────────────────────────────────────────
        _existing = agent_pool._entries.get(req_session_id)
        _model_name = (_existing.model_info.get("name", "") if _existing else "")
        _write_session_config(req_session_id, req_mode, working_dir=effective_cwd, model_name=_model_name)
        cmd = [_bzcode, "--stdio", "--resume", req_session_id]
        api_keys = _read_api_keys()
        if not api_keys.get("BZ_API_KEY"):
            print("[ws] BZ_API_KEY not found — rejecting session spawn", file=sys.stderr)
            await websocket.send_text(json.dumps({
                "type": "result", "status": "error",
                "error": "No BZ_API_KEY configured. Please set it in Settings → AI API Key.",
            }))
            return
        env = {**os.environ, **api_keys, "BZ_PYTHON": sys.executable,
               **( {"BZ_HOME": _bz_home} if _bz_home else {} )}

        print(f"[ws] connect  cwd={effective_cwd}  sessionId={req_session_id}  mode={req_mode}",
              file=sys.stderr)
        _active_sessions.add(req_session_id)

        try:
            entry = await agent_pool.get_or_create(
                req_session_id, effective_cwd, req_mode,
                _bzcode, cmd, env)
        except (FileNotFoundError, PermissionError) as exc:
            _active_sessions.discard(req_session_id)
            await websocket.send_text(json.dumps({
                "type": "result", "status": "error",
                "error": f"Failed to start bzcode ({type(exc).__name__}): {exc}",
            }))
            return

        ws_shim = _WsShim()

        try:
            send_task, relay_task = await entry.attach_ws(ws_shim)
            # Wait until either the WS disconnects or the process dies
            done, pending = await asyncio.wait(
                [send_task, relay_task],
                return_when=asyncio.FIRST_COMPLETED,
            )
            for t in pending:
                t.cancel()
                try:
                    await t
                except (asyncio.CancelledError, Exception):
                    pass
        except (BrokenPipeError, ConnectionResetError, asyncio.CancelledError, WebSocketDisconnect):
            pass
        finally:
            await entry.detach_ws()
            _active_sessions.discard(req_session_id)
            if entry.is_dead:
                try:
                    await websocket.send_text(json.dumps({
                        "type": "system",
                        "message": f"⚠ bzcode exited unexpectedly (code {entry.proc.returncode}). Reconnecting…",
                    }))
                except Exception:
                    pass
            else:
                print(f"[ws] detached  sessionId={req_session_id}  pid={entry.proc.pid}",
                      file=sys.stderr)

    # ── § 2 · Auth & Credentials ─────────────────────────────────────────────

    @auth_router.post("/auth")
    async def auth(body: AuthBody):
        try:
            _write_bzcode_credentials(
                access_token=body.accessToken,
                refresh_token=body.refreshToken or "",
                expires_at=body.expiresAt,
                auth_url=body.authUrl,
            )
        except Exception as exc:
            print(f"[auth] failed to write credentials: {exc}", file=sys.stderr)
            raise HTTPException(status_code=500, detail="credential_write_failed")
        return {"ok": True}

    # ── § 10 · Misc (version + home live here; rest follows below) ───────────

    @misc_router.get("/api/version")
    async def api_version():
        return {
            "backend":       BACKEND_VERSION,
            "bzcode":        getattr(app.state, "bzcode_version", None),
            "bzcode_latest": getattr(app.state, "bzcode_latest",  None),
        }

    @misc_router.get("/api/apikey-status")
    async def apikey_status():
        keys_file = Path(app.state.bz_home) / "api_keys.json"
        if not keys_file.exists():
            return {"present": False, "last4": None}
        try:
            data = json.loads(keys_file.read_text())
            key = data.get("BZ_API_KEY", "")
            if key and isinstance(key, str):
                return {"present": True, "last4": key[-4:]}
        except Exception:
            pass
        return {"present": False, "last4": None}

    @misc_router.get("/api/user/me")
    async def user_me():
        """Fetch current user info from BoltzHub using the stored API key."""
        keys_file = Path(app.state.bz_home) / "api_keys.json"
        try:
            data = json.loads(keys_file.read_text())
            api_key = data.get("BZ_API_KEY", "")
        except Exception:
            api_key = ""
        if not api_key:
            return JSONResponse({"error": "no_key"}, status_code=401)
        import aiohttp as _aio_me
        try:
            async with _aio_me.ClientSession() as session:
                async with session.get(
                    "https://boltzhub.com/bz-appstore-api/v1/users/me",
                    headers={"X-API-Key": api_key},
                    timeout=_aio_me.ClientTimeout(total=8),
                ) as resp:
                    body = await resp.json()
                    if resp.status != 200:
                        return JSONResponse({"error": "upstream", "status": resp.status}, status_code=resp.status)
                    return {
                        "displayName": body.get("displayName") or body.get("username") or "",
                        "email": body.get("email") or "",
                        "username": body.get("username") or "",
                        "isVerified": body.get("isVerified", False),
                    }
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=502)

    @misc_router.post("/api/classify-mode")
    async def classify_mode(request: Request):
        """Classify a user message into one of the four agent modes."""
        body = await request.json()
        message = (body.get("message") or "").strip()
        if not message:
            return {"mode": "general"}

        keys_file = Path(app.state.bz_home) / "api_keys.json"
        try:
            data = json.loads(keys_file.read_text())
            api_key = data.get("BZ_API_KEY", "")
        except Exception:
            api_key = ""
        if not api_key:
            return {"mode": "general"}

        system_prompt = """\
You are a routing classifier for an AI agent. Given the user request, output the single best mode id.

Modes — pick the most specific match, use "general" only as a last resort:

widget  → building something visual and interactive: a clock, timer, stopwatch, calculator,
          to-do list, game, chart, form, canvas mini-app, data visualisation, or any
          self-contained browser UI component
worker  → file or document tasks: create/edit Excel, CSV, PDF, or Word files; extract,
          compare, or summarise document content; data processing, ETL, or automation scripts
coder   → software development on an existing project: write, debug, or refactor code;
          build APIs, backends, CLIs, services, or deploy applications
general → everything else: open-ended questions, explanations, research, writing prose,
          or tasks that produce no file, UI, or code artifact

Examples:
  "create a clock widget"               → widget
  "build a countdown timer"             → widget
  "make a to-do list app"               → widget
  "build a bar chart for sales data"    → widget
  "create an Excel file to sum numbers" → worker
  "extract data from this PDF"          → worker
  "summarise the attached report"       → worker
  "compare these two documents"         → worker
  "convert this CSV to JSON"            → worker
  "add auth to my Express app"          → coder
  "fix the bug in my Python script"     → coder
  "write unit tests for this module"    → coder
  "review my pull request"              → coder
  "explain how quicksort works"         → general
  "what are the pros and cons of X"     → general
  "write a poem about the ocean"        → general

Reply with ONLY one word: widget, worker, coder, or general.\
"""
        import aiohttp as _aio_cls
        payload = {
            "model": "anthropic-claude-4.5-sonnet",
            "messages": [{"role": "user", "content": message}],
            "stream": False,
            "system": system_prompt,
            "genOptions": {"maxTokens": 10, "temperature": 0},
        }
        try:
            async with _aio_cls.ClientSession() as session:
                async with session.post(
                    "https://flow.boltzbit.com/bz-api/v1/ai/messages",
                    headers={"X-API-KEY": api_key, "Content-Type": "application/json"},
                    json=payload,
                    timeout=_aio_cls.ClientTimeout(total=10),
                ) as resp:
                    if resp.status != 200:
                        return {"mode": "general"}
                    data = await resp.json()
                    text = ""
                    for block in data.get("content", []):
                        if block.get("type") == "text":
                            text = block.get("text", "").strip().lower()
                            break
                    valid = {"general", "widget", "worker", "coder"}
                    mode = text if text in valid else "general"
                    return {"mode": mode}
        except Exception:
            return {"mode": "general"}

    @misc_router.get("/api/home")
    async def api_home():
        home = str(Path.home())
        cwd = app.state.default_cwd or os.getcwd()
        return {
            "home": home,
            "defaultCwd": cwd if os.path.isdir(cwd) else home,
        }

    @misc_router.get("/api/apikey-verify")
    async def apikey_verify():
        """Verify BZ_API_KEY by spawning a fresh bzcode probe session."""
        api_keys = _read_api_keys()
        if not api_keys.get("BZ_API_KEY"):
            return {"status": "missing"}

        _bzcode = app.state.bzcode_path
        _bz_home = app.state.bz_home
        import shutil as _sh, secrets as _sec
        _bzcode = _sh.which(_bzcode) or _bzcode
        if not os.path.isfile(_bzcode):
            return {"status": "error", "reason": "bzcode not found"}

        # Fresh session ID every call — no stale context that could trigger compaction
        probe_id = f"bz-probe-{_sec.token_hex(4)}"
        env = {**os.environ, **api_keys, "BZ_PYTHON": sys.executable,
               "BZ_HOME": _bz_home}
        result_status = "unverified"
        result_reason = ""
        proc = None
        try:
            proc = await asyncio.create_subprocess_exec(
                _bzcode, "--stdio", "--resume", probe_id,
                stdin=asyncio.subprocess.PIPE,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.DEVNULL,
                env=env,
                limit=1024 * 1024,
            )

            got_idle = False
            deadline = asyncio.get_event_loop().time() + 20

            while asyncio.get_event_loop().time() < deadline:
                remaining = deadline - asyncio.get_event_loop().time()
                try:
                    line = await asyncio.wait_for(
                        proc.stdout.readline(), timeout=min(remaining, 2))
                except asyncio.TimeoutError:
                    continue  # keep waiting until overall deadline
                if not line:
                    break
                try:
                    msg = json.loads(line)
                except Exception:
                    continue

                mtype = msg.get("type", "")

                if mtype == "status" and msg.get("status") == "idle" and not got_idle:
                    got_idle = True
                    proc.stdin.write(b'{"type":"setMode","mode":"yolo"}\n')
                    probe = json.dumps({"type": "user", "content": "hi"}) + "\n"
                    proc.stdin.write(probe.encode())
                    await proc.stdin.drain()

                elif mtype == "status" and msg.get("status") == "running":
                    # bzcode accepted the prompt and is making the AI call — key is valid
                    result_status = "verified"
                    break

                elif mtype == "result":
                    err = msg.get("error", "")
                    if msg.get("status") == "error" and (
                            "403" in err or "not authorised" in err
                            or "unauthorized" in err.lower()):
                        result_status = "invalid"
                        result_reason = err[:300]
                    else:
                        result_status = "verified"
                    break

                elif mtype in ("assistant", "text"):
                    # Got an actual response — definitely verified
                    result_status = "verified"
                    break

        except Exception as exc:
            return {"status": "error", "reason": str(exc)}
        finally:
            if proc:
                try:
                    proc.kill()
                    await asyncio.wait_for(proc.wait(), timeout=3)
                except Exception:
                    pass
            # Clean up probe session files so probes don't appear in the session list
            import shutil as _probe_sh
            _probe_jsonl = Path(_bz_home) / "sessions" / f"{probe_id}.jsonl"
            _probe_dir   = Path(_bz_home) / "sessions" / probe_id
            for _p in (_probe_jsonl, _probe_dir):
                try:
                    if _p.is_dir():
                        _probe_sh.rmtree(_p, ignore_errors=True)
                    elif _p.exists():
                        _p.unlink()
                except Exception:
                    pass

        return {"status": result_status, "reason": result_reason}

    @misc_router.get("/api/server/log")
    async def server_log(lines: int = 200):
        bz_home_path = app.state.bz_home
        log_file = Path(bz_home_path) / "server.log"
        log_lines: list[str] = []
        if log_file.exists():
            try:
                text = log_file.read_text(encoding='utf-8', errors='replace')
                log_lines = text.splitlines()[-lines:]
            except Exception:
                pass
        return {
            "bzHome":   bz_home_path,
            "logFile":  str(log_file),
            "lines":    log_lines,
        }

    # ── Proxy (misc) ──────────────────────────────────────────────────────────

    @misc_router.post("/proxy")
    async def proxy(body: ProxyBody):
        import aiohttp, re as _re
        creds_file = SERVER_DATA_DIR / "credentials.json"
        creds: dict = {}
        if creds_file.exists():
            try:
                creds = json.loads(creds_file.read_text())
            except Exception:
                pass
        _ph = re.compile(r'\{\{(\w+)\}\}')
        def _resolve(text):
            return _ph.sub(lambda m: creds.get(m.group(1), m.group(0)), text)
        resolved_headers = {k: _resolve(str(v)) for k, v in body.headers.items()}
        resolved_body    = _resolve(body.body) if isinstance(body.body, str) else body.body
        if not body.url.startswith("http"):
            raise HTTPException(400, "url must start with http")
        try:
            connector = aiohttp.TCPConnector(ssl=False)
            async with aiohttp.ClientSession(connector=connector) as session:
                async with session.request(
                    body.method, body.url,
                    headers=resolved_headers, data=resolved_body,
                    allow_redirects=True,
                ) as resp:
                    content = await resp.read()
                    return JSONResponse(
                        content=json.loads(content) if resp.content_type == "application/json" else content.decode(errors="replace"),
                        status_code=resp.status,
                    )
        except Exception as exc:
            raise HTTPException(502, str(exc))

    # ── Credentials ───────────────────────────────────────────────────────────

    @auth_router.get("/credentials")
    async def get_credential_keys():
        f = SERVER_DATA_DIR / "credentials.json"
        if not f.exists():
            return {"keys": []}
        try:
            return {"keys": list(json.loads(f.read_text()).keys())}
        except Exception:
            return {"keys": []}

    @auth_router.post("/credentials")
    async def post_credential(body: CredentialBody):
        SERVER_DATA_DIR.mkdir(parents=True, exist_ok=True)
        f = SERVER_DATA_DIR / "credentials.json"
        data = json.loads(f.read_text()) if f.exists() else {}
        data[body.key] = body.value
        f.write_text(json.dumps(data, indent=2, ensure_ascii=False))
        return {"ok": True, "key": body.key}

    @auth_router.delete("/credentials/{key}")
    async def delete_credential(key: str):
        f = SERVER_DATA_DIR / "credentials.json"
        if not f.exists():
            raise HTTPException(404, "not found")
        data = json.loads(f.read_text())
        if key not in data:
            raise HTTPException(404, "not found")
        del data[key]
        f.write_text(json.dumps(data, indent=2, ensure_ascii=False))
        return {"ok": True, "deleted": key}

    @auth_router.get("/auth/status")
    async def auth_status():
        valid, reason = _credentials_valid()
        return {"valid": valid, "reason": reason}

    @auth_router.post("/auth/logout")
    async def logout(body: LogoutBody):
        bz_home = app.state.bz_home
        creds_file = Path(bz_home) / "credentials.json"
        try:
            if creds_file.exists():
                creds_file.unlink()
                print(f"[auth] credentials deleted from {creds_file}", file=sys.stderr)
            else:
                print(f"[auth] logout: no credentials file at {creds_file}", file=sys.stderr)
        except Exception as exc:
            print(f"[auth] logout error: {exc}", file=sys.stderr)
        return {"ok": True}

    @auth_router.post("/agent-key")
    async def set_api_key(request: Request):
        body = await request.json()
        key_name  = body.get("name", "")
        key_value = body.get("value", "")
        if key_name != "BZ_API_KEY":
            raise HTTPException(400, "Only BZ_API_KEY is accepted")
        if not key_value:
            raise HTTPException(400, "value is required")
        bz_home = app.state.bz_home
        keys_file = Path(bz_home) / "api_keys.json"
        keys: dict = {"BZ_API_KEY": key_value}  # overwrite; only BZ_API_KEY is stored
        with open(keys_file, "w") as f:
            json.dump(keys, f, indent=2)
        flushed = await agent_pool.flush_all(reason="api_key_reset")
        print(f"[api-key] BZ_API_KEY updated, flushed {flushed} session(s)", file=sys.stderr)
        return {"ok": True, "flushed": flushed}

    @auth_router.get("/agent-keys")
    async def list_api_keys():
        return {"keys": list(_read_api_keys().keys())}

    @auth_router.delete("/agent-key/{name}")
    async def delete_api_key(name: str):
        if name != "BZ_API_KEY":
            raise HTTPException(400, "Only BZ_API_KEY can be deleted")
        bz_home = app.state.bz_home
        keys_file = Path(bz_home) / "api_keys.json"
        if not keys_file.exists():
            return {"ok": True}
        try:
            with open(keys_file, "w") as f:
                json.dump({}, f)
            flushed = await agent_pool.flush_all(reason="api_key_deleted")
            print(f"[api-key] BZ_API_KEY deleted, flushed {flushed} session(s)", file=sys.stderr)
        except Exception as exc:
            raise HTTPException(500, str(exc))
        return {"ok": True}

    @auth_router.delete("/sessions-history")
    async def clear_sessions_history():
        """Delete all session JSONL files so every session starts with fresh context."""
        flushed = await agent_pool.flush_all(reason="history_cleared")
        deleted, errors = 0, 0
        if SESSIONS_DIR.exists():
            for f in SESSIONS_DIR.glob("*.jsonl"):
                try:
                    f.unlink()
                    deleted += 1
                except Exception as e:
                    print(f"[history] could not delete {f}: {e}", file=sys.stderr)
                    errors += 1
        print(f"[history] cleared {deleted} session files, flushed {flushed} processes", file=sys.stderr)
        return {"ok": True, "deleted": deleted, "flushed": flushed, "errors": errors}

    # ── § 3 · File System ─────────────────────────────────────────────────────

    @files_router.get("/shell")
    async def shell(cmd: str = Query(""), cwd: str = Query("")):
        if not cmd:
            raise HTTPException(400, "cmd is required")
        try:
            proc = await asyncio.create_subprocess_shell(
                cmd,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.STDOUT,
                cwd=cwd if os.path.isdir(cwd) else os.getcwd(),
            )
            stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=30)
            return {"output": stdout.decode(errors="replace"), "returncode": proc.returncode}
        except asyncio.TimeoutError:
            raise HTTPException(408, "Command timed out (30 s)")
        except Exception as exc:
            raise HTTPException(500, str(exc))

    @files_router.get("/files")
    async def list_files(path: str = Query("")):
        if not path:
            path = app.state.default_cwd or os.getcwd()
        p = Path(path)
        # Relative paths (server-stripped prefix) → reconstruct against parent of default_cwd
        if not p.is_absolute() and app.state.default_cwd:
            p = Path(app.state.default_cwd).parent / p
        if not p.exists() or not p.is_dir():
            raise HTTPException(404, "path not found or not a directory")
        entries = []
        for entry in sorted(p.iterdir(), key=lambda e: (e.is_file(), e.name.lower())):
            if entry.name.startswith('.'):
                continue  # hide dotfiles (including pptx / excel JSON sidecars)
            # Hide doc sidecars: {name}.{ext}.json where ext is a known document type
            _n = entry.name.lower()
            if _n.endswith('.json') and any(
                _n.endswith(s) for s in (
                    '.docx.json', '.doc.json', '.pdf.json',
                    '.html.json', '.htm.json',
                    '.md.json', '.markdown.json',
                )
            ):
                continue
            try:
                stat = entry.stat()
                entries.append({"name": entry.name, "path": str(entry),
                                 "isDir": entry.is_dir(), "size": stat.st_size,
                                 "modified": stat.st_mtime})
            except (PermissionError, OSError):
                pass
        return {"path": str(p), "entries": entries}

    @files_router.post("/files/mkdir")
    async def mkdir(body: MkdirBody):
        name = body.name.strip()
        if not body.parent or not name:
            raise HTTPException(400, "parent and name required")
        if "/" in name or "\\" in name or name in (".", ".."):
            raise HTTPException(400, "invalid folder name")
        new_dir = Path(body.parent) / name
        parent_path = Path(body.parent)
        if not parent_path.exists() or not parent_path.is_dir():
            raise HTTPException(400, f"parent directory not found: {body.parent}")
        if not os.access(parent_path, os.W_OK):
            raise HTTPException(403, f"no write permission on {body.parent}")
        try:
            new_dir.mkdir(parents=False, exist_ok=False)
            return {"path": str(new_dir)}
        except FileExistsError:
            raise HTTPException(409, "folder already exists")
        except PermissionError as exc:
            raise HTTPException(403, str(exc))
        except OSError as exc:
            raise HTTPException(500, str(exc))

    # ── § 5 · Canvas & Widgets ───────────────────────────────────────────────

    @canvas_router.get("/canvas")
    async def get_canvas(cwd: str = Query(""), sessionId: str = Query("")):
        f = _canvas_file(sessionId, cwd)
        if not f.exists():
            return {"widgets": []}
        return json.loads(f.read_text())

    @canvas_router.post("/canvas")
    async def post_canvas(request: Request, cwd: str = Query(""), sessionId: str = Query("")):
        body = await request.json()
        f = _canvas_file(sessionId, cwd)
        f.parent.mkdir(parents=True, exist_ok=True)
        f.write_text(json.dumps(body, indent=2, ensure_ascii=False))
        return {"ok": True, "file": str(f)}

    # ── Widgets ───────────────────────────────────────────────────────────────

    @canvas_router.get("/widgets")
    async def get_widgets():
        data = _load_index()
        return {"widgets": [
            {**e, "code": _load_code(e["id"])}
            for e in data.get("widgets", []) if not e.get("archived", False)
        ]}

    @canvas_router.get("/widgets/template")
    async def get_widget_template(name: str = Query(...)):
        safe = "".join(c if c.isalnum() or c in "-_" else "_" for c in name)
        p = WIDGETS_DIR / f"{safe}.js"
        if not p.exists():
            raise HTTPException(404, f"template not found: {name}")
        from fastapi.responses import Response as _Resp
        return _Resp(content=p.read_text(encoding="utf-8"), media_type="application/javascript")

    @canvas_router.get("/widgets/{widget_id}")
    async def get_widget(widget_id: str = FPath(...)):
        data = _load_index()
        entry = next((w for w in data.get("widgets", []) if w.get("id") == widget_id), None)
        if entry is None:
            raise HTTPException(404, "widget not found")
        return {**entry, "code": _load_code(widget_id)}

    @canvas_router.post("/widgets")
    async def post_widget(request: Request):
        body = await request.json()
        widget_id = body.get("id", "").strip()
        if not widget_id:
            raise HTTPException(400, "'id' is required")
        code = body.pop("code", None)
        data = _load_index()
        widgets = data.get("widgets", [])
        now = _now()
        idx = next((i for i, w in enumerate(widgets) if w.get("id") == widget_id), None)
        if idx is not None:
            entry = {**widgets[idx], **body, "updatedAt": now}
            entry.setdefault("createdAt", now)
            widgets[idx] = entry
        else:
            entry = {**body, "archived": False, "createdAt": now, "updatedAt": now}
            widgets.append(entry)
        data["widgets"] = widgets
        _save_index(data)
        if code is not None:
            _save_code(widget_id, code)
        return {**entry, "code": _load_code(widget_id)}

    @canvas_router.post("/widgets/seed")
    async def seed_widgets(request: Request):
        body = await request.json()
        incoming = body if isinstance(body, list) else (body.get("widgets") or [])
        data = _load_index()
        widgets = data.get("widgets", [])
        existing_ids = {w["id"] for w in widgets}
        now = _now()
        seeded = 0
        for w in incoming:
            wid = w.get("id", "")
            if not wid:
                continue
            code = w.pop("code", "")
            if wid in existing_ids:
                idx = next(i for i, x in enumerate(widgets) if x["id"] == wid)
                if not widgets[idx].get("isBuiltin", False):
                    continue
                widgets[idx] = {**widgets[idx], **w, "updatedAt": now}
                _save_code(wid, code)
                seeded += 1
            else:
                entry = {**w, "id": wid, "isBuiltin": True, "archived": False,
                         "createdAt": now, "updatedAt": now}
                widgets.append(entry)
                existing_ids.add(wid)
                _save_code(wid, code)
                seeded += 1
        data["widgets"] = widgets
        _save_index(data)
        print(f"[widgets] seeded {seeded} built-in widget(s)", file=sys.stderr)
        return {"seeded": seeded}

    @canvas_router.delete("/widgets/{widget_id}")
    async def delete_widget(widget_id: str = FPath(...)):
        data = _load_index()
        widgets = data.get("widgets", [])
        found = False
        for w in widgets:
            if w.get("id") == widget_id:
                w["archived"] = True
                w["updatedAt"] = _now()
                found = True
                break
        if not found:
            raise HTTPException(404, "widget not found")
        data["widgets"] = widgets
        _save_index(data)
        return {"ok": True, "archived": widget_id}

    @canvas_router.post("/canvas/deploy-widget")
    async def deploy_widget(body: DeployWidgetBody):
        import secrets as _sec, datetime as _dt
        if not body.sessionId and (not body.cwd or not os.path.isdir(body.cwd)):
            raise HTTPException(400, "sessionId or valid cwd required")
        if not body.code:
            raise HTTPException(400, "code is required")
        canvas_id = _sec.token_hex(5)
        widget_code_dir = _custom_widgets_dir(body.sessionId)
        widget_code_dir.mkdir(parents=True, exist_ok=True)
        (widget_code_dir / f"{canvas_id}.js").write_text(body.code, encoding="utf-8")
        if body.initialData:
            widget_data_dir = (SESSIONS_DIR / body.sessionId / "widget_data") if body.sessionId else (SERVER_DATA_DIR / "widget_data")
            widget_data_dir.mkdir(parents=True, exist_ok=True)
            records, next_id = [], 1
            for row in body.initialData:
                row = {k: v for k, v in row.items() if k not in ("id", "created_at")}
                row["id"] = next_id
                row["created_at"] = _dt.datetime.utcnow().isoformat() + "Z"
                records.append(row)
                next_id += 1
            (widget_data_dir / f"{canvas_id}.json").write_text(
                json.dumps({"_next_id": next_id, "records": records}, indent=2), encoding="utf-8")
        canvas_file = _canvas_file(body.sessionId, body.cwd)
        canvas_data: dict = {"version": 1, "widgets": []}
        if canvas_file.exists():
            try:
                canvas_data = json.loads(canvas_file.read_text(encoding="utf-8"))
            except Exception:
                pass
        existing = canvas_data.get("widgets", [])
        pad = 24
        x = body.x if body.x is not None else (pad if not existing else pad)
        y = body.y if body.y is not None else (pad if not existing else max((e.get("y", 0) + e.get("h", 0)) for e in existing) + pad)
        new_entry = {"canvasId": canvas_id, "widgetId": canvas_id, "kind": "custom",
                     "title": body.title, "x": x, "y": y, "w": body.w, "h": body.h}
        existing.append(new_entry)
        canvas_data["widgets"] = existing
        canvas_file.write_text(json.dumps(canvas_data, indent=2, ensure_ascii=False), encoding="utf-8")
        return {"ok": True, "canvasId": canvas_id, "widgetId": canvas_id,
                "title": body.title, "x": x, "y": y, "w": body.w, "h": body.h,
                "canvasFile": str(canvas_file)}

    @canvas_router.get("/custom-widgets/{canvas_id}")
    async def get_custom_widget(canvas_id: str = FPath(...), sessionId: str = Query("")):
        cwd_dir = _custom_widgets_dir(sessionId)
        p = cwd_dir / f"{canvas_id}.js"
        if not p.exists() and sessionId:
            p = CUSTOM_WIDGETS_DIR / f"{canvas_id}.js"
        if not p.exists():
            raise HTTPException(404, "not found")
        return {"canvasId": canvas_id, "code": p.read_text(encoding="utf-8")}

    @canvas_router.put("/custom-widgets/{canvas_id}")
    async def put_custom_widget(canvas_id: str = FPath(...), body: CustomWidgetCodeBody = CustomWidgetCodeBody(), sessionId: str = Query("")):
        if not canvas_id:
            raise HTTPException(400, "canvasId required")
        dest = _custom_widgets_dir(sessionId)
        dest.mkdir(parents=True, exist_ok=True)
        (dest / f"{canvas_id}.js").write_text(body.code, encoding="utf-8")
        return {"ok": True, "canvasId": canvas_id}

    @canvas_router.delete("/custom-widgets/{canvas_id}")
    async def delete_custom_widget(canvas_id: str = FPath(...), sessionId: str = Query("")):
        for d in [_custom_widgets_dir(sessionId), CUSTOM_WIDGETS_DIR]:
            p = d / f"{canvas_id}.js"
            if p.exists():
                p.unlink()
        return {"ok": True}

    # ── § 4 · Sessions ────────────────────────────────────────────────────────

    @sessions_router.get("/sessions")
    async def get_sessions(cwd: str = Query("")):
        if not SESSIONS_DIR.exists():
            return {"sessions": []}
        all_sessions = [m for p in SESSIONS_DIR.glob("*.jsonl") if (m := _read_session_file(p))]
        if cwd:
            all_sessions = [s for s in all_sessions if s["workingDir"] == cwd]
        sessions = sorted(all_sessions, key=lambda s: s["lastModified"], reverse=True)
        defaults = _load_defaults()
        for s in sessions:
            sid = s["sessionId"]
            wd  = s["workingDir"]
            s["isActive"]         = sid in _active_sessions
            s["isRunning"]        = sid in _running_sessions
            s["isDefault"]        = defaults.get(wd) == sid
            s["defaultSessionId"] = defaults.get(wd)
        # Strip absolute path prefix before sending to client — expose only the
        # path relative to the parent of default_cwd (e.g. "workspace/session-dir").
        _default_cwd = app.state.default_cwd or ""
        _path_base = str(Path(_default_cwd).parent) if _default_cwd else ""
        for s in sessions:
            wd = s["workingDir"]
            if _path_base and wd.startswith(_path_base + "/"):
                s["workingDir"] = wd[len(_path_base) + 1:]
        return {"sessions": sessions}

    @sessions_router.post("/session-default")
    async def set_default_session(body: SetDefaultBody):
        if not body.cwd:
            raise HTTPException(400, "cwd required")
        if body.sessionId:
            _save_default(body.cwd, body.sessionId)
        else:
            _clear_default(body.cwd)
        return {"ok": True}

    @sessions_router.delete("/sessions/{session_id}")
    async def delete_session(session_id: str = FPath(...)):
        if "/" in session_id or ".." in session_id:
            raise HTTPException(400, "invalid sessionId")
        p = SESSIONS_DIR / f"{session_id}.jsonl"
        if not p.exists():
            raise HTTPException(404, "not found")
        p.unlink()
        return {"ok": True}

    @sessions_router.post("/sessions/{session_id}/title")
    async def update_session_title(session_id: str, body: SessionTitleBody):
        if "/" in session_id or ".." in session_id:
            raise HTTPException(400, "invalid sessionId")
        _save_title(session_id, body.title[:100])
        return {"ok": True}

    @sessions_router.post("/sessions/create")
    async def create_session_with_handshake(body: CreateSessionBody):
        """Pre-create a session: write config, pre-warm bzcode in background, return immediately."""
        import secrets as _sec2
        effective_cwd = (body.cwd if body.cwd and os.path.isdir(body.cwd)
                         else app.state.default_cwd)
        mode = body.mode or _load_mode_config().get("default", "general")
        session_id = f"bz-{_sec2.token_hex(6)}"
        _write_session_config(session_id, mode, working_dir=effective_cwd)

        # Pre-warm: start bzcode now so it's ready when the browser redirects and connects.
        _bzcode = app.state.bzcode_path
        _bz_home = getattr(app.state, 'bz_home', None)
        cmd = [_bzcode, "--stdio", "--resume", session_id]
        env = {**os.environ, **_read_api_keys(), "BZ_PYTHON": sys.executable,
               **({"BZ_HOME": _bz_home} if _bz_home else {})}

        async def _prewarm():
            try:
                await agent_pool.get_or_create(session_id, effective_cwd, mode, _bzcode, cmd, env)
            except Exception as _e:
                print(f"[session] prewarm failed for {session_id}: {_e}", file=sys.stderr)

        asyncio.create_task(_prewarm())
        return {"ok": True, "sessionId": session_id}

    # ── Search / token-stats / agent-modes / settings (misc) ─────────────────

    @misc_router.get("/search")
    async def search(q: str = Query(""), key: str = Query(""), num: int = Query(10)):
        if not q:
            raise HTTPException(400, "q is required")
        if not key:
            raise HTTPException(400, "key is required")
        import aiohttp as _aio
        params = {"engine": "google", "q": q, "api_key": key, "num": num, "hl": "en", "gl": "us"}
        try:
            async with _aio.ClientSession() as session:
                async with session.get("https://serpapi.com/search.json", params=params) as resp:
                    body = await resp.json(content_type=None)
                    if not resp.ok:
                        raise HTTPException(resp.status, body.get("error", "SerpAPI error"))
        except HTTPException:
            raise
        except Exception as exc:
            raise HTTPException(502, str(exc))
        organic = body.get("organic_results", [])
        results = [{"title": r.get("title",""), "link": r.get("link",""),
                    "displayLink": r.get("displayed_link", r.get("link","")),
                    "snippet": r.get("snippet",""), "favicon": r.get("favicon",""),
                    "position": r.get("position", i+1)} for i, r in enumerate(organic)]
        return {"results": results, "meta": body.get("search_information", {})}

    # ── Token stats / agent modes ─────────────────────────────────────────────

    @misc_router.get("/token-stats")
    async def token_stats():
        return _token_stats

    @misc_router.get("/api/pool/status")
    async def pool_status():
        return {"agents": agent_pool.status()}

    @misc_router.post("/api/pool/connect")
    async def pool_connect(request: Request):
        """Create or reuse an agent. Returns session info with conversation history."""
        body = await request.json()
        req_cwd = body.get("cwd", "")
        req_session_id = body.get("sessionId", "")
        req_mode = body.get("mode", "") or _load_mode_config().get("default", "general")

        _bzcode = app.state.bzcode_path
        _cwd = app.state.default_cwd
        _bz_home = app.state.bz_home

        # Resolve effective_cwd from req_cwd:
        # - absolute path that still exists → use directly
        # - relative path (server stripped the prefix) → reconstruct against parent of default_cwd
        # - fallback to default_cwd
        if req_cwd and os.path.isabs(req_cwd) and os.path.isdir(req_cwd):
            effective_cwd = req_cwd
        elif req_cwd and not os.path.isabs(req_cwd) and _cwd:
            _resolved = os.path.join(str(Path(_cwd).parent), req_cwd)
            effective_cwd = _resolved if os.path.isdir(_resolved) else _cwd
        else:
            effective_cwd = _cwd

        # Validate / generate session ID
        if req_session_id:
            if req_session_id not in agent_pool._entries:
                if not (SESSIONS_DIR / f"{req_session_id}.jsonl").exists():
                    req_session_id = ""
            # For existing sessions, prefer the server-stored workingDir — but only
            # when it still exists on disk. If the folder was renamed the client-resolved
            # path above is more up-to-date, so we keep that instead.
            if req_session_id:
                _sf = SESSIONS_DIR / f"{req_session_id}.jsonl"
                if _sf.exists():
                    _saved = _read_session_file(_sf)
                    _stored_wd = _saved.get("workingDir") if _saved else None
                    if _stored_wd and os.path.isdir(_stored_wd):
                        effective_cwd = _stored_wd
                # Recover saved agent mode from meta.json for existing sessions
                _meta_file = SESSIONS_DIR / req_session_id / "meta.json"
                if _meta_file.exists():
                    try:
                        _meta = json.loads(_meta_file.read_text())
                        req_mode = _meta.get("mode", req_mode)
                    except Exception:
                        pass
        if not req_session_id:
            import secrets as _sec
            req_session_id = f"bz-{_sec.token_hex(6)}"

        # Pre-flight checks
        import shutil as _sh
        _bzcode_resolved = _sh.which(_bzcode) or _bzcode
        if not os.path.isfile(_bzcode_resolved):
            raise HTTPException(500, f"bzcode not found: {_bzcode}")
        _bzcode = _bzcode_resolved
        if not os.access(_bzcode, os.X_OK):
            try:
                os.chmod(_bzcode, 0o755)
            except OSError:
                raise HTTPException(500, f"bzcode not executable: {_bzcode}")

        cred_ok, cred_reason = _credentials_valid()
        if not cred_ok:
            raise HTTPException(401, cred_reason)

        _existing2 = agent_pool._entries.get(req_session_id)
        _model_name2 = (_existing2.model_info.get("name", "") if _existing2 else "")
        _write_session_config(req_session_id, req_mode, working_dir=effective_cwd, model_name=_model_name2)
        cmd = [_bzcode, "--stdio", "--resume", req_session_id]
        api_keys = _read_api_keys()
        if not api_keys.get("BZ_API_KEY"):
            print("[pool] BZ_API_KEY not found — rejecting session spawn", file=sys.stderr)
            raise HTTPException(401, "No BZ_API_KEY configured. Please set it in Settings → AI API Key.")
        env = {**os.environ, **api_keys, "BZ_PYTHON": sys.executable,
               **({"BZ_HOME": _bz_home} if _bz_home else {})}

        is_reuse = req_session_id in agent_pool._entries
        try:
            entry = await agent_pool.get_or_create(
                req_session_id, effective_cwd, req_mode, _bzcode, cmd, env)
        except (FileNotFoundError, PermissionError) as exc:
            raise HTTPException(500, f"Failed to start bzcode: {exc}")

        # Always read history from .jsonl — works for both new (empty) and reused sessions.
        # The session message from bzcode startup is consumed by the dispatcher,
        # not by this endpoint.
        messages = entry._read_session_messages()

        _active_sessions.add(req_session_id)
        return {
            "sessionId": req_session_id,
            "cwd": effective_cwd,
            "mode": req_mode,
            "sessionMode": entry.session_mode,
            "agentStatus": entry.agent_status,
            "pid": entry.proc.pid if entry.proc else None,
            "reused": is_reuse,
            "messages": messages,
            "modes": entry.available_modes,
            "commands": entry.available_commands,
        }

    @misc_router.get("/api/pool/{session_id}/stream")
    async def pool_stream(session_id: str, request: Request):
        """SSE stream of all bzcode output for a pooled agent."""
        entry = agent_pool._entries.get(session_id)
        if entry is None:
            raise HTTPException(404, f"No agent in pool with sessionId={session_id}")
        if entry.is_dead:
            raise HTTPException(410, f"Agent {session_id} process is dead")

        q = entry.subscribe(replay=True)

        async def event_generator():
            try:
                while True:
                    try:
                        raw = await asyncio.wait_for(q.get(), timeout=15)
                    except asyncio.TimeoutError:
                        yield ": ping\n\n"
                        continue
                    if raw is None:
                        yield f"data: {json.dumps({'type': 'system', 'message': 'Agent process exited'})}\n\n"
                        break
                    if raw and raw[0] == "{":
                        yield f"data: {raw}\n\n"
            except asyncio.CancelledError:
                pass
            finally:
                entry.unsubscribe(q)

        from starlette.responses import StreamingResponse
        return StreamingResponse(
            event_generator(),
            media_type="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "Connection": "keep-alive",
                "X-Accel-Buffering": "no",
            },
        )

    @misc_router.post("/api/pool/{session_id}/send")
    async def pool_send(session_id: str, request: Request):
        """Send a message directly to a pooled agent's stdin."""
        body = await request.json()

        entry = agent_pool._entries.get(session_id)
        if entry is None:
            raise HTTPException(404, f"No agent in pool with sessionId={session_id}")
        if entry.is_dead:
            raise HTTPException(410, f"Agent {session_id} process is dead")
        if entry.proc.stdin is None:
            raise HTTPException(500, "Agent stdin is not available")

        # Accept full message JSON (type, content, subtype, etc.) or legacy {message} field
        if "type" in body:
            msg = body
        else:
            message = body.get("message", "")
            if not message:
                raise HTTPException(400, "message or type is required")
            msg = {"type": "user", "content": message}

        # Seed the SSE replay buffer with a new user prompt so a client that connects
        # or reconnects mid-turn can render it (it isn't in the .jsonl transcript yet
        # and bzcode doesn't echo it on stdout). Skip permission/input replies — those
        # are mid-turn responses, not the start of a new turn. The buffered/echoed copy
        # keeps the client's `clientId` so the sender can dedup its optimistic bubble.
        if msg.get("type") == "user" and not msg.get("subtype"):
            entry.seed_user_turn(json.dumps(msg))

        # bzcode never sees the client-only clientId field.
        stdin_msg = {k: v for k, v in msg.items() if k != "clientId"}
        payload = json.dumps(stdin_msg) + "\n"
        entry.proc.stdin.write(payload.encode())
        await entry.proc.stdin.drain()
        return {"ok": True, "pid": entry.proc.pid, "sessionId": session_id}

    @misc_router.get("/agent-modes")
    async def agent_modes():
        return _load_mode_config()

    # Cache for the live model list fetched from the BoltzBit API
    _models_cache: list = []
    _models_cache_ts: float = 0.0
    _MODELS_CACHE_TTL = 300  # 5 minutes
    _BZ_MODELS_URL = "https://flow.boltzbit.com/bz-api/v1/ai/models"

    async def _fetch_bz_models(api_key: str) -> list:
        import time, aiohttp
        nonlocal _models_cache, _models_cache_ts
        now = time.monotonic()
        if _models_cache and (now - _models_cache_ts) < _MODELS_CACHE_TTL:
            return _models_cache
        try:
            async with aiohttp.ClientSession() as sess:
                async with sess.get(
                    _BZ_MODELS_URL,
                    headers={"x-api-key": api_key},
                    timeout=aiohttp.ClientTimeout(total=8),
                ) as resp:
                    if resp.status == 200:
                        data = await resp.json()
                        _models_cache = [
                            {"id": m["codename"], "displayName": m["displayName"]}
                            for m in data if m.get("codename") and m.get("displayName")
                        ]
                        _models_cache_ts = now
        except Exception as exc:
            print(f"[models] failed to fetch model list: {exc}", file=sys.stderr)
        return _models_cache

    @misc_router.get("/api/models")
    async def list_models(session_id: str = Query("")):
        api_keys = _read_api_keys()
        bz_api_key = api_keys.get("BZ_API_KEY", "")
        models = await _fetch_bz_models(bz_api_key) if bz_api_key else []
        current = ""
        if session_id:
            entry = agent_pool._entries.get(session_id)
            if entry:
                current = entry.model_info.get("name", "") or entry.model_info.get("displayName", "")
            if not current:
                _meta_file = SESSIONS_DIR / session_id / "meta.json"
                if _meta_file.exists():
                    try:
                        current = json.loads(_meta_file.read_text()).get("model", "")
                    except Exception:
                        pass
        return {"models": models, "current": current}

    @misc_router.post("/api/sessions/{session_id}/model")
    async def set_session_model(session_id: str, request: Request):
        body = await request.json()
        model_id = body.get("model", "")
        if not model_id:
            raise HTTPException(400, "model is required")
        entry = agent_pool._entries.get(session_id)
        req_mode = "general"
        effective_cwd = str(app.state.default_cwd)
        if entry:
            req_mode = entry.mode
            effective_cwd = entry.cwd
        else:
            _meta_file = SESSIONS_DIR / session_id / "meta.json"
            if _meta_file.exists():
                try:
                    _meta = json.loads(_meta_file.read_text())
                    req_mode = _meta.get("mode", req_mode)
                    effective_cwd = _meta.get("workingDir", effective_cwd)
                except Exception:
                    pass
        _write_session_config(session_id, req_mode, working_dir=effective_cwd, model_name=model_id)
        # Update live settings.json with model field so bzcode picks it up on next turn
        _cfg_dir = SESSIONS_DIR / session_id
        _settings_path = _cfg_dir / "settings.json"
        try:
            _settings = {}
            if _settings_path.exists():
                _settings = json.loads(_settings_path.read_text())
            _settings["model"] = model_id
            _settings_path.write_text(json.dumps(_settings, indent=2))
        except Exception as exc:
            print(f"[model] failed to patch settings.json: {exc}", file=sys.stderr)
        return {"ok": True, "sessionId": session_id, "model": model_id}

    # ── File read / write ─────────────────────────────────────────────────────

    @files_router.get("/api/file")
    async def read_file(path: str = Query("")):
        if not path:
            raise HTTPException(400, "path required")
        p = Path(path)
        if not p.exists() or not p.is_file():
            raise HTTPException(404, "file not found")
        return {"path": str(p), "content": p.read_text(errors="replace")}

    @files_router.put("/api/file")
    async def write_file(body: WriteFileBody):
        p = Path(body.path)
        p.parent.mkdir(parents=True, exist_ok=True)
        p.write_text(body.content, encoding="utf-8")
        return {"ok": True, "path": str(p)}

    @files_router.post("/api/file/rename")
    async def file_rename(body: FileRenameBody):
        p = Path(body.path)
        if not p.exists():
            raise HTTPException(404, "path not found")
        dest = p.parent / body.newName
        if dest.exists():
            raise HTTPException(409, "destination already exists")
        try:
            p.rename(dest)
            return {"ok": True, "path": str(dest)}
        except Exception as exc:
            raise HTTPException(500, str(exc))

    @files_router.post("/api/file/duplicate")
    async def file_duplicate(body: FileDuplicateBody):
        import shutil as _sh
        p = Path(body.path)
        if not p.exists():
            raise HTTPException(404, "path not found")
        if p.is_dir():
            raise HTTPException(400, "directory duplication not supported")
        stem, suffix = p.stem, p.suffix
        dest = p.parent / f"{stem} copy{suffix}"
        n = 2
        while dest.exists():
            dest = p.parent / f"{stem} copy {n}{suffix}"
            n += 1
        try:
            _sh.copy2(str(p), str(dest))
            return {"ok": True, "path": str(dest)}
        except Exception as exc:
            raise HTTPException(500, str(exc))

    @files_router.get("/api/file/download")
    async def file_download(path: str = Query(...)):
        import mimetypes as _mt
        p = Path(path)
        if not p.exists() or p.is_dir():
            raise HTTPException(404, "file not found")
        mime, _ = _mt.guess_type(str(p))
        mime = mime or "application/octet-stream"
        from fastapi.responses import Response as _Resp
        return _Resp(
            content=p.read_bytes(),
            media_type=mime,
            headers={"Content-Disposition": f'attachment; filename="{p.name}"'},
        )

    @files_router.get("/api/file/view")
    async def file_view(path: str = Query(...)):
        import mimetypes as _mt
        p = Path(path)
        if not p.exists() or p.is_dir():
            raise HTTPException(404, "file not found")
        mime, _ = _mt.guess_type(str(p))
        mime = mime or "application/octet-stream"
        from fastapi.responses import Response as _Resp
        return _Resp(
            content=p.read_bytes(),
            media_type=mime,
            headers={"Content-Disposition": f'inline; filename="{p.name}"'},
        )

    @files_router.post("/api/file/upload")
    async def file_upload(request: Request, background_tasks: BackgroundTasks):
        form = await request.form()
        upload = form.get("file")
        dir_path = str(form.get("dir", "")).strip() or app.state.default_cwd
        if upload is None:
            raise HTTPException(400, "expected field 'file'")
        filename = getattr(upload, "filename", None) or "upload"
        data = await upload.read()
        dest = Path(dir_path) / filename
        dest.parent.mkdir(parents=True, exist_ok=True)
        # Avoid clobbering: append (2), (3), … if file exists
        if dest.exists():
            stem, suffix = dest.stem, dest.suffix
            n = 2
            while dest.exists():
                dest = dest.parent / f"{stem} ({n}){suffix}"
                n += 1
        dest.write_bytes(data)
        ppt_processing = dest.suffix.lower() in ('.pptx', '.ppt')
        if ppt_processing:
            # Parse and write sidecar in background so it's ready before first open
            background_tasks.add_task(ppt_load, path=str(dest))
        return {"ok": True, "path": str(dest), "name": dest.name, "pptProcessing": ppt_processing}

    @files_router.post("/api/file/mkdir")
    async def file_mkdir(body: dict):
        path_str = str(body.get("path", "")).strip()
        if not path_str:
            raise HTTPException(400, "path required")
        p = Path(path_str)
        p.mkdir(parents=True, exist_ok=True)
        return {"ok": True, "path": str(p)}

    @files_router.delete("/api/file")
    async def file_delete(path: str = Query(...)):
        import shutil as _sh
        p = Path(path)
        if not p.exists():
            raise HTTPException(404, "path not found")
        if p.is_dir():
            _sh.rmtree(str(p))
        else:
            p.unlink()
        return {"ok": True}

    @files_router.get("/api/doc/cursor")
    async def get_cursor(path: str = Query(...)):
        if not path:
            raise HTTPException(400, "path required")
        return _cursor_store.get(path, {"selStart": 0, "selEnd": 0})

    @files_router.put("/api/doc/cursor")
    async def put_cursor(body: CursorBody):
        _cursor_store[body.path] = {"selStart": body.selStart, "selEnd": body.selEnd}
        return {"ok": True}

    # ── Settings ──────────────────────────────────────────────────────────────

    @misc_router.get("/settings/resources")
    async def settings_resources():
        import shutil as _sh
        session_count = session_bytes = 0
        if SESSIONS_DIR.exists():
            for f in SESSIONS_DIR.glob("*.jsonl"):
                try:
                    session_bytes += f.stat().st_size
                    session_count += 1
                except Exception:
                    pass
        server_data_bytes = 0
        if SERVER_DATA_DIR.exists():
            for f in SERVER_DATA_DIR.rglob("*"):
                try:
                    if f.is_file():
                        server_data_bytes += f.stat().st_size
                except Exception:
                    pass
        try:
            disk = _sh.disk_usage(Path.home())
            disk_info = {"total": disk.total, "used": disk.used, "free": disk.free}
        except Exception:
            disk_info = {"total": 0, "used": 0, "free": 0}
        return {"sessions": {"count": session_count, "bytes": session_bytes},
                "serverData": {"bytes": server_data_bytes}, "disk": disk_info}

    @misc_router.delete("/settings/sessions/clear")
    async def clear_sessions(olderThanDays: int = Query(30)):
        import time as _t
        cutoff = _t.time() - max(1, olderThanDays) * 86_400
        deleted = 0
        if SESSIONS_DIR.exists():
            for f in SESSIONS_DIR.glob("*.jsonl"):
                try:
                    if f.stat().st_mtime < cutoff:
                        f.unlink()
                        deleted += 1
                except Exception:
                    pass
        return {"deleted": deleted}

    # ── Document / Office ─────────────────────────────────────────────────────

    def _doc_sidecar_path(p: Path) -> Path:
        return Path(str(p) + '.json')

    def _extract_images(blocks: list) -> tuple:
        """Move imageUrl out of block styles into a shared images dict keyed by stable ID.
        Returns (mutated_blocks, images_dict).  Mutates block style dicts in-place."""
        import hashlib as _hl
        images: dict = {}
        for block in blocks:
            for style in (block.get("styles") or []):
                url = style.get("imageUrl")
                if not url:
                    continue
                prefix = url[:200] if isinstance(url, str) else ""
                key = "img-" + _hl.md5(prefix.encode()).hexdigest()[:8]
                images[key] = {
                    "url": url,
                    **({"width":  style["imageWidth"]}  if style.get("imageWidth")  else {}),
                    **({"height": style["imageHeight"]} if style.get("imageHeight") else {}),
                }
                style.pop("imageUrl",    None)
                style.pop("imageWidth",  None)
                style.pop("imageHeight", None)
                style["imageId"] = key
        return blocks, images

    def _expand_images(blocks: list, images: dict) -> list:
        """Inline imageUrl/Width/Height back into block styles from the images dict."""
        for block in blocks:
            for style in (block.get("styles") or []):
                img_id = style.get("imageId")
                if img_id and img_id in images:
                    img = images[img_id]
                    style["imageUrl"]    = img["url"]
                    if img.get("width"):  style["imageWidth"]  = img["width"]
                    if img.get("height"): style["imageHeight"] = img["height"]
        return blocks

    @misc_router.post("/api/doc/parse")
    async def doc_parse(request: Request):
        ct = request.headers.get("content-type", "")
        _MAX_DOC_BYTES = 50 * 1024 * 1024
        try:
            if "multipart" in ct:
                form = await request.form()
                upload = form.get("file")
                if upload is None:
                    raise HTTPException(400, "expected field 'file'")
                filename = getattr(upload, "filename", None) or "upload"
                data = await upload.read()
                # Uploaded files have no path, parse directly
                return _detect_and_parse(filename, data)
            else:
                body = await request.json()
                path_str = str(body.get("path", "")).strip()
                if not path_str:
                    raise HTTPException(400, "path required")
                p = Path(path_str)
                if not p.exists():
                    raise HTTPException(404, "file not found")

                force_refresh = bool(body.get("force"))

                # --- Sidecar check (DOCX only) ---
                if p.suffix.lower() in (".docx", ".doc"):
                    sc_path = _doc_sidecar_path(p)
                    # On force_refresh, only discard the sidecar when the DOCX file itself
                    # is newer — meaning the file was edited externally (e.g. by the agent).
                    # If the sidecar is newer, it is the authoritative source (written by the
                    # last auto-save) and must be preserved; re-parsing from DOCX would lose
                    # bz-office fields (imagePlacedX/Y, imageWrap) that have no DOCX equivalent.
                    if force_refresh and sc_path.exists():
                        try:
                            if p.stat().st_mtime > sc_path.stat().st_mtime:
                                sc_path.unlink()
                        except Exception:
                            pass
                    if sc_path.exists():
                        sidecar = json.loads(sc_path.read_text(encoding='utf-8'))
                        # Accept sidecar if any style carries fontFamily (text content is
                        # up-to-date) OR imageUrl/imageId (image-only blocks have no fontFamily).
                        # Old sidecars that predate fontFamily support will have neither,
                        # causing a re-parse from DOCX to pick up the new field.
                        has_valid = any(
                            s.get("fontFamily") or s.get("imageUrl") or s.get("imageId")
                            for b in (sidecar.get("blocks") or [])
                            for s in (b.get("styles") or [])
                        )
                        if has_valid or not sidecar.get("blocks"):
                            # Expand imageId references back to full imageUrl for the frontend
                            if sidecar.get("images"):
                                _expand_images(sidecar["blocks"], sidecar["images"])
                            return sidecar

                if p.stat().st_size > _MAX_DOC_BYTES:
                    raise HTTPException(413, "file too large (max 50 MB)")
                data = p.read_bytes()
                filename = p.name
                result = _detect_and_parse(filename, data)

                # Write sidecar for DOCX files so subsequent opens are instant
                if p.suffix.lower() in (".docx", ".doc"):
                    sc_path = _doc_sidecar_path(p)
                    try:
                        sc_path.write_text(
                            json.dumps(result, ensure_ascii=False, indent=2),
                            encoding='utf-8'
                        )
                    except Exception:
                        pass  # sidecar write failure is non-fatal

                return result
        except HTTPException:
            raise
        except ValueError as exc:
            raise HTTPException(400, str(exc))
        except Exception as exc:
            raise HTTPException(422, f"could not parse: {exc}")

    @misc_router.put("/api/doc/save")
    async def doc_save(body: DocSaveBody):
        p = Path(body.path)
        if p.suffix.lower() not in (".docx", ".doc"):
            raise HTTPException(400, "only DOCX files can be saved")
        try:
            sc_path = _doc_sidecar_path(p)
            word_count = sum(len(b.get("text", "").split()) for b in body.blocks)
            # Load existing sidecar to preserve metadata (defaultFont, pages, etc.)
            if sc_path.exists():
                sidecar = json.loads(sc_path.read_text(encoding='utf-8'))
            else:
                sidecar = {"filename": p.name, "type": "docx", "truncated": False}
            # Extract imageUrl fields from styles into a shared images registry so
            # the sidecar doesn't embed large base64 blobs inside style ranges.
            import copy as _copy
            save_blocks = _copy.deepcopy(body.blocks)
            save_blocks, new_images = _extract_images(save_blocks)
            # Merge into any images already in the sidecar (preserves old entries)
            existing_images = sidecar.get("images") or {}
            existing_images.update(new_images)
            if existing_images:
                sidecar["images"] = existing_images
            sidecar["blocks"] = save_blocks
            sidecar["wordCount"] = word_count
            # Write DOCX first, then sidecar — so sidecar always has the newest mtime.
            # The mtime ordering is used by force_refresh to decide whether the DOCX
            # was externally modified (agent edit) vs. saved by bz-office (sidecar newer).
            try:
                docx_bytes = _blocks_to_docx(body.blocks)  # original blocks carry imageUrl
                p.write_bytes(docx_bytes)
            except Exception:
                pass  # non-fatal: sidecar is the source of truth
            sc_path.write_text(
                json.dumps(sidecar, ensure_ascii=False, indent=2),
                encoding='utf-8'
            )
            return {"ok": True, "path": str(p), "wordCount": word_count}
        except Exception as exc:
            raise HTTPException(500, f"could not save: {exc}")

    @misc_router.get("/api/doc/download")
    async def doc_download(path: str):
        """Convert sidecar → DOCX and stream back as a file download."""
        from fastapi.responses import Response as _Resp
        p = Path(path)
        sc_path = _doc_sidecar_path(p)
        if not sc_path.exists():
            raise HTTPException(404, "sidecar not found — open the file first")
        try:
            sidecar = json.loads(sc_path.read_text(encoding='utf-8'))
            blocks = sidecar.get("blocks", [])
            if sidecar.get("images"):
                _expand_images(blocks, sidecar["images"])
            docx_bytes = _blocks_to_docx(blocks)
            return _Resp(
                content=docx_bytes,
                media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                headers={"Content-Disposition": f'attachment; filename="{p.name}"'},
            )
        except HTTPException:
            raise
        except Exception as exc:
            raise HTTPException(500, f"could not generate DOCX: {exc}")

    def _sidecar_to_api(sidecar: dict) -> dict:
        """Convert sidecar JSON to the API schema the frontend expects."""
        sheets = []
        for sheet in sidecar.get('sheets', []):
            api_cells: dict = {}
            for ref, cd in sheet.get('cells', {}).items():
                api_cd: dict = {}
                v = cd.get('v')
                if v is not None:
                    api_cd['value'] = v
                if 'f' in cd:
                    api_cd['formula'] = cd['f']
                s = cd.get('s', {})
                if s.get('bold'):   api_cd['fontBold']  = True
                if s.get('italic'): api_cd['fontItalic'] = True
                if s.get('fg'):
                    _fg = s['fg']
                    if _fg.startswith('#') and len(_fg) == 7:
                        _fg = 'FF' + _fg[1:]
                    api_cd['fontColor'] = _fg
                if s.get('bg'):
                    _bg = s['bg']
                    if _bg.startswith('#') and len(_bg) == 7:
                        _bg = 'FF' + _bg[1:]
                    api_cd['bgColor'] = _bg
                if s.get('align'):  api_cd['align']     = s['align']
                if s.get('format') is not None: api_cd['dataFormatString'] = s['format']
                if s.get('wrap') is not None:   api_cd['wrapText'] = bool(s['wrap'])
                if s.get('fontSize') is not None: api_cd['fontSize'] = s['fontSize']
                if api_cd:
                    api_cells[ref] = api_cd
            # Grid dimensions: prefer explicit grid field, fall back to col_widths array
            grid_obj = sheet.get('grid', {})
            col_widths: dict = {str(k): v for k, v in (grid_obj.get('columnIndexToWidth') or {}).items()}
            if not col_widths:
                for i, w in enumerate(sheet.get('col_widths', [])):
                    if w:
                        col_widths[str(i)] = max(30, int(w * 7.5) if w < 50 else w)
            row_heights: dict = {str(k): v for k, v in (grid_obj.get('rowIndexToHeight') or {}).items()}
            sheets.append({
                "sheetName": sheet['name'],
                "cells": api_cells,
                "images": [],
                "columnIndexToWidth": col_widths,
                "rowIndexToHeight": row_heights,
                "hiddenColIndices": [],
                "hiddenRowIndices": [],
                "mergedCellIndices": [],
                "mergedCellRanges": sheet.get('mergedCells', []),
            })
        p = Path(sidecar.get('xlsx_path', ''))
        return {"id": p.stem, "name": p.stem, "sheets": sheets, "sources": []}

    def _find_excel_worker() -> Path | None:
        candidates = [
            Path(__file__).parent / "bzcode_assets" / "scripts" / "excel-worker.py",
            Path.cwd() / "bzcode_assets" / "scripts" / "excel-worker.py",
        ]
        for c in candidates:
            if c.exists():
                return c
        return None

    @misc_router.get("/api/excel/load")
    async def excel_load(path: str = Query(...)):
        p = Path(path)
        if not p.exists():
            raise HTTPException(404, "file not found")

        # Fast path: sidecar JSON created by excel-worker.py has correct computed values
        json_path = p.parent / f".{p.name}.excel.json"
        if json_path.exists():
            try:
                sidecar = json.loads(json_path.read_text(encoding='utf-8'))
                return _sidecar_to_api(sidecar)
            except Exception:
                pass  # fall through to openpyxl

        # Slow path: parse xlsx with openpyxl (for files not created by excel-worker.py)
        try:
            import openpyxl
            wb_vals  = openpyxl.load_workbook(p, data_only=True)
            wb_forms = openpyxl.load_workbook(p, data_only=False)
            formula_map: dict = {}
            for ws_f in wb_forms.worksheets:
                for row in ws_f.iter_rows():
                    for cell in row:
                        if isinstance(cell.value, str) and cell.value.startswith("="):
                            formula_map[(ws_f.title, cell.coordinate)] = cell.value
            sheets = []
            for ws in wb_vals.worksheets:
                cells: dict = {}
                col_widths: dict = {}
                for row in ws.iter_rows():
                    for cell in row:
                        cell_id    = cell.coordinate
                        formula    = formula_map.get((ws.title, cell_id))
                        cached_val = cell.value
                        if cached_val is None and not formula:
                            continue
                        cd: dict = {}
                        if cached_val is not None:
                            cd["value"] = str(cached_val) if not isinstance(cached_val, (int, float, bool)) else cached_val
                        if formula:
                            cd["formula"] = formula
                        try:
                            f = cell.font
                            if f.bold:   cd["fontBold"]  = True
                            if f.italic: cd["fontItalic"] = True
                            # Read explicit font color; skip Excel "auto" colors (alpha=00)
                            if f.color and f.color.type == 'rgb':
                                rgb = f.color.rgb  # AARRGGBB string
                                if rgb[:2].upper() not in ('00', ''):
                                    cd["fontColor"] = rgb
                        except Exception:
                            pass
                        try:
                            fill = cell.fill
                            if fill and fill.fill_type == "solid" and fill.fgColor and fill.fgColor.type == "rgb":
                                rgb = fill.fgColor.rgb
                                if rgb not in ("FF000000", "00000000", "FFFFFFFF"):
                                    cd["bgColor"] = rgb  # FFRRGGBB format expected by renderer
                        except Exception:
                            pass
                        if cd:
                            cells[cell_id] = cd
                for col_letter, dim in (ws.column_dimensions or {}).items():
                    if dim.width:
                        idx = openpyxl.utils.column_index_from_string(col_letter) - 1
                        col_widths[str(idx)] = max(30, int(dim.width * 7.5))
                merged_ranges = [str(r) for r in ws.merged_cells.ranges]
                sheets.append({"sheetName": ws.title, "cells": cells, "images": [],
                               "columnIndexToWidth": col_widths, "rowIndexToHeight": {},
                               "hiddenColIndices": [], "hiddenRowIndices": [], "mergedCellIndices": [],
                               "mergedCellRanges": merged_ranges})
            return {"id": p.stem, "name": p.stem, "sheets": sheets, "sources": []}
        except Exception as exc:
            raise HTTPException(500, str(exc))

    @misc_router.put("/api/excel/patch")
    async def excel_patch(body: ExcelPatchBody):
        """
        Merge cell updates into the sidecar JSON, re-evaluate all formulas via
        excel-worker.py --recalc, rewrite the xlsx, and return the updated cells
        (all formula cells + the patched cells) so the frontend can re-render.
        """
        p = Path(body.path)
        json_path = p.parent / f".{p.name}.excel.json"
        if not json_path.exists():
            raise HTTPException(404, "sidecar JSON not found — file must be created with excel-worker.py first")

        script = _find_excel_worker()
        if not script:
            raise HTTPException(500, "excel-worker.py not found on server")

        patch_json = json.dumps({"sheet": body.sheet, "cells": body.cells})
        import asyncio as _asyncio
        proc = await _asyncio.create_subprocess_exec(
            sys.executable, str(script),
            "--recalc", str(json_path),
            "--out",    str(p),
            "--patch",  patch_json,
            stdout=_asyncio.subprocess.PIPE,
            stderr=_asyncio.subprocess.PIPE,
        )
        stdout, _ = await proc.communicate()
        try:
            result = json.loads(stdout)
        except Exception:
            raise HTTPException(500, "excel-worker.py returned invalid output")
        if not result.get("ok"):
            raise HTTPException(500, result.get("error", "recalculation failed"))

        # Return updated sidecar converted to API schema
        sidecar = json.loads(json_path.read_text(encoding='utf-8'))
        return _sidecar_to_api(sidecar)

    @misc_router.get("/api/excel/download")
    async def excel_download(path: str = Query(...)):
        from fastapi.responses import FileResponse
        p = Path(path)
        if not p.exists():
            raise HTTPException(404, "file not found")
        return FileResponse(
            path=str(p),
            filename=p.name,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )

    @misc_router.put("/api/excel/grid")
    async def excel_grid(body: ExcelGridBody):
        """Persist column/row dimension changes from frontend drag-resize."""
        p = Path(body.path)
        json_path = p.parent / f".{p.name}.excel.json"
        if not json_path.exists():
            raise HTTPException(404, "sidecar not found")
        sidecar = json.loads(json_path.read_text(encoding='utf-8'))
        target = body.sheet or (sidecar['sheets'][0]['name'] if sidecar.get('sheets') else "")
        for sheet in sidecar.get('sheets', []):
            if sheet['name'] == target:
                if 'grid' not in sheet:
                    sheet['grid'] = {}
                sheet['grid']['columnIndexToWidth'] = body.columnIndexToWidth
                sheet['grid']['rowIndexToHeight']   = body.rowIndexToHeight
                break
        json_path.write_text(json.dumps(sidecar, indent=2, ensure_ascii=False), encoding='utf-8')
        return {"ok": True}

    @misc_router.put("/api/excel/renamesheet")
    async def excel_renamesheet(body: ExcelRenameSheetBody):
        """Rename a sheet in the sidecar JSON."""
        p = Path(body.path)
        json_path = p.parent / f".{p.name}.excel.json"
        if not json_path.exists():
            raise HTTPException(404, "sidecar not found")
        sidecar = json.loads(json_path.read_text(encoding='utf-8'))
        names = {s['name'] for s in sidecar.get('sheets', [])}
        if body.oldName not in names:
            raise HTTPException(404, "sheet not found")
        if body.newName in names and body.newName != body.oldName:
            raise HTTPException(409, "name already in use")
        for sheet in sidecar.get('sheets', []):
            if sheet['name'] == body.oldName:
                sheet['name'] = body.newName
                break
        json_path.write_text(json.dumps(sidecar, indent=2, ensure_ascii=False), encoding='utf-8')
        return {"ok": True, "name": body.newName}

    @misc_router.post("/api/excel/addsheet")
    async def excel_addsheet(body: ExcelAddSheetBody):
        """Append a new empty sheet to the sidecar JSON."""
        p = Path(body.path)
        json_path = p.parent / f".{p.name}.excel.json"
        if not json_path.exists():
            raise HTTPException(404, "sidecar not found")
        sidecar = json.loads(json_path.read_text(encoding='utf-8'))
        existing = {s['name'] for s in sidecar.get('sheets', [])}
        if body.sheetName in existing:
            raise HTTPException(409, "sheet already exists")
        sidecar.setdefault('sheets', []).append({
            "name": body.sheetName,
            "col_widths": [],
            "cells": {},
            "grid": {"columnIndexToWidth": {}, "rowIndexToHeight": {}},
        })
        json_path.write_text(json.dumps(sidecar, indent=2, ensure_ascii=False), encoding='utf-8')
        return {"ok": True, "sheetName": body.sheetName}

    @misc_router.put("/api/excel/merge")
    async def excel_merge(body: ExcelMergeBody):
        """Replace the mergedCells list for a sheet, rebuild xlsx, return updated data."""
        p = Path(body.path)
        json_path = p.parent / f".{p.name}.excel.json"
        if not json_path.exists():
            raise HTTPException(404, "sidecar not found")

        script = _find_excel_worker()
        if not script:
            raise HTTPException(500, "excel-worker.py not found on server")

        sidecar = json.loads(json_path.read_text(encoding='utf-8'))
        target = body.sheet or (sidecar['sheets'][0]['name'] if sidecar.get('sheets') else "")
        for sheet in sidecar.get('sheets', []):
            if sheet['name'] == target:
                sheet['mergedCells'] = body.mergedCells
                break
        json_path.write_text(json.dumps(sidecar, indent=2, ensure_ascii=False), encoding='utf-8')

        import asyncio as _asyncio
        proc = await _asyncio.create_subprocess_exec(
            sys.executable, str(script),
            "--recalc", str(json_path),
            "--out",    str(p),
            stdout=_asyncio.subprocess.PIPE,
            stderr=_asyncio.subprocess.PIPE,
        )
        await proc.communicate()

        sidecar = json.loads(json_path.read_text(encoding='utf-8'))
        return _sidecar_to_api(sidecar)

    @misc_router.put("/api/excel/save")
    async def excel_save(body: ExcelSaveBody):
        if not body.path:
            raise HTTPException(400, "path required")
        try:
            import openpyxl
            from openpyxl.styles import Font, PatternFill
            wb = openpyxl.Workbook()
            wb.remove(wb.active)
            for sheet in body.sheets:
                ws = wb.create_sheet(title=sheet.get("sheetName", "Sheet"))
                for cell_id, cd in sheet.get("cells", {}).items():
                    try:
                        cell = ws[cell_id]
                        formula = cd.get("formula")
                        v = cd.get("value")
                        if formula and isinstance(formula, str) and formula.startswith("="):
                            cell.value = formula
                        elif v is not None:
                            try:
                                cell.value = float(v) if isinstance(v, str) and v.replace(".", "", 1).lstrip("-").isdigit() else v
                            except Exception:
                                cell.value = v
                        font_kw = {}
                        if cd.get("fontBold"):   font_kw["bold"] = True
                        if cd.get("fontItalic"): font_kw["italic"] = True
                        if cd.get("fontFamily"): font_kw["name"] = cd["fontFamily"]
                        if cd.get("fontSize"):   font_kw["size"] = cd["fontSize"] / 20
                        if font_kw: cell.font = Font(**font_kw)
                        if cd.get("dataFormatString"): cell.number_format = cd["dataFormatString"]
                    except Exception:
                        pass
            p = Path(body.path)
            p.parent.mkdir(parents=True, exist_ok=True)
            wb.save(p)
            return {"ok": True, "path": str(p)}
        except Exception as exc:
            raise HTTPException(500, str(exc))

    @misc_router.get("/api/ppt/load")
    async def ppt_load(path: str = Query(...)):
        p = Path(path)
        if not p.exists():
            raise HTTPException(404, "file not found")
        try:
            import json as _json
            # Ground truth: sidecar JSON takes priority over the binary .pptx
            json_path = p.parent / f".{p.name}.json"
            if json_path.exists():
                data = _json.loads(json_path.read_text())
                slides = data.get("slides", [])
                # Re-export the .pptx whenever the sidecar is newer (nanosecond
                # precision) so the file on disk stays in sync with the app.
                if json_path.stat().st_mtime_ns >= p.stat().st_mtime_ns:
                    try:
                        _pptx_export(p, slides)
                    except Exception:
                        pass
                return {"slides": slides}
            # Fall back: parse .pptx with full extraction
            from pptx import Presentation
            from pptx.oxml.ns import qn as _qn
            import secrets as _sec2
            import base64 as _b64

            prs = Presentation(str(p))
            sw = prs.slide_width
            sh = prs.slide_height
            CW, CH = 896, 504
            sx = CW / sw if sw else 1
            sy = CH / sh if sh else 1

            # ── Build theme color lookup via PPTX zip (works across pptx versions) ─
            theme_colors: dict = {}
            try:
                import zipfile as _zf
                from lxml import etree as _etree
                _NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                with _zf.ZipFile(str(p)) as _z:
                    _theme_files = sorted(
                        [n for n in _z.namelist() if '/theme/' in n and n.endswith('.xml')]
                    )
                    if _theme_files:
                        _tel = _etree.fromstring(_z.read(_theme_files[0]))
                        _clr = _tel.find(f'.//{{{_NS_A}}}clrScheme')
                        if _clr is not None:
                            for _child in _clr:
                                _tag = _child.tag.split('}')[1]
                                _srgb = _child.find(f'{{{_NS_A}}}srgbClr')
                                _sysc = _child.find(f'{{{_NS_A}}}sysClr')
                                if _srgb is not None:
                                    _v = _srgb.get('val', '')
                                    if len(_v) == 6:
                                        theme_colors[_tag] = f'#{_v.lower()}'
                                elif _sysc is not None:
                                    _last = _sysc.get('lastClr', '')
                                    if len(_last) == 6:
                                        theme_colors[_tag] = f'#{_last.lower()}'
            except Exception:
                pass

            # python-pptx theme color enum → XML slot name
            _THEME_SLOT = {
                1: 'dk1', 2: 'lt1', 3: 'dk2', 4: 'lt2',
                5: 'accent1', 6: 'accent2', 7: 'accent3',
                8: 'accent4', 9: 'accent5', 10: 'accent6',
                11: 'hlink', 12: 'folHlink',
            }

            _NS_A_FULL = 'http://schemas.openxmlformats.org/drawingml/2006/main'

            def _resolve_color(color_obj):
                """Return '#rrggbb' from a pptx color object, or None."""
                try:
                    rgb = color_obj.rgb
                    # RGBColor.__str__ returns 6-char hex (e.g. 'CEC4B6')
                    hex6 = str(rgb)
                    if len(hex6) == 6:
                        return f'#{hex6.lower()}'
                    # Fallback: index access
                    return f'#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}'
                except Exception:
                    pass
                try:
                    tc = color_obj.theme_color
                    tc_int = int(getattr(tc, 'real', tc))
                    slot = _THEME_SLOT.get(tc_int)
                    if slot and slot in theme_colors:
                        raw = theme_colors[slot]
                        # Apply luminance modifier (lum_mod in 1/1000 units)
                        try:
                            lum_mod = color_obj._element.find(
                                f'.//{{{_NS_A_FULL}}}lumMod')
                            if lum_mod is not None:
                                factor = int(lum_mod.get('val', '100000')) / 100000
                                h = raw.lstrip('#')
                                r2, g2, b2 = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
                                r2 = int(min(255, r2 * factor))
                                g2 = int(min(255, g2 * factor))
                                b2 = int(min(255, b2 * factor))
                                return f'#{r2:02x}{g2:02x}{b2:02x}'
                        except Exception:
                            pass
                        return raw
                except Exception:
                    pass
                return None

            def _get_fill(fill_obj, shape_element=None):
                """Return {'type': 'solid'|'gradient'|'none', 'color', 'opacity', 'gradient'}."""
                try:
                    # Check for gradient via XML first (more reliable than python-pptx API)
                    if shape_element is not None:
                        _sp = shape_element.find(_qn('p:spPr'))
                        if _sp is not None:
                            _gf = _sp.find('.//' + _qn('a:gradFill'))
                            if _gf is not None:
                                _lin = _gf.find(_qn('a:lin'))
                                _ang = int(_lin.get('ang', '0')) / 60000.0 if _lin is not None else 0.0
                                _stops = []
                                for _gs in sorted(_gf.findall('.//' + _qn('a:gs')), key=lambda g: int(g.get('pos', '0'))):
                                    _pos = int(_gs.get('pos', '0')) / 100000.0
                                    _c = _color_from_solid_elem(_gs)
                                    if _c:
                                        _stops.append({'pos': _pos, 'color': _c})
                                if _stops:
                                    return {'type': 'gradient', 'color': _stops[0]['color'],
                                            'gradient': {'angle': _ang, 'stops': _stops}}
                    ft = str(fill_obj.type)
                    if ft == 'None' or 'BACKGROUND' in ft:
                        return {'type': 'none'}
                    if 'SOLID' in ft or ft == '1':
                        c = _resolve_color(fill_obj.fore_color)
                        opacity = 1.0
                        try:
                            if shape_element is not None:
                                _sp = shape_element.find(_qn('p:spPr'))
                                _solid = _sp.find('.//' + _qn('a:solidFill')) if _sp is not None else None
                                if _solid is not None:
                                    for _child in _solid:
                                        _alpha = _child.find(f'{{{_NS_A_FULL}}}alpha')
                                        if _alpha is not None:
                                            opacity = int(_alpha.get('val', '100000')) / 100000.0
                                            break
                        except Exception:
                            pass
                        result = {'type': 'solid', 'color': c}
                        if opacity < 0.999:
                            result['opacity'] = round(opacity, 4)
                        return result
                except Exception:
                    pass
                return {'type': 'none'}

            def _color_from_solid_elem(solid_elem):
                """Return '#rrggbb' from an a:solidFill element, or None."""
                _SCHEME_MAP = {'bg1': 'lt1', 'tx1': 'dk1', 'bg2': 'lt2', 'tx2': 'dk2'}
                for _c in solid_elem:
                    _ctag = _c.tag.split('}')[-1]
                    if _ctag == 'srgbClr':
                        _v = _c.get('val', '')
                        return f'#{_v.lower()}' if len(_v) == 6 else None
                    elif _ctag == 'schemeClr':
                        _key = _SCHEME_MAP.get(_c.get('val', ''), _c.get('val', ''))
                        return theme_colors.get(_key)
                    elif _ctag == 'sysClr':
                        _v = _c.get('lastClr', '')
                        return f'#{_v.lower()}' if len(_v) == 6 else None
                return None

            def _get_bg(slide):
                """Resolve background, walking slide → layout → master.
                Returns (color_str, gradient_dict_or_None)."""
                _SCHEME_MAP = {'bg1': 'lt1', 'tx1': 'dk1', 'bg2': 'lt2', 'tx2': 'dk2'}

                def _resolve_scheme(val):
                    return theme_colors.get(_SCHEME_MAP.get(val, val))

                def _color_from_node(node):
                    """Resolve color from schemeClr/srgbClr/sysClr element."""
                    tag = node.tag.split('}')[-1]
                    if tag == 'srgbClr':
                        v = node.get('val', '')
                        return f'#{v.lower()}' if len(v) == 6 else None
                    if tag == 'schemeClr':
                        return _resolve_scheme(node.get('val', ''))
                    if tag == 'sysClr':
                        v = node.get('lastClr', '')
                        return f'#{v.lower()}' if len(v) == 6 else None
                    return None

                def _fill_from_csld(cSld_elem):
                    """Returns (color, gradient) from p:cSld, both may be None."""
                    bg = cSld_elem.find(_qn('p:bg'))
                    if bg is None:
                        return None, None
                    bgPr = bg.find(_qn('p:bgPr'))
                    if bgPr is not None:
                        solid = bgPr.find(_qn('a:solidFill'))
                        if solid is not None:
                            for child in solid:
                                c = _color_from_node(child)
                                if c:
                                    return c, None
                        grad = bgPr.find(_qn('a:gradFill'))
                        if grad is not None:
                            # Angle: DrawingML ang is in 60000ths of a degree,
                            # 0 = left-to-right, increasing clockwise
                            lin = grad.find(_qn('a:lin'))
                            angle_deg = 0.0
                            if lin is not None:
                                angle_deg = int(lin.get('ang', '0')) / 60000.0
                            gs_list = sorted(
                                grad.findall('.//' + _qn('a:gs')),
                                key=lambda g: int(g.get('pos', '0'))
                            )
                            stops = []
                            for gs in gs_list:
                                pos = int(gs.get('pos', '0')) / 100000.0
                                for child in gs:
                                    c = _color_from_node(child)
                                    if c:
                                        stops.append({'pos': pos, 'color': c})
                                        break
                            if stops:
                                fallback = stops[0]['color']
                                return fallback, {'angle': angle_deg, 'stops': stops}
                        return None, None
                    bgRef = bg.find(_qn('p:bgRef'))
                    if bgRef is not None:
                        for child in bgRef:
                            c = _color_from_node(child)
                            if c:
                                return c, None
                    return None, None

                # 1. Slide's own fill
                try:
                    fd = _get_fill(slide.background.fill)
                    if fd['type'] == 'solid' and fd.get('color'):
                        return fd['color'], None
                except Exception:
                    pass

                # 2. Walk up hierarchy
                for src in [slide, slide.slide_layout, slide.slide_layout.slide_master]:
                    try:
                        c, grad = _fill_from_csld(src.background._element)
                        if c:
                            return c, grad
                    except Exception:
                        pass

                return '#ffffff', None

            slides_out = []
            for slide in prs.slides:
                # Slide background
                bg_color = '#ffffff'
                bg_gradient = None
                try:
                    bg_color, bg_gradient = _get_bg(slide)
                except Exception:
                    pass

                def _iter_shapes_flat(shapes):
                    for s in shapes:
                        if s.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                            try:
                                yield from _iter_shapes_flat(s.shapes)
                            except Exception:
                                pass
                        else:
                            yield s

                def _extract_geometry(shape):
                    st, cr = 'rect', 0
                    try:
                        sp_pr = shape._element.find(_qn('p:spPr'))
                        if sp_pr is not None:
                            pg = sp_pr.find('.//' + _qn('a:prstGeom'))
                            if pg is not None:
                                prst = pg.get('prst', 'rect')
                                st = prst  # pass through actual preset name
                                if prst == 'roundRect':
                                    cr = 33.334  # OOXML default adj=16667
                                    av = pg.find(_qn('a:avLst'))
                                    if av is not None:
                                        for gd in av:
                                            fmla = gd.get('fmla', '')
                                            if fmla.startswith('val '):
                                                try:
                                                    cr = int(fmla.split()[1]) / 500
                                                except Exception:
                                                    pass
                    except Exception:
                        pass
                    return st, cr

                def _get_layout_ph_defaults(shape):
                    """Extract text style defaults from the matching layout placeholder."""
                    defaults = {}
                    try:
                        ph_fmt = shape.placeholder_format
                        if ph_fmt is None:
                            return defaults
                        layout = shape.part.slide.slide_layout
                        for lshp in layout.shapes:
                            try:
                                lpf = lshp.placeholder_format
                                if lpf is None or lpf.idx != ph_fmt.idx:
                                    continue
                                sp = lshp._element
                                tx = sp.find(_qn('p:txBody'))
                                if tx is None:
                                    break
                                bpr = tx.find(_qn('a:bodyPr'))
                                if bpr is not None:
                                    anchor = bpr.get('anchor', '')
                                    if anchor:
                                        defaults['textAnchor'] = anchor
                                lst = tx.find(_qn('a:lstStyle'))
                                if lst is not None:
                                    lvl1 = lst.find(_qn('a:lvl1pPr'))
                                    if lvl1 is not None:
                                        drpr = lvl1.find(_qn('a:defRPr'))
                                        if drpr is not None:
                                            sz = drpr.get('sz')
                                            if sz:
                                                defaults['fontSize'] = int(sz) // 100
                                            if drpr.get('b') == '1':
                                                defaults['bold'] = True
                                            if drpr.get('cap', '') == 'all':
                                                defaults['allCaps'] = True
                                            solid = drpr.find(_qn('a:solidFill'))
                                            if solid is not None:
                                                c = _color_from_solid_elem(solid)
                                                if c:
                                                    defaults['color'] = c
                                break
                            except Exception:
                                continue
                    except Exception:
                        pass
                    return defaults

                boxes = []

                # ── Layout non-placeholder shapes (connectors, rects) ─────────
                try:
                    for lshp in slide.slide_layout.shapes:
                        try:
                            if lshp.shape_type == 14:  # placeholder — skip
                                continue
                            lx = int((lshp.left or 0) * sx)
                            ly = int((lshp.top or 0) * sy)
                            sw_emu = lshp.width or 0
                            sh_emu = lshp.height or 0

                            if lshp.shape_type == 9:  # connector / line
                                ln_w_px = max(1, int(25400 * sx))  # default 2pt
                                ln_color = '#888888'
                                spPr2 = lshp._element.find(_qn('p:spPr'))
                                if spPr2 is not None:
                                    ln_el = spPr2.find(_qn('a:ln'))
                                    if ln_el is not None:
                                        ln_w_px = max(1, int(int(ln_el.get('w', '25400')) * sx))
                                        sol2 = ln_el.find(_qn('a:solidFill'))
                                        if sol2 is not None:
                                            c2 = _color_from_solid_elem(sol2)
                                            if c2:
                                                ln_color = c2
                                if sw_emu == 0:
                                    bw, bh = ln_w_px, max(1, int(sh_emu * sy))
                                    bx = lx - ln_w_px // 2
                                    by_ = ly
                                elif sh_emu == 0:
                                    bw, bh = max(1, int(sw_emu * sx)), ln_w_px
                                    bx = lx
                                    by_ = ly - ln_w_px // 2
                                else:
                                    bx, by_ = lx, ly
                                    bw = max(1, int(sw_emu * sx))
                                    bh = max(1, int(sh_emu * sy))
                                boxes.append({
                                    'id': _sec2.token_hex(4),
                                    'x': bx, 'y': by_, 'w': bw, 'h': bh,
                                    'rotation': 0, 'shapeType': 'rect',
                                    'text': '', 'styles': [],
                                    'boxStyle': {'bgColor': ln_color},
                                })
                            elif lshp.shape_type == 1:  # rectangle
                                _fd2 = {'type': 'none'}
                                _grad2 = None
                                _spPr2 = lshp._element.find(_qn('p:spPr'))
                                if _spPr2 is not None:
                                    _sol2 = _spPr2.find('.//' + _qn('a:solidFill'))
                                    if _sol2 is not None:
                                        _c2 = _color_from_solid_elem(_sol2)
                                        if _c2:
                                            _fd2 = {'type': 'solid', 'color': _c2}
                                    _gf2 = _spPr2.find('.//' + _qn('a:gradFill'))
                                    if _gf2 is not None and _fd2['type'] == 'none':
                                        _lin2 = _gf2.find(_qn('a:lin'))
                                        _ang2 = int(_lin2.get('ang', '0')) / 60000.0 if _lin2 is not None else 0.0
                                        _stops2 = []
                                        for _gs2 in sorted(_gf2.findall('.//' + _qn('a:gs')), key=lambda g: int(g.get('pos','0'))):
                                            _pos2 = int(_gs2.get('pos', '0')) / 100000.0
                                            # gs children are color nodes — reuse _color_from_solid_elem trick
                                            _c2 = _color_from_solid_elem(_gs2)
                                            if _c2:
                                                _stops2.append({'pos': _pos2, 'color': _c2})
                                        if _stops2:
                                            _fd2 = {'type': 'solid', 'color': _stops2[0]['color']}
                                            _grad2 = {'angle': _ang2, 'stops': _stops2}
                                if _fd2.get('type') == 'solid' or _grad2:
                                    _bs2 = {'bgColor': _fd2.get('color', 'transparent')}
                                    if _grad2:
                                        _bs2['bgGradient'] = _grad2
                                    boxes.append({
                                        'id': _sec2.token_hex(4),
                                        'x': lx, 'y': ly,
                                        'w': max(1, int(sw_emu * sx)),
                                        'h': max(1, int(sh_emu * sy)),
                                        'rotation': 0, 'shapeType': 'rect',
                                        'text': '', 'styles': [],
                                        'boxStyle': _bs2,
                                    })
                        except Exception:
                            pass
                except Exception:
                    pass

                for shape in _iter_shapes_flat(slide.shapes):
                    try:
                        box_id = _sec2.token_hex(4)
                        x = int((shape.left or 0) * sx)
                        y = int((shape.top or 0) * sy)
                        w = max(1, int((shape.width or 100) * sx))
                        h = max(1, int((shape.height or 50) * sy))
                        rotation = float(getattr(shape, 'rotation', 0) or 0)

                        # ── Image shapes (type=13 or picture placeholder type=14) ──
                        _is_picture = shape.shape_type == 13
                        if not _is_picture and not shape.has_text_frame:
                            try:
                                _blob_test = shape.image.blob
                                _is_picture = True
                            except Exception:
                                pass
                        if _is_picture:
                            try:
                                img_bytes = shape.image.blob
                                ct = shape.image.content_type or 'image/png'
                                img_b64 = _b64.b64encode(img_bytes).decode()
                                boxes.append({
                                    'id': box_id, 'x': x, 'y': y, 'w': w, 'h': h,
                                    'rotation': rotation,
                                    'imageData': f'data:{ct};base64,{img_b64}',
                                    'shapeType': 'rect',
                                    'text': '', 'styles': [], 'boxStyle': {},
                                })
                            except Exception:
                                pass
                            continue

                        # ── Pure geometric shapes (no text frame) ──────────
                        if not shape.has_text_frame:
                            _st, _cr = _extract_geometry(shape)
                            _fd = {'type': 'none'}
                            try:
                                _fd = _get_fill(shape.fill, shape._element)
                            except Exception:
                                pass
                            if _fd.get('type') != 'none' or _st != 'rect':
                                _bs = {'bgColor': 'transparent'}
                                if _fd.get('type') in ('solid', 'gradient') and _fd.get('color'):
                                    _bs['bgColor'] = _fd['color']
                                if _fd.get('type') == 'gradient' and _fd.get('gradient'):
                                    _bs['bgGradient'] = _fd['gradient']
                                # Read border/line
                                try:
                                    _spPr = shape._element.find(_qn('p:spPr'))
                                    if _spPr is not None:
                                        _ln = _spPr.find(_qn('a:ln'))
                                        if _ln is not None and _ln.find(_qn('a:noFill')) is None:
                                            _ln_w_emu = int(_ln.get('w', '0') or '0')
                                            if _ln_w_emu > 0:
                                                _ln_sol = _ln.find(_qn('a:solidFill'))
                                                _ln_color = '#000000'
                                                if _ln_sol is not None:
                                                    _lc = _color_from_solid_elem(_ln_sol)
                                                    if _lc:
                                                        _ln_color = _lc
                                                _bs['borderColor'] = _ln_color
                                                _bs['borderWidth'] = round(_ln_w_emu / 12700, 1)
                                except Exception:
                                    pass
                                boxes.append({
                                    'id': box_id, 'x': x, 'y': y, 'w': w, 'h': h,
                                    'rotation': rotation,
                                    'shapeType': _st,
                                    'cornerRadius': _cr,
                                    'fill': _fd,
                                    'text': '', 'styles': [], 'boxStyle': _bs,
                                })
                            continue

                        # ── Geometry ───────────────────────────────────────
                        shape_type, corner_radius = _extract_geometry(shape)

                        # ── Fill ──────────────────────────────────────────
                        fill_dict = {'type': 'none'}
                        try:
                            fill_dict = _get_fill(shape.fill, shape._element)
                        except Exception:
                            pass

                        # ── Rich text paragraphs ──────────────────────────
                        tf = shape.text_frame
                        paras_out = []
                        box_style: dict = {'bgColor': 'transparent'}
                        first_done = False

                        for para in tf.paragraphs:
                            # Alignment
                            align_str = 'left'
                            try:
                                from pptx.enum.text import PP_ALIGN
                                al = para.alignment
                                if al == PP_ALIGN.CENTER:
                                    align_str = 'center'
                                elif al == PP_ALIGN.RIGHT:
                                    align_str = 'right'
                            except Exception:
                                pass

                            # Space before (EMU → points: 1pt = 12700 EMU)
                            space_before = 0
                            try:
                                sb = para.space_before
                                if sb is not None:
                                    space_before = round(float(sb) / 12700, 1)
                            except Exception:
                                pass

                            runs_out = []
                            para_text_parts = []
                            for run in para.runs:
                                rt = run.text or ''
                                if not rt:
                                    continue
                                para_text_parts.append(rt)
                                rs: dict = {'text': rt}
                                try:
                                    if run.font.size:
                                        rs['fontSize'] = int(run.font.size / 12700)
                                except Exception:
                                    pass
                                try:
                                    if run.font.bold:
                                        rs['bold'] = True
                                except Exception:
                                    pass
                                try:
                                    if run.font.italic:
                                        rs['italic'] = True
                                except Exception:
                                    pass
                                try:
                                    if run.font.underline:
                                        rs['underline'] = True
                                except Exception:
                                    pass
                                try:
                                    fn = run.font.name
                                    if not fn:
                                        # run.font.name misses <a:latin> set via XML directly
                                        _rpr = run._r.find(_qn('a:rPr'))
                                        if _rpr is not None:
                                            _lat = _rpr.find(_qn('a:latin'))
                                            if _lat is not None:
                                                _tf = _lat.get('typeface', '')
                                                # skip theme-font aliases (+mj-lt / +mn-lt)
                                                if _tf and not _tf.startswith('+'):
                                                    fn = _tf
                                    if fn:
                                        rs['fontFamily'] = fn
                                except Exception:
                                    pass
                                try:
                                    # spc = character spacing in hundredths of a point (can be negative)
                                    _rpr2 = run._r.find(_qn('a:rPr'))
                                    if _rpr2 is not None:
                                        _spc = _rpr2.get('spc')
                                        if _spc is not None:
                                            rs['letterSpacing'] = int(_spc) / 100
                                except Exception:
                                    pass
                                try:
                                    c = _resolve_color(run.font.color)
                                    if c:
                                        rs['color'] = c
                                except Exception:
                                    pass
                                runs_out.append(rs)

                                # Build boxStyle from first run
                                if not first_done:
                                    if 'fontSize' in rs:
                                        box_style['fontSize'] = rs['fontSize']
                                    if rs.get('bold'):
                                        box_style['fontWeight'] = 'bold'
                                    if rs.get('italic'):
                                        box_style['fontStyle'] = 'italic'
                                    if 'color' in rs:
                                        box_style['color'] = rs['color']
                                    if 'fontFamily' in rs:
                                        box_style['fontFamily'] = rs['fontFamily']
                                    box_style['textAlign'] = align_str
                                    first_done = True

                            para_text = ''.join(para_text_parts)
                            if para_text or runs_out:
                                paras_out.append({
                                    'text': para_text,
                                    'align': align_str,
                                    'spaceBefore': space_before,
                                    'runs': runs_out,
                                })

                        # Fill color to boxStyle bgColor
                        if fill_dict.get('type') in ('solid', 'gradient') and fill_dict.get('color'):
                            box_style['bgColor'] = fill_dict['color']
                        if fill_dict.get('type') == 'gradient' and fill_dict.get('gradient'):
                            box_style['bgGradient'] = fill_dict['gradient']

                        # Border/line from shape XML
                        try:
                            _spPr_t = shape._element.find(_qn('p:spPr'))
                            if _spPr_t is not None:
                                _ln_t = _spPr_t.find(_qn('a:ln'))
                                if _ln_t is not None and _ln_t.find(_qn('a:noFill')) is None:
                                    _ln_w_emu_t = int(_ln_t.get('w', '0') or '0')
                                    if _ln_w_emu_t > 0:
                                        _ln_sol_t = _ln_t.find(_qn('a:solidFill'))
                                        _ln_color_t = '#000000'
                                        if _ln_sol_t is not None:
                                            _lct = _color_from_solid_elem(_ln_sol_t)
                                            if _lct:
                                                _ln_color_t = _lct
                                        box_style['borderColor'] = _ln_color_t
                                        box_style['borderWidth'] = round(_ln_w_emu_t / 12700, 1)
                        except Exception:
                            pass

                        # Text vertical anchor and internal margins from bodyPr
                        try:
                            _txBody_t = shape._element.find(_qn('p:txBody'))
                            if _txBody_t is not None:
                                _bpr_t = _txBody_t.find(_qn('a:bodyPr'))
                                if _bpr_t is not None:
                                    _anchor = _bpr_t.get('anchor', '')
                                    if _anchor:
                                        box_style['textAnchor'] = _anchor
                                    # Store internal margins as canvas pixels (default 91440 EMU = 0.1in each side)
                                    _l = int(_bpr_t.get('lIns', '91440') or '91440')
                                    _r = int(_bpr_t.get('rIns', '91440') or '91440')
                                    _t = int(_bpr_t.get('tIns', '45720') or '45720')
                                    _b = int(_bpr_t.get('bIns', '45720') or '45720')
                                    box_style['padL'] = round(_l * sx, 2)
                                    box_style['padR'] = round(_r * sx, 2)
                                    box_style['padT'] = round(_t * sy, 2)
                                    box_style['padB'] = round(_b * sy, 2)
                                    # normAutofit: box was sized to fit text, no soft word-wrap needed
                                    if _bpr_t.find(_qn('a:normAutofit')) is not None:
                                        box_style['normAutofit'] = True
                        except Exception:
                            pass

                        # Apply layout placeholder defaults for missing styles
                        try:
                            if shape.shape_type == 14:
                                ld = _get_layout_ph_defaults(shape)
                                if ld:
                                    if 'textAnchor' in ld:
                                        box_style['textAnchor'] = ld['textAnchor']
                                    if 'fontSize' in ld and 'fontSize' not in box_style:
                                        box_style['fontSize'] = ld['fontSize']
                                    if ld.get('bold') and 'fontWeight' not in box_style:
                                        box_style['fontWeight'] = 'bold'
                                    if ld.get('allCaps'):
                                        box_style['allCaps'] = True
                                    if 'color' in ld and 'color' not in box_style:
                                        box_style['color'] = ld['color']
                        except Exception:
                            pass

                        full_text = '\n'.join(p['text'] for p in paras_out)
                        boxes.append({
                            'id': box_id, 'x': x, 'y': y, 'w': w, 'h': h,
                            'rotation': rotation,
                            'shapeType': shape_type,
                            'cornerRadius': corner_radius,
                            'fill': fill_dict,
                            'text': full_text,
                            'paragraphs': paras_out,
                            'styles': [],
                            'boxStyle': box_style,
                        })
                    except Exception:
                        pass

                slide_width_pt = int(sw / 12700)
                # Pre-compute text layout (line breaks, positions, overflow) server-side
                try:
                    from ppt_layout import compute_slide_layouts
                    compute_slide_layouts(boxes, slide_width_pt)
                except Exception as _le:
                    import sys as _sys
                    print(f'[ppt_layout] skipped: {_le}', file=_sys.stderr)
                slide_out = {'bgColor': bg_color, 'boxes': boxes,
                            'slideWidthPt': slide_width_pt}
                if bg_gradient:
                    slide_out['bgGradient'] = bg_gradient
                slides_out.append(slide_out)
            # Write sidecar atomically so future loads are instant and edits persist
            import json as _json_w, tempfile as _tmpw, shutil as _shuw, os as _osw
            try:
                _tfd, _tname = _tmpw.mkstemp(dir=p.parent)
                with _osw.fdopen(_tfd, 'w') as _tf:
                    _tf.write(_json_w.dumps({"slides": slides_out}))
                _shuw.move(_tname, str(json_path))
            except Exception:
                pass
            return {"slides": slides_out}
        except Exception as exc:
            raise HTTPException(500, str(exc))

    @misc_router.get("/api/ppt/status")
    async def ppt_status(path: str = Query(...)):
        p = Path(path)
        json_path = p.parent / f".{p.name}.json"
        return {"ready": json_path.exists(), "hasSidecar": json_path.exists()}

    class CheckFitBody(BaseModel):
        path: str                # absolute path to .pptx
        slide: int = 0           # 0-based slide index
        box_id: str = ""         # box id from sidecar (empty → first text box)
        text: str = ""           # proposed text (newlines → paragraph breaks)

    @misc_router.post("/api/ppt/checkfit")
    async def ppt_checkfit(body: CheckFitBody):
        """
        Check whether `text` fits inside a specific box on a slide.
        Uses the pre-computed sidecar for box geometry and style, then
        re-runs the Pillow layout engine with the new text.
        Returns: fits, lines, textHeight, boxHeight, overflowBy (canvas px), layoutLines.
        """
        import json as _jf
        from ppt_layout import checkfit as _checkfit
        p = Path(body.path)
        json_path = p.parent / f".{p.name}.json"
        if not json_path.exists():
            raise HTTPException(404, "sidecar not found — open the PPTX first")
        sidecar = _jf.loads(json_path.read_text())
        slides = sidecar.get("slides", [])
        if body.slide >= len(slides):
            raise HTTPException(400, f"slide {body.slide} out of range")
        slide = slides[body.slide]
        slide_width_pt = slide.get("slideWidthPt", 720)
        boxes = slide.get("boxes", [])
        box = None
        if body.box_id:
            box = next((b for b in boxes if b.get("id") == body.box_id), None)
        if box is None:
            # Fall back to first box that has paragraphs
            box = next((b for b in boxes if b.get("paragraphs")), None)
        if box is None:
            raise HTTPException(404, "no text box found on that slide")
        result = _checkfit(box, body.text, slide_width_pt)
        return result

    @misc_router.put("/api/ppt/save")
    async def ppt_save(body: PptSaveBody):
        if not body.path:
            raise HTTPException(400, "path required")
        if not body.slides:
            raise HTTPException(400, "slides cannot be empty")
        try:
            import json as _json, tempfile, shutil as _shutil, os as _os
            from pptx import Presentation
            from pptx.util import Emu, Pt
            from pptx.dml.color import RGBColor

            CW, CH = 896, 504
            p = Path(body.path)
            json_path = p.parent / f".{p.name}.json"
            p.parent.mkdir(parents=True, exist_ok=True)

            def hex_to_rgb(hex_str: str):
                h = hex_str.lstrip("#")
                if len(h) == 6:
                    try:
                        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
                    except Exception:
                        pass
                return None

            # 1. Write JSON sidecar atomically (ground truth)
            tmp_fd, tmp_name = tempfile.mkstemp(suffix=".json", dir=p.parent)
            try:
                _os.close(tmp_fd)
                Path(tmp_name).write_text(_json.dumps({"slides": body.slides}))
                _shutil.move(tmp_name, str(json_path))
            except Exception:
                try: _os.unlink(tmp_name)
                except OSError: pass
                raise

            # 2. Re-export to .pptx for compatibility
            _pptx_export(p, body.slides)
            return {"ok": True, "path": str(p)}
        except Exception as exc:
            raise HTTPException(500, str(exc))

    @misc_router.post("/api/dev-server/start")
    async def dev_server_start(body: DevServerBody, request: Request):
        cwd = body.cwd or app.state.default_cwd
        if not cwd or not Path(cwd).is_dir():
            raise HTTPException(400, "invalid cwd")
        if cwd in _dev_servers:
            entry = _dev_servers[cwd]
            if entry["proc"].returncode is None:
                return {"url": entry["url"]}
        port = await _find_free_port()
        host_header = request.headers.get("host", "")
        if ".workspaces.boltzhub.com" in host_header:
            workspace_id = host_header.split(".")[0]
            url = f"https://{workspace_id}-{port}.workspaces.boltzhub.com"
        else:
            url = f"http://localhost:{port}"
        pkg_dir = Path(cwd)
        if (pkg_dir / "pnpm-lock.yaml").exists():
            cmd = ["pnpm", "dev", "--port", str(port), "--host", "0.0.0.0"]
        elif (pkg_dir / "yarn.lock").exists():
            cmd = ["yarn", "dev", "--port", str(port), "--host", "0.0.0.0"]
        else:
            cmd = ["npm", "run", "dev", "--", "--port", str(port), "--host", "0.0.0.0"]
        import os as _os
        import shutil as _shutil
        env = _os.environ.copy()
        extra_paths = [
            "/usr/local/bin", "/usr/bin",
            str(Path.home() / ".local/node/bin"),
            str(Path.home() / ".local/share/pnpm"),
            str(Path.home() / ".nvm/versions/node/current/bin"),
            "/root/.local/share/pnpm",
            "/root/.local/node/bin",
        ]
        env["PATH"] = ":".join(extra_paths) + ":" + env.get("PATH", "")
        resolved = _shutil.which(cmd[0], path=env["PATH"])
        if resolved:
            cmd = [resolved] + cmd[1:]
        print(f"[dev-server] starting {' '.join(cmd)} in {cwd}", file=sys.stderr)
        try:
            proc = await asyncio.create_subprocess_exec(
                *cmd, cwd=cwd,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.STDOUT,
                env=env,
            )
        except FileNotFoundError as exc:
            print(f"[dev-server] command not found: {exc}", file=sys.stderr)
            raise HTTPException(500, f"command not found: {exc}")

        async def _pipe_output(p: asyncio.subprocess.Process, label: str) -> None:
            if p.stdout is None:
                return
            try:
                async for raw in p.stdout:
                    line = raw.decode(errors="replace").rstrip()
                    print(f"[{label}] {line}", file=sys.stderr)
            except Exception:
                pass

        asyncio.create_task(_pipe_output(proc, f"dev-server:{Path(cwd).name}"))
        _dev_servers[cwd] = {"proc": proc, "url": url}
        await asyncio.sleep(2)
        if proc.returncode is not None:
            print(f"[dev-server] exited immediately (rc={proc.returncode})", file=sys.stderr)
            _dev_servers.pop(cwd, None)
            raise HTTPException(500, "dev server exited immediately — check package.json")
        print(f"[dev-server] started pid={proc.pid} url={url}", file=sys.stderr)
        return {"url": url, "pid": proc.pid}

    @misc_router.post("/api/dev-server/stop")
    async def dev_server_stop(body: DevServerBody):
        cwd = body.cwd
        entry = _dev_servers.pop(cwd, None)
        if entry:
            print(f"[dev-server] stopping pid={entry['proc'].pid} cwd={cwd}", file=sys.stderr)
            try:
                entry["proc"].terminate()
            except Exception:
                pass
        return {"ok": True}

    # ── § 6 · Database ────────────────────────────────────────────────────────

    @db_router.get("/health")
    async def db_health(request: Request):
        pool = getattr(request.app.state, "db", None)
        if pool is None:
            return JSONResponse({"ok": False, "error": "Database not connected"}, status_code=503)
        try:
            async with pool.acquire() as conn:
                version = await conn.fetchval("SELECT version()")
            return {"ok": True, "version": version}
        except Exception as exc:
            return JSONResponse({"ok": False, "error": str(exc)}, status_code=503)

    @db_router.get("/widget/{canvas_id}/schema")
    async def widget_schema_get(canvas_id: str = FPath(...)):
        if not _CANVAS_ID_RE.match(canvas_id):
            raise HTTPException(400, f"Invalid canvasId: {canvas_id!r}")
        data = _widget_load(canvas_id)
        return {"columns": data.get("_schema", []), "rowCount": len(data.get("records", []))}

    @db_router.post("/widget/{canvas_id}/schema")
    async def widget_schema_ensure(canvas_id: str, body: WidgetSchemaBody):
        if not _CANVAS_ID_RE.match(canvas_id):
            raise HTTPException(400, f"Invalid canvasId: {canvas_id!r}")
        with _widget_lock(canvas_id):
            data = _widget_load(canvas_id)
            existing = {c["name"]: c for c in data.get("_schema", []) if "name" in c}
            for col in body.columns:
                if "name" in col and col["name"] not in existing:
                    existing[col["name"]] = col
            data["_schema"] = list(existing.values())
            _widget_save(canvas_id, data)
        return {"columns": data["_schema"], "rowCount": len(data.get("records", []))}

    @db_router.get("/widget/{canvas_id}/rows")
    async def widget_query(canvas_id: str = FPath(...),
                           order: str = Query("id"), dir: str = Query("asc"),
                           limit: int = Query(1000), offset: int = Query(0)):
        if not _CANVAS_ID_RE.match(canvas_id):
            raise HTTPException(400, f"Invalid canvasId: {canvas_id!r}")
        data = _widget_load(canvas_id)
        records = data["records"]
        desc = dir.upper() == "DESC"
        records = sorted(records, key=lambda r: r.get(order, 0), reverse=desc)
        page = records[offset: offset + min(limit, 10000)]
        return {"rows": page, "total": len(records), "limit": limit, "offset": offset}

    @db_router.post("/widget/{canvas_id}/rows")
    async def widget_insert(canvas_id: str, body: WidgetRowBody):
        if not _CANVAS_ID_RE.match(canvas_id):
            raise HTTPException(400, f"Invalid canvasId")
        import datetime as _dt
        rows = body.rows or ([body.row] if body.row else [])
        if not rows:
            raise HTTPException(400, "Provide 'row' or 'rows'")
        with _widget_lock(canvas_id):
            data = _widget_load(canvas_id)
            inserted = []
            for row in rows:
                row = {k: v for k, v in row.items() if k not in ("id", "created_at")}
                row["id"] = data["_next_id"]
                row["created_at"] = _dt.datetime.utcnow().isoformat() + "Z"
                data["_next_id"] += 1
                data["records"].append(row)
                inserted.append(row)
            _widget_save(canvas_id, data)
        return {"inserted": inserted}

    @db_router.put("/widget/{canvas_id}/rows/{row_id}")
    async def widget_update(canvas_id: str, row_id: int, body: WidgetUpdateBody):
        if not _CANVAS_ID_RE.match(canvas_id):
            raise HTTPException(400, "Invalid canvasId")
        patch = {k: v for k, v in body.data.items() if k not in ("id", "created_at")}
        if not patch:
            raise HTTPException(400, "'data' required")
        with _widget_lock(canvas_id):
            data = _widget_load(canvas_id)
            for r in data["records"]:
                if r.get("id") == row_id:
                    r.update(patch)
                    _widget_save(canvas_id, data)
                    return {"updated": r}
        raise HTTPException(404, "Row not found")

    @db_router.delete("/widget/{canvas_id}/rows/{row_id}")
    async def widget_delete(canvas_id: str, row_id: int):
        if not _CANVAS_ID_RE.match(canvas_id):
            raise HTTPException(400, "Invalid canvasId")
        with _widget_lock(canvas_id):
            data = _widget_load(canvas_id)
            before = len(data["records"])
            data["records"] = [r for r in data["records"] if r.get("id") != row_id]
            if len(data["records"]) == before:
                raise HTTPException(404, "Row not found")
            _widget_save(canvas_id, data)
        return {"deleted": row_id}

    @db_router.post("/widget/{canvas_id}/exec")
    async def widget_exec(canvas_id: str, body: WidgetExecBody):
        if not _CANVAS_ID_RE.match(canvas_id):
            raise HTTPException(400, "Invalid canvasId")
        data = _widget_load(canvas_id)
        ns = {"records": data["records"], "result": None}
        try:
            exec(compile(body.code, "<widget-exec>", "exec"), ns)  # nosec
        except Exception as exc:
            raise HTTPException(400, str(exc))
        return {"result": ns.get("result")}

    # ── § 8 · Batch Execution ─────────────────────────────────────────────────

    @batch_router.post("/batch")
    async def batch_run(body: BatchRunBody):
        import uuid as _uuid, time as _t
        if not body.cwds or not body.message:
            raise HTTPException(400, "cwds and message required")
        _bzcode = app.state.bzcode_path
        batch_id = _uuid.uuid4().hex[:12]
        items = [_BatchItem(cwd, _bzcode, resume_session_id=body.sessions.get(cwd, ""))
                 for cwd in body.cwds]
        _batch_store[batch_id] = {"items": items, "created": _t.time()}

        async def _run():
            await asyncio.gather(*[item.run(body.message) for item in items],
                                 return_exceptions=True)

        asyncio.create_task(_run())
        return {"batchId": batch_id}

    @batch_router.get("/batch/{batch_id}")
    async def batch_status(batch_id: str = FPath(...)):
        batch = _batch_store.get(batch_id)
        if not batch:
            raise HTTPException(404, "not found")
        items = [item.to_dict() for item in batch["items"]]
        done  = all(i["status"] in ("done", "error") for i in items)
        return {"batchId": batch_id, "done": done, "items": items}

    # ── § 7 · BoltzHub ────────────────────────────────────────────────────────

    @boltzhub_router.get("/check")
    async def boltzhub_check(cwd: str = Query("")):
        if not cwd:
            cwd = app.state.default_cwd
        token     = _boltzhub_token()
        cfg       = _read_app_config(cwd)
        bzhub_dir = Path(cwd) / ".bzhub"
        return {
            "isLoggedIn":   bool(token),
            "hasAppConfig": bool(cfg),
            "appConfig":    cfg,
            "hasBzhubDir":  bzhub_dir.is_dir(),
            "configPath":   str(bzhub_dir / "app_config.json"),
            "dirName":      Path(cwd).name,
            "cwd":          cwd,
        }

    @boltzhub_router.post("/create-app")
    async def boltzhub_create_app(body: CreateAppBody):
        import aiohttp as _aio
        token = _boltzhub_token()
        if not token:
            raise HTTPException(401, "Not logged in to BoltzHub")
        cwd = body.cwd or app.state.default_cwd
        api_body = {"name": body.name, "visibility": body.visibility}
        if body.description:
            api_body["description"] = body.description
        if body.priceMonthly:
            api_body["priceMonthly"] = body.priceMonthly
        connector = _aio.TCPConnector(ssl=False)
        async with _aio.ClientSession(connector=connector) as sess:
            async with sess.post(f"{BOLTZHUB_API}/v1/creator/apps",
                                 json=api_body, headers=_bz_headers(token)) as resp:
                if resp.status not in (200, 201):
                    raise HTTPException(resp.status, await resp.text())
                result = await resp.json()
        cfg = {"id": result["id"], "name": result["name"],
               "description": result.get("description"),
               "visibility": result.get("visibility", "private"),
               "buildCommand": body.buildCommand,
               "createdAt": result.get("createdAt")}
        _write_app_config(cwd, cfg)
        return {"ok": True, "appConfig": cfg}

    @boltzhub_router.post("/push")
    async def boltzhub_push(body: PushBody):
        """SSE streaming response for push progress."""
        import aiohttp as _aio, zipfile as _zf, io as _io
        cwd = body.cwd or app.state.default_cwd
        token = _boltzhub_token()

        async def _stream():
            def emit(step, message, **kw):
                return f"data: {json.dumps({'step':step,'message':message,**kw})}\n\n"
            try:
                if not token:
                    yield emit("error", "Not logged in to BoltzHub"); return
                cfg = _read_app_config(cwd)
                if not cfg:
                    yield emit("error", "No .bzhub/app_config.json found"); return
                app_id    = cfg["id"]
                _sync_env_oauth_client_id(cwd, app_id)
                build_cmd = cfg.get("buildCommand") or "pnpm build"
                yield emit("build", f"Running: {build_cmd}")
                proc = await asyncio.create_subprocess_shell(
                    build_cmd, cwd=cwd,
                    stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE)
                _, stderr = await proc.communicate()
                if proc.returncode != 0:
                    yield emit("error", f"Build failed: {stderr.decode()[:300]}"); return
                yield emit("archive", "Archiving project…")
                bzhub_dir = Path(cwd) / ".bzhub"
                bzhub_dir.mkdir(parents=True, exist_ok=True)
                zip_path = bzhub_dir / "project.zip"
                if zip_path.exists():
                    zip_path.unlink()
                zip_cmd = f'cd "{cwd}" && zip -r "{zip_path}" . -x "node_modules/*" ".bzhub/*" ".git/*"'
                proc2 = await asyncio.create_subprocess_shell(
                    zip_cmd, stdout=asyncio.subprocess.DEVNULL, stderr=asyncio.subprocess.PIPE)
                _, zip_err = await proc2.communicate()
                if proc2.returncode not in (0, 12):
                    yield emit("error", f"Archive failed: {zip_err.decode()[:300]}"); return
                zip_bytes = zip_path.read_bytes()
                auth = _bz_auth(token)
                connector = _aio.TCPConnector(ssl=False)
                async with _aio.ClientSession(connector=connector) as sess:
                    yield emit("upload", f"Uploading {len(zip_bytes)//1024} KB…")
                    form = _aio.FormData()
                    form.add_field("archiveFile", zip_bytes, filename="project.zip",
                                   content_type="application/zip")
                    async with sess.post(f"{BOLTZHUB_API}/v1/creator/apps/{app_id}/code",
                                         data=form, headers=auth) as r:
                        if r.status not in (200, 201):
                            yield emit("error", f"Upload failed ({r.status}): {await r.text()}"); return
                    yield emit("deploy", "Deploying…")
                    async with sess.put(f"{BOLTZHUB_API}/v1/creator/apps/{app_id}/deploy",
                                        headers=auth) as r:
                        if r.status not in (200, 201):
                            yield emit("error", f"Deploy trigger failed: {await r.text()}"); return
                    service_url = None
                    for attempt in range(60):
                        if attempt:
                            await asyncio.sleep(5)
                        async with sess.get(f"{BOLTZHUB_API}/v1/creator/apps/{app_id}/status",
                                            headers=auth) as r:
                            if r.status != 200:
                                continue
                            st = await r.json()
                            service_url = st.get("serviceUrl")
                            dep_status  = st.get("status")
                            yield emit("deploy", st.get("stepMessage", f"Deploying… ({attempt*5}s)"))
                            if dep_status == "deployed":
                                break
                            if dep_status == "failed":
                                yield emit("error", "Deployment failed"); return
                    yield emit("publish", "Publishing version…")
                    if body.releaseNotes:
                        async with sess.post(
                            f"{BOLTZHUB_API}/v1/creator/apps/{app_id}/versions",
                            json={"releaseNotes": body.releaseNotes, "versionNumber": body.versionNumber},
                            headers=auth,
                        ) as r:
                            pass
                yield emit("done", "Deployed!", serviceUrl=service_url or "", appId=app_id)
            except Exception as exc:
                yield emit("error", str(exc))

        return StreamingResponse(_stream(),
                                 media_type="text/event-stream",
                                 headers={"Cache-Control": "no-cache",
                                          "X-Accel-Buffering": "no"})

    @boltzhub_router.get("/apps")
    async def boltzhub_apps():
        import aiohttp as _aio
        token = _boltzhub_token()
        if not token:
            raise HTTPException(401, "Not logged in")
        connector = _aio.TCPConnector(ssl=False)
        async with _aio.ClientSession(connector=connector) as sess:
            async with sess.get(f"{BOLTZHUB_API}/v1/creator/apps",
                                 headers=_bz_auth(token)) as r:
                if r.status != 200:
                    raise HTTPException(r.status, await r.text())
                return await r.json()

    @boltzhub_router.get("/versions")
    async def boltzhub_versions(appId: str = Query(...)):
        import aiohttp as _aio
        token = _boltzhub_token()
        if not token:
            raise HTTPException(401, "Not logged in")
        connector = _aio.TCPConnector(ssl=False)
        async with _aio.ClientSession(connector=connector) as sess:
            async with sess.get(f"{BOLTZHUB_API}/v1/creator/apps/{appId}/versions",
                                 headers=_bz_auth(token)) as r:
                if r.status != 200:
                    raise HTTPException(r.status, await r.text())
                data = await r.json()
                items = data if isinstance(data, list) else data.get("items", [])
                items.sort(key=lambda v: v.get("createdAt", ""), reverse=True)
                latest = items[0]["versionNumber"] if items else "0.0.0"
                parts  = latest.split(".")
                try:
                    suggested = f"{parts[0]}.{parts[1]}.{int(parts[2])+1}"
                except Exception:
                    suggested = "1.0.0"
                return {"versions": items, "suggestedNext": suggested}

    @boltzhub_router.get("/token-usage")
    async def boltzhub_token_usage(period: str = Query("30d")):
        import aiohttp as _aio
        token = _boltzhub_token()
        if not token:
            raise HTTPException(401, "Not logged in")
        connector = _aio.TCPConnector(ssl=False)
        async with _aio.ClientSession(connector=connector) as sess:
            async with sess.get(
                f"{BOLTZHUB_API}/v1/creator/tokens/usage/history?period={period}&limit=100",
                headers=_bz_auth(token),
            ) as r:
                if r.status != 200:
                    raise HTTPException(r.status, await r.text())
                return await r.json()

    @boltzhub_router.post("/sync")
    async def boltzhub_sync(body: BzHubSyncBody):
        import aiohttp as _aio, io, zipfile as _zf
        cwd    = body.cwd or app.state.default_cwd
        app_id = body.appId
        token  = _boltzhub_token()

        async def _stream():
            def emit(step: str, message: str, **kw):
                return f"data: {json.dumps({'step': step, 'message': message, **kw})}\n\n"
            try:
                if not token:
                    yield emit("error", "Not logged in to BoltzHub"); return
                if not app_id:
                    cfg = _read_app_config(cwd)
                    if not cfg:
                        yield emit("error", "No .bzhub/app_config.json found"); return
                    _app_id = cfg["id"]
                else:
                    _app_id = app_id
                connector = _aio.TCPConnector(ssl=False)
                yield emit("download", "Downloading project…")
                async with _aio.ClientSession(connector=connector) as sess:
                    async with sess.get(
                        f"{BOLTZHUB_API}/v1/creator/apps/{_app_id}/code",
                        headers=_bz_auth(token),
                    ) as r:
                        if r.status != 200:
                            yield emit("error", f"Download failed ({r.status})"); return
                        zip_bytes = await r.read()
                yield emit("extract", "Extracting files…")
                buf = io.BytesIO(zip_bytes)
                with _zf.ZipFile(buf) as z:
                    z.extractall(cwd)
                yield emit("install", "Installing dependencies…")
                lock_pnpm = (Path(cwd) / "pnpm-lock.yaml").exists()
                install_cmd = "pnpm install" if lock_pnpm else "npm install"
                if (Path(cwd) / "package.json").exists():
                    proc = await asyncio.create_subprocess_shell(
                        install_cmd, cwd=cwd,
                        stdout=asyncio.subprocess.DEVNULL, stderr=asyncio.subprocess.DEVNULL)
                    await proc.wait()
                yield emit("done", "Project synced successfully!")
            except Exception as exc:
                yield emit("error", str(exc))

        return StreamingResponse(_stream(), media_type="text/event-stream",
                                 headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})

    @boltzhub_router.post("/create-version")
    async def boltzhub_create_version(body: BzHubVersionBody):
        import aiohttp as _aio
        token = _boltzhub_token()
        if not token:
            raise HTTPException(401, "Not logged in")
        connector = _aio.TCPConnector(ssl=False)
        async with _aio.ClientSession(connector=connector) as sess:
            async with sess.post(
                f"{BOLTZHUB_API}/v1/creator/apps/{body.appId}/versions",
                json={"releaseNotes": body.releaseNotes, "versionNumber": body.versionNumber},
                headers=_bz_auth(token),
            ) as r:
                result = await r.json() if r.content_type == "application/json" else {"status": r.status}
                if r.status not in (200, 201):
                    raise HTTPException(r.status, str(result))
                return result

    @boltzhub_router.post("/publish")
    async def boltzhub_publish(body: BzHubPublishBody):
        import aiohttp as _aio
        token = _boltzhub_token()
        if not token:
            raise HTTPException(401, "Not logged in")
        connector = _aio.TCPConnector(ssl=False)
        async with _aio.ClientSession(connector=connector) as sess:
            async with sess.put(
                f"{BOLTZHUB_API}/v1/creator/apps/{body.appId}/publish",
                headers=_bz_headers(token),
            ) as r:
                result = await r.json() if r.content_type == "application/json" else {"status": r.status}
                if r.status not in (200, 201):
                    raise HTTPException(r.status, str(result))
                return result

    # ── § 9 · WhatsApp ────────────────────────────────────────────────────────

    @whatsapp_router.post("/incoming")
    async def whatsapp_incoming(request: Request):
        data   = await request.form()
        from_  = data.get("From", "").strip()
        body_t = data.get("Body", "").strip()
        if not from_ or not body_t:
            return "<Response/>"
        _bzcode = app.state.bzcode_path
        _cwd    = app.state.default_cwd
        whatsapp_dir = Path(_cwd) / "whatsapp"
        whatsapp_dir.mkdir(parents=True, exist_ok=True)
        async with _whatsapp_lock:
            if from_ not in _whatsapp_sessions:
                _whatsapp_sessions[from_] = _WASess(from_, _bzcode, str(whatsapp_dir))
            sess = _whatsapp_sessions[from_]
        from server import _send_whatsapp, _load_creds
        async def _process():
            reply = await sess.chat(body_t)
            await _send_whatsapp(from_, reply, _load_creds())
        asyncio.create_task(_process())
        from fastapi.responses import Response
        return Response(content="<Response/>", media_type="text/xml")

    @whatsapp_router.post("/status")
    async def whatsapp_status(request: Request):
        data   = await request.form()
        status = data.get("MessageStatus", "unknown")
        sid    = data.get("MessageSid", "")
        print(f"[whatsapp] delivery {sid}: {status}", file=sys.stderr)
        from fastapi.responses import Response
        return Response(content="<Response/>", media_type="text/xml")

    # ── Mount all routers ─────────────────────────────────────────────────────
    for _router in [ws_router, auth_router, files_router, sessions_router,
                    canvas_router, db_router, boltzhub_router,
                    batch_router, whatsapp_router, misc_router]:
        app.include_router(_router)

    return app


# ── Static file serving helper ────────────────────────────────────────────────

def mount_frontend(app: FastAPI, dist_dir: Path) -> None:
    """Serve the Vite production build as a SPA from the same server."""
    if not dist_dir.is_dir():
        print(f"[frontend] dist dir not found: {dist_dir}", file=sys.stderr)
        return
    index_html = dist_dir / "index.html"
    if not index_html.exists():
        print(f"[frontend] index.html not found — run 'pnpm build' first", file=sys.stderr)
        return

    # Serve /assets and other static sub-directories
    for entry in dist_dir.iterdir():
        if entry.is_dir():
            app.mount(f"/{entry.name}", StaticFiles(directory=str(entry)), name=entry.name)

    # SPA catch-all — must be last
    @app.get("/{full_path:path}", include_in_schema=False)
    async def spa_fallback(full_path: str):
        # Serve real files, fall back to index.html.
        # Guard against OSError ENAMETOOLONG (e.g. gateway injects JWT into URL path).
        try:
            candidate = dist_dir / full_path
            if candidate.is_file():
                return FileResponse(str(candidate))
        except OSError:
            pass
        return FileResponse(str(index_html))

    print(f"[frontend] serving {dist_dir}", file=sys.stderr)


# ── Entry point ───────────────────────────────────────────────────────────────

def main() -> None:
    parser = argparse.ArgumentParser(description="BoltzAgent FastAPI server")
    parser.add_argument("--bzcode",  default="bzcode")
    parser.add_argument("--host",    default="localhost")
    parser.add_argument("--port",    type=int, default=18789)
    parser.add_argument("--cwd",     default=os.getcwd())
    parser.add_argument("--bz-home", default="", dest="bz_home",
                        help="BZ_HOME directory for bzcode credentials, sessions and settings")
    parser.add_argument("--dist",    default="", metavar="DIR",
                        help="Vite dist/ directory to serve as the SPA frontend")
    args = parser.parse_args()

    import shutil as _shutil
    bzcode_path = _shutil.which(args.bzcode) or args.bzcode
    default_cwd = os.path.abspath(args.cwd)
    Path(default_cwd).mkdir(parents=True, exist_ok=True)
    _bz_home_raw = args.bz_home or os.environ.get("BZ_HOME", "")
    bz_home      = os.path.abspath(_bz_home_raw) if _bz_home_raw else ""

    fastapi_app = create_app(bzcode_path=bzcode_path, default_cwd=default_cwd,
                              bz_home=bz_home, port=args.port)

    if args.dist:
        mount_frontend(fastapi_app, Path(args.dist).resolve())

    print(f"BoltzAgent FastAPI server", flush=True)
    print(f"  WebSocket : ws://{args.host}:{args.port}/ws", flush=True)
    print(f"  HTTP API  : http://{args.host}:{args.port}", flush=True)
    print(f"  Docs      : http://{args.host}:{args.port}/docs", flush=True)
    if args.dist:
        print(f"  Frontend  : http://{args.host}:{args.port}/", flush=True)

    # Suppress access-log noise from high-frequency polling endpoints.
    class _SuppressPolling(logging.Filter):
        _MUTED = frozenset(['/sessions', '/api/apikey-status', '/api/apikey-verify'])
        def filter(self, record: logging.LogRecord) -> bool:
            if isinstance(record.args, tuple) and len(record.args) >= 3:
                path = str(record.args[2]).split('?')[0]
                if path in self._MUTED:
                    return False
            return True
    logging.getLogger('uvicorn.access').addFilter(_SuppressPolling())

    uvicorn.run(
        fastapi_app,
        host=args.host,
        port=args.port,
        log_level="info",
    )


# Allow `uvicorn app:app` to work without any arguments (uses defaults).
# Set environment variables to configure without CLI flags:
#   BZCODE_PATH   path to the bzcode binary  (default: bzcode, resolved via PATH)
#   BZCODE_CWD    default working directory   (default: current directory)
#   BZ_HOME       bzcode home dir for credentials/sessions/settings (recommended on servers)
#   PORT          HTTP port                   (default: 18789)
#   BZCODE_DIST   path to Vite dist/ folder   (default: ./dist if it exists)
import re, shutil as _shutil_mod
app = create_app(
    bzcode_path=_shutil_mod.which(os.environ.get("BZCODE_PATH", "bzcode")) or os.environ.get("BZCODE_PATH", "bzcode"),
    default_cwd=os.path.abspath(os.environ.get("BZCODE_CWD", os.getcwd())),
    bz_home=os.environ.get("BZ_HOME", ""),
    port=int(os.environ.get("PORT", "18789")),
)

# Auto-mount the frontend if BZCODE_DIST is set, or if ./dist/index.html exists
_dist_env  = os.environ.get("BZCODE_DIST", "")
_dist_auto = Path("./dist")
_dist_path = Path(_dist_env).resolve() if _dist_env else (_dist_auto.resolve() if (_dist_auto / "index.html").exists() else None)
if _dist_path:
    mount_frontend(app, _dist_path)

if __name__ == "__main__":
    main()
