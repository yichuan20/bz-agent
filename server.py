#!/usr/bin/env python3
"""
Unified server:
  ws://localhost:8765?cwd=/path   — start a fresh bzcode session in that dir
  ws://localhost:8765?sessionId=X — resume an existing session
  http://localhost:8766/sessions  — list sessions (one per directory)
  http://localhost:8766/search    — SerpAPI proxy
"""

BACKEND_VERSION = "0.6.2"

import asyncio
import json
import os
import re
import sys
from datetime import datetime, timezone
from pathlib import Path
from typing import Optional, Tuple

try:
    import websockets  # kept for legacy compatibility; WS now runs inside aiohttp
except ImportError:
    websockets = None  # type: ignore[assignment]

try:
    import aiohttp
    from aiohttp import web
except ImportError:
    sys.exit("Missing dependency: pip install aiohttp")

# ══════════════════════════════════════════════════════════════════════════════
# § 3 · MODE CONFIG
# ══════════════════════════════════════════════════════════════════════════════
# agent_modes.json lives next to server.py. Re-read on every connection so
# edits take effect without a server restart.

_MODES_CONFIG_FILE = Path(__file__).parent / "agent_modes.json"


def _load_mode_config() -> dict:
    try:
        return json.loads(_MODES_CONFIG_FILE.read_text())
    except Exception:
        return {"default": "general", "modes": {}}


def _mode_entry(mode: str) -> dict:
    cfg = _load_mode_config()
    modes = cfg.get("modes", {})
    return modes.get(mode) or modes.get(cfg.get("default", "general"), {}) or {}


def _build_widget_template_table() -> str:
    """Build a markdown table of widget templates from server_data/widgets/index.json.
    Used to inject an up-to-date alias table into the new-widget skill at session-write time."""
    try:
        index_path = SERVER_DATA_DIR / "widgets" / "index.json"
        data = json.loads(index_path.read_text(encoding="utf-8"))
        widgets = [w for w in data.get("widgets", []) if not w.get("archived")]
    except Exception:
        return "(template index unavailable)"

    lines = ["| Template | Matches requests like… | Default size |", "|---|---|---|"]
    for w in widgets:
        name = w.get("id", "")
        label = w.get("label", name)
        keywords = ", ".join(w.get("keywords", [])[:6])
        dw, dh = w.get("defaultW", 380), w.get("defaultH", 280)
        lines.append(f"| `{name}` | {label}: {keywords} | {dw}×{dh} |")
    return "\n".join(lines)


# ══════════════════════════════════════════════════════════════════════════════
# § 4 · SESSION MANAGEMENT
# ══════════════════════════════════════════════════════════════════════════════


def _write_session_config(session_id: str, mode: str, working_dir: str = "", model_name: str = "") -> None:
    """Write IDENTITY.md, SOUL.md, settings.json, and meta.json into the session
    config directory before spawning bzcode.  bzcode picks these up at startup and
    on every --resume, so they are re-applied on reconnect too.
    meta.json is our own metadata (not read by bzcode) used by _read_session_file.
    model_name: reported model from bzcode status (empty = unknown, treated as boltzbit)."""
    import shutil as _shutil

    is_boltzbit = not model_name or model_name.lower().startswith("boltzbit")
    _asset_suffix = "" if is_boltzbit else "_generic"

    entry = _mode_entry(mode)
    cfg_dir = SESSIONS_DIR / session_id
    cfg_dir.mkdir(parents=True, exist_ok=True)

    # Purge sub-agent session files/dirs that bzcode's Agent tool leaves inside
    # the config dir (e.g. cozy-hopping-comet.jsonl, tool-results/).  They are
    # not part of our config and prevent bzcode from resuming cleanly.
    _OWNED_NAMES = {
        "meta.json",
        "IDENTITY.md",
        "SOUL.md",
        "AGENTS.md",
        "settings.json",
        "skills",
        "scripts",
        "templates",
        "custom_widgets",
        "widget_data",
        ".bzcanvas.json",
    }
    for item in list(cfg_dir.iterdir()):
        if item.name not in _OWNED_NAMES:
            if item.is_dir():
                _shutil.rmtree(item, ignore_errors=True)
            else:
                item.unlink(missing_ok=True)

    # Detect model change and log when assets switch variant
    _meta_path = cfg_dir / "meta.json"
    _prev_model = ""
    try:
        _prev_meta = json.loads(_meta_path.read_text(encoding="utf-8"))
        _prev_model = _prev_meta.get("model", "")
    except Exception:
        pass
    _variant = "boltzbit" if is_boltzbit else "generic"
    _prev_variant = "boltzbit" if (not _prev_model or _prev_model.lower().startswith("boltzbit")) else "generic"
    if model_name and _prev_model and _prev_model != model_name:
        print(
            f"[session] {session_id} model changed: {_prev_model!r} → {model_name!r} "
            f"({_prev_variant} → {_variant} assets)",
            file=sys.stderr,
        )
    elif not _prev_model and model_name:
        print(f"[session] {session_id} model set to {model_name!r} — using {_variant} assets", file=sys.stderr)

    # Our own metadata — used by _read_session_file since new bzcode no longer
    # writes a session header line into the .jsonl
    meta = {"sessionId": session_id, "workingDir": working_dir, "mode": mode, "model": model_name}
    _meta_path.write_text(json.dumps(meta, indent=2), encoding="utf-8")

    # IDENTITY.md — who the agent is
    identity = entry.get("identity", "")
    if identity:
        (cfg_dir / "IDENTITY.md").write_text(f"# Identity\n\n{identity}\n", encoding="utf-8")
    else:
        # Remove stale config from a previous mode
        (cfg_dir / "IDENTITY.md").unlink(missing_ok=True)

    # SOUL.md — how the agent behaves (optional, falls back to project/user default)
    soul = entry.get("soul")
    if soul:
        (cfg_dir / "SOUL.md").write_text(soul, encoding="utf-8")
    else:
        (cfg_dir / "SOUL.md").unlink(missing_ok=True)

    # Copy agent scripts — variant chosen by model (boltzbit vs generic).
    _server_dir = Path(__file__).resolve().parent
    import shutil as _sh

    _scripts_folder = f"scripts{_asset_suffix}"
    src_scripts = _server_dir / "bzcode_assets" / _scripts_folder
    if not src_scripts.is_dir():
        # Fallback 1: try without suffix (use boltzbit scripts if generic not present yet)
        if _asset_suffix:
            src_scripts = _server_dir / "bzcode_assets" / "scripts"
        # Fallback 2: check CWD
        if not src_scripts.is_dir():
            _cwd_candidate = Path.cwd() / "bzcode_assets" / _scripts_folder
            if _cwd_candidate.is_dir():
                src_scripts = _cwd_candidate
            else:
                _cwd_candidate2 = Path.cwd() / "bzcode_assets" / "scripts"
                if _cwd_candidate2.is_dir():
                    src_scripts = _cwd_candidate2
                else:
                    print(
                        "[session] WARNING: bzcode_assets/scripts not found — agent scripts unavailable",
                        file=sys.stderr,
                    )
    dst_scripts = cfg_dir / "scripts"
    dst_scripts.mkdir(exist_ok=True)
    for script in src_scripts.glob("*.py"):
        dest = dst_scripts / script.name
        if not dest.exists() or script.stat().st_mtime > dest.stat().st_mtime:
            _sh.copy2(script, dest)
    print(
        f"[session] scripts ({_scripts_folder}): {src_scripts} → {dst_scripts} ({len(list(dst_scripts.glob('*.py')))} files)",
        file=sys.stderr,
    )

    # Copy templates — variant chosen by model.
    _templates_folder = f"templates{_asset_suffix}"
    src_templates = _server_dir / "bzcode_assets" / _templates_folder
    if not src_templates.is_dir():
        if _asset_suffix:
            src_templates = _server_dir / "bzcode_assets" / "templates"
        if not src_templates.is_dir():
            _cwd_tmpl = Path.cwd() / "bzcode_assets" / _templates_folder
            if _cwd_tmpl.is_dir():
                src_templates = _cwd_tmpl
            else:
                src_templates = Path.cwd() / "bzcode_assets" / "templates"
    if src_templates.is_dir():
        dst_templates = cfg_dir / "templates"
        if dst_templates.exists():
            _shutil.rmtree(dst_templates)
        _sh.copytree(src_templates, dst_templates)
        print(f"[session] templates ({_templates_folder}): {src_templates} → {dst_templates}", file=sys.stderr)
    else:
        print("[session] WARNING: bzcode_assets/templates not found", file=sys.stderr)

    # {scripts_path} resolves to the session-local scripts directory.
    # Using the session config dir means no absolute paths leak into templates.
    _session_scripts = str(dst_scripts)

    def _resolve(text: str) -> str:
        return (
            text.replace("{server_data_path}", str(SERVER_DATA_DIR))
            .replace("{scripts_path}", _session_scripts)
            .replace("{session_dir}", str(cfg_dir))
            .replace("{working_dir}", working_dir)
            .replace("{widget_template_table}", _build_widget_template_table())
        )

    # AGENTS.md — pick boltzbit or generic variant; fall back to the other if absent.
    _agents_md_key = "agents_md" if is_boltzbit else "agents_md_generic"
    agents_md = entry.get(_agents_md_key) or entry.get("agents_md", "")
    if agents_md:
        (cfg_dir / "AGENTS.md").write_text(_resolve(agents_md), encoding="utf-8")
    else:
        (cfg_dir / "AGENTS.md").unlink(missing_ok=True)

    # settings.json — tools, model, permissions, etc.
    _settings_key = "settings" if is_boltzbit else "settings_generic"
    settings = entry.get(_settings_key) or entry.get("settings")
    if settings:
        (cfg_dir / "settings.json").write_text(json.dumps(settings, indent=2), encoding="utf-8")
    else:
        (cfg_dir / "settings.json").unlink(missing_ok=True)

    # skills/{name}/SKILL.md — pick boltzbit or generic variant.
    skills_dir = cfg_dir / "skills"
    if skills_dir.exists():
        _shutil.rmtree(skills_dir)
    _skills_key = "skills" if is_boltzbit else "skills_generic"
    skills = entry.get(_skills_key) or entry.get("skills", {})
    for skill_name, skill_content in skills.items():
        skill_path = skills_dir / skill_name / "SKILL.md"
        skill_path.parent.mkdir(parents=True, exist_ok=True)
        skill_path.write_text(_resolve(skill_content), encoding="utf-8")


# ══════════════════════════════════════════════════════════════════════════════
# § 1 · CONFIGURATION & CONSTANTS
# ══════════════════════════════════════════════════════════════════════════════

# ══════════════════════════════════════════════════════════════════════════════
# § 2 · IN-MEMORY STATE
# ══════════════════════════════════════════════════════════════════════════════

# Bzcode session files live under BZ_HOME/sessions/
SESSIONS_DIR = Path(os.environ.get("BZ_HOME") or "/usr/local/boltzbit").expanduser() / "sessions"

# ── Secret redaction ──────────────────────────────────────────────────────────
# Applied to every raw JSON string before it reaches the frontend WebSocket.
# Patterns cover the credential formats that appear in Bash tool output when
# the agent cats credential files or prints env vars.

_SECRET_KEY_RE = re.compile(
    r"(?i)\b(BZ_API_KEY|DYNAS_API_KEY|DPYES_API_KEY|ANKSY_API_KEY"
    r"|accessToken|refreshToken|access_token|refresh_token"
    r"|api_key|apikey|apiKey)",
    re.IGNORECASE,
)

# JSON double-quoted value:  "KEY": "VALUE{16+}"
_JSON_SECRET_RE = re.compile(
    r'(?i)(")(BZ_API_KEY|DYNAS_API_KEY|DPYES_API_KEY|ANKSY_API_KEY'
    r"|accessToken|refreshToken|access_token|refresh_token"
    r'|api_?key|apiKey)(")\s*:\s*"([^"]{16,})"',
)
# Env-var / plain  KEY=VALUE  or  KEY: VALUE  (unquoted values)
_PLAIN_SECRET_RE = re.compile(
    r"(?i)\b(BZ_API_KEY|DYNAS_API_KEY|DPYES_API_KEY|ANKSY_API_KEY"
    r"|access_?token|refresh_?token|api_?key|apiKey)\s*[=:]\s*([A-Za-z0-9._~+/\-]{16,})",
)
# BoltzBit API key literal (always safe to redact wherever it appears)
_BZ_KEY_RE = re.compile(r"bz_[A-Za-z0-9]{15,}")
# Bearer / Basic auth headers
_BEARER_RE = re.compile(r"(?i)(Bearer|Basic)\s+[A-Za-z0-9._~+/\-]{16,}=*")


def _redact(raw: str) -> str:
    raw = _BZ_KEY_RE.sub("[REDACTED]", raw)
    raw = _JSON_SECRET_RE.sub(lambda m: f'{m.group(1)}{m.group(2)}{m.group(3)}: "[REDACTED]"', raw)
    raw = _PLAIN_SECRET_RE.sub(lambda m: f"{m.group(1)}=[REDACTED]", raw)
    raw = _BEARER_RE.sub(lambda m: f"{m.group(1)} [REDACTED]", raw)
    return raw


# Per-file cursor positions: abs_path -> {selStart, selEnd}
# Stored in-memory (survives tab switches, cleared on server restart).
_cursor_store: dict = {}

# Tracks session IDs with an active WebSocket / bzcode process
_active_sessions = set()  # type: ignore[var-annotated]
# Tracks session IDs where bzcode is actively processing a request (status: running)
_running_sessions = set()  # type: ignore[var-annotated]
_TITLES_FILE = SESSIONS_DIR / "_titles.json"
_DEFAULTS_FILE = SESSIONS_DIR / "_defaults.json"  # cwd -> sessionId

# Accumulated token usage since server start (counts every result message)
_token_stats: dict = {"input": 0, "output": 0, "total": 0}


def _add_tokens(usage: dict) -> None:
    inp = int(usage.get("inputTokens", 0) or 0)
    out = int(usage.get("outputTokens", 0) or 0)
    _token_stats["input"] += inp
    _token_stats["output"] += out
    _token_stats["total"] += inp + out


def _load_titles() -> dict:
    try:
        return json.loads(_TITLES_FILE.read_text())
    except Exception:
        return {}


def _save_title(session_id: str, title: str) -> None:
    titles = _load_titles()
    titles[session_id] = title
    SESSIONS_DIR.mkdir(parents=True, exist_ok=True)
    _TITLES_FILE.write_text(json.dumps(titles, indent=2))


def _load_defaults() -> dict:
    try:
        return json.loads(_DEFAULTS_FILE.read_text())
    except Exception:
        return {}


def _save_default(cwd: str, session_id: str) -> None:
    defaults = _load_defaults()
    defaults[cwd] = session_id
    SESSIONS_DIR.mkdir(parents=True, exist_ok=True)
    _DEFAULTS_FILE.write_text(json.dumps(defaults, indent=2))


def _clear_default(cwd: str) -> None:
    defaults = _load_defaults()
    defaults.pop(cwd, None)
    _DEFAULTS_FILE.write_text(json.dumps(defaults, indent=2))


# Server-local data — lives alongside server.py so it travels with the project
SERVER_DATA_DIR = Path(__file__).resolve().parent / "server_data"

# ══════════════════════════════════════════════════════════════════════════════
# § 10 · CANVAS & CUSTOM WIDGETS
# ══════════════════════════════════════════════════════════════════════════════
# Canvas-edited widget code lives here, keyed by canvasId.
# Separate from server_data/widgets/ (toolbar templates) so edited instances
# don't pollute the template list.
CUSTOM_WIDGETS_DIR = SERVER_DATA_DIR / "custom_widgets"


def _custom_widgets_dir(session_id: str) -> Path:
    """Return the custom_widgets directory for a session (or global fallback)."""
    if session_id:
        return SESSIONS_DIR / session_id / "custom_widgets"
    return CUSTOM_WIDGETS_DIR


# ── Session reader ────────────────────────────────────────────────────────────


def _read_session_file(path: Path) -> Optional[dict]:
    """Parse a session JSONL and return metadata.  Handles two formats:

    Old format (bzcode < new version):
      Line 0: {"type": "session", "sessionId": "...", "workingDir": "..."}
      Lines 1+: message objects

    New format (bzcode >= new version):
      No header — messages start on line 0.
      workingDir/sessionId are read from {sessionId}/meta.json written by us.
    """
    try:
        with open(path, encoding="utf-8") as f:
            lines = [l.strip() for l in f if l.strip()]
        if not lines:
            return None

        # Detect format from first line
        first = json.loads(lines[0])
        session_id = path.stem
        working_dir = ""
        created = ""

        if first.get("type") == "session":
            # ── Old format: first line is the session header ──────────────────
            working_dir = first.get("workingDir", "")
            session_id = first.get("sessionId", path.stem)
            created = first.get("created", "")
            msg_lines = lines[1:]
        else:
            # ── New format: no header, look up our meta.json ──────────────────
            meta_file = SESSIONS_DIR / session_id / "meta.json"
            if meta_file.exists():
                try:
                    meta = json.loads(meta_file.read_text())
                    working_dir = meta.get("workingDir", "")
                    session_id = meta.get("sessionId", session_id)
                except Exception:
                    pass
            if not working_dir:
                return None  # can't place this session in any project
            msg_lines = lines

        # Walk messages to extract title and last preview
        title = ""
        last_preview = ""
        msg_count = 0
        _SKIP_EXACT = {"Hi, hand shake, say yes", "[Request interrupted by user]"}
        import re as _re
        _SYSREM_RE = _re.compile(r"<system-reminder>.*?</system-reminder>", _re.DOTALL)
        for line in msg_lines:
            try:
                msg = json.loads(line)
                if msg.get("role") == "user":
                    msg_count += 1
                    content = msg.get("content", "")
                    text = ""
                    if isinstance(content, str):
                        text = content
                    elif isinstance(content, list):
                        for block in content:
                            if isinstance(block, dict) and block.get("type") == "text":
                                text = block.get("text", "")
                                break
                    # Strip embedded system-reminder blocks then check for noise
                    clean = _SYSREM_RE.sub("", text).strip()
                    _t = clean
                    if not title and _t and _t not in _SKIP_EXACT:
                        title = _t[:60]
                    if _t and _t not in _SKIP_EXACT:
                        last_preview = _t[:150]
            except json.JSONDecodeError:
                pass

        # Read agent mode from our meta.json
        agent_mode = "general"
        meta_file = SESSIONS_DIR / session_id / "meta.json"
        if meta_file.exists():
            try:
                agent_mode = json.loads(meta_file.read_text()).get("mode", "general")
            except Exception:
                pass

        stat = path.stat()
        custom_titles = _load_titles()
        return {
            "sessionId": session_id,
            "workingDir": working_dir,
            "dirName": Path(working_dir).name if working_dir else "Unknown",
            "messageCount": msg_count,
            "title": custom_titles.get(session_id) or title or "(empty)",
            "lastMessage": last_preview,
            "lastModified": stat.st_mtime,
            "created": created,
            "mode": agent_mode,
        }
    except Exception:
        return None


# ══════════════════════════════════════════════════════════════════════════════
# § 5 · AUTHENTICATION & CREDENTIALS
# ══════════════════════════════════════════════════════════════════════════════

CORS_HEADERS = {
    "Access-Control-Allow-Origin": "*",
    "Access-Control-Allow-Methods": "GET, POST, PUT, DELETE, OPTIONS",
    "Access-Control-Allow-Headers": "*",
}


def _write_bzcode_credentials(
    access_token: str,
    refresh_token: str = "",
    expires_at=None,
    auth_url: str = "https://boltzhub.com",
) -> None:
    """Write access token to $BZ_HOME/credentials.json in the format bzcode expects."""
    creds_dir = Path(os.environ.get("BZ_HOME") or "/usr/local/boltzbit").expanduser()
    creds_file = creds_dir / "credentials.json"
    creds_dir.mkdir(parents=True, exist_ok=True)
    existing: dict = {}
    if creds_file.exists():
        try:
            with open(creds_file) as f:
                existing = json.load(f)
        except Exception:
            pass
    prev_entry: dict = existing.get(auth_url, {})
    entry: dict = {"accessToken": access_token}
    # Preserve the existing refreshToken when the caller doesn't supply one.
    resolved_refresh = refresh_token or prev_entry.get("refreshToken", "")
    if resolved_refresh:
        entry["refreshToken"] = resolved_refresh
    # bzcode expects expiresAt in milliseconds. If the caller didn't supply it,
    # parse the exp claim from the JWT itself so we always write a valid expiry.
    if expires_at is None:
        try:
            import base64 as _b64

            seg = access_token.split(".")[1]
            seg += "=" * (4 - len(seg) % 4)
            payload = json.loads(_b64.b64decode(seg.replace("-", "+").replace("_", "/")))
            exp = payload.get("exp")
            if exp:
                expires_at = exp * 1000  # seconds → milliseconds
        except Exception:
            pass
    if expires_at is None:
        # Preserve whatever was there before rather than dropping it.
        expires_at = prev_entry.get("expiresAt")
    if expires_at is not None:
        ms_val = int(expires_at)
        if ms_val < 10_000_000_000:  # looks like seconds → convert
            ms_val *= 1000
        entry["expiresAt"] = ms_val
    existing[auth_url] = entry
    with open(creds_file, "w") as f:
        json.dump(existing, f, indent=2)
    print(f"[auth] credentials written to {creds_file} for {auth_url}", file=sys.stderr)


def _read_api_keys() -> dict:
    """Return only BZ_API_KEY from BZ_HOME/api_keys.json.
    No other keys are passed to bzcode spawns — this is intentional."""
    keys_file = Path(os.environ.get("BZ_HOME") or "/usr/local/boltzbit").expanduser() / "api_keys.json"
    if not keys_file.exists():
        return {}
    try:
        with open(keys_file) as f:
            data = json.load(f)
        val = data.get("BZ_API_KEY", "")
        return {"BZ_API_KEY": val} if val and isinstance(val, str) else {}
    except Exception:
        return {}


# ══════════════════════════════════════════════════════════════════════════════
# § 11 · WIDGET SYSTEM
# ══════════════════════════════════════════════════════════════════════════════
#
# Layout inside server_data/widgets/:
#   index.json          — metadata for all widgets (no code)
#   {id}.js             — code for each widget (one file per widget)
#
# This keeps the search index lean and lets code files be edited directly.

WIDGETS_DIR = SERVER_DATA_DIR / "widgets"
WIDGETS_INDEX = WIDGETS_DIR / "index.json"


def _now() -> str:
    return datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")


def _load_index() -> dict:
    """Load index.json; returns empty structure if missing or corrupt."""
    if not WIDGETS_INDEX.exists():
        return {"version": 1, "widgets": []}
    try:
        with open(WIDGETS_INDEX, encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {"version": 1, "widgets": []}


def _save_index(data: dict) -> None:
    WIDGETS_DIR.mkdir(parents=True, exist_ok=True)
    with open(WIDGETS_INDEX, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)


def _code_path(widget_id: str) -> Path:
    """Return the .js file path for a widget id."""
    # Sanitise id so it's safe as a filename
    safe = "".join(c if c.isalnum() or c in "-_" else "_" for c in widget_id)
    return WIDGETS_DIR / f"{safe}.js"


def _load_code(widget_id: str) -> str:
    path = _code_path(widget_id)
    if not path.exists():
        return ""
    with open(path, encoding="utf-8") as f:
        return f.read()


def _save_code(widget_id: str, code: str) -> None:
    WIDGETS_DIR.mkdir(parents=True, exist_ok=True)
    with open(_code_path(widget_id), "w", encoding="utf-8") as f:
        f.write(code)


def _load_creds() -> dict:
    creds_file = SERVER_DATA_DIR / "credentials.json"
    if not creds_file.exists():
        return {}
    try:
        with open(creds_file, encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}


_PLACEHOLDER_RE = re.compile(r"\{\{(\w+)\}\}")


def _resolve(text: str, creds: dict) -> str:
    """Replace {{KEY}} placeholders with stored credential values."""
    return _PLACEHOLDER_RE.sub(lambda m: creds.get(m.group(1), m.group(0)), text)


# ══════════════════════════════════════════════════════════════════════════════
# § 6 · FILE SYSTEM
# ══════════════════════════════════════════════════════════════════════════════


def _canvas_file(session_id: str, cwd: str) -> Path:
    """Return the .bzcanvas.json path — session dir preferred, cwd as fallback."""
    if session_id:
        return SESSIONS_DIR / session_id / ".bzcanvas.json"
    return Path(cwd) / ".bzcanvas.json"


# ══════════════════════════════════════════════════════════════════════════════
# § 14 · WHATSAPP
# ══════════════════════════════════════════════════════════════════════════════
#
# Setup:
#   1. Create a Twilio account and activate the WhatsApp Sandbox.
#   2. Add these credentials via POST /credentials:
#        TWILIO_ACCOUNT_SID   your-account-sid
#        TWILIO_AUTH_TOKEN    your-auth-token
#        TWILIO_FROM          whatsapp:+14155238886   (sandbox number)
#   3. Set the sandbox webhook URL to:
#        http://YOUR_VM_IP:5081/whatsapp/incoming
#   4. Send a message from your phone → bzcode replies in the {cwd}/whatsapp/ dir.

# One persistent bzcode process per WhatsApp phone number.
_whatsapp_sessions: dict = {}
_whatsapp_lock: Optional[asyncio.Lock] = None


def _get_whatsapp_lock() -> asyncio.Lock:
    global _whatsapp_lock
    if _whatsapp_lock is None:
        _whatsapp_lock = asyncio.Lock()
    return _whatsapp_lock


class _WASess:
    """Persistent bzcode process for one WhatsApp contact."""

    def __init__(self, phone: str, bzcode_path: str, cwd: str):
        self.phone = phone
        self.bzcode_path = bzcode_path
        self.cwd = cwd
        self.proc: Optional[asyncio.subprocess.Process] = None
        self._buf: list = []
        self._done = asyncio.Event()
        self._msg_lock = asyncio.Lock()
        import hashlib as _hl

        self.session_id = f"wa-{_hl.md5(phone.encode()).hexdigest()[:12]}"

    async def _start(self) -> None:
        """Spawn (or restart) the bzcode process."""
        self.proc = await asyncio.create_subprocess_exec(
            self.bzcode_path,
            "--stdio",
            "--resume",
            self.session_id,
            stdin=asyncio.subprocess.PIPE,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            cwd=self.cwd,
            env={**os.environ, **_read_api_keys(), "BZ_PYTHON": sys.executable},
            limit=16 * 1024 * 1024,
        )
        asyncio.create_task(self._read_loop())
        # Drain the session message and switch to yolo so tools auto-approve
        await asyncio.sleep(0.3)
        self.proc.stdin.write(b'{"type":"setMode","mode":"yolo"}\n')
        await self.proc.stdin.drain()
        await asyncio.sleep(0.2)

    async def _read_loop(self) -> None:
        assert self.proc and self.proc.stdout
        while True:
            line = await self.proc.stdout.readline()
            if not line:
                break
            raw = line.decode().strip()
            if not raw:
                continue
            try:
                msg = json.loads(raw)
            except json.JSONDecodeError:
                continue
            t = msg.get("type")
            if t == "assistant":
                for block in msg.get("content", []):
                    if block.get("type") == "text" and block.get("text"):
                        self._buf.append(block["text"])
            elif t == "result" and msg.get("output"):
                self._buf.append(msg["output"])
            elif t == "status" and msg.get("status") == "idle":
                self._done.set()

    async def chat(self, text: str, timeout: int = 90) -> str:
        """Send a message and wait for the full response."""
        async with self._msg_lock:
            if self.proc is None or self.proc.returncode is not None:
                await self._start()
            self._buf.clear()
            self._done.clear()
            payload = json.dumps({"type": "user", "content": text}) + "\n"
            self.proc.stdin.write(payload.encode())
            await self.proc.stdin.drain()
            try:
                await asyncio.wait_for(self._done.wait(), timeout=timeout)
            except asyncio.TimeoutError:
                return "⏱ Response timed out. Try again."
            return "\n\n".join(self._buf).strip() or "✓ Done."


async def _send_whatsapp(to: str, body: str, creds: dict) -> None:
    """Send a WhatsApp message back via Twilio REST API."""
    sid = creds.get("TWILIO_ACCOUNT_SID", "")
    token = creds.get("TWILIO_AUTH_TOKEN", "")
    from_ = creds.get("TWILIO_FROM", "")
    if not (sid and token and from_):
        print("[whatsapp] missing Twilio credentials", file=sys.stderr)
        return
    url = f"https://api.twilio.com/2010-04-01/Accounts/{sid}/Messages.json"
    data = {"From": from_, "To": to, "Body": body[:1500]}  # WhatsApp 1600-char limit
    connector = aiohttp.TCPConnector(ssl=False)
    async with aiohttp.ClientSession(connector=connector) as session:
        async with session.post(url, data=data, auth=aiohttp.BasicAuth(sid, token)) as resp:
            if resp.status not in (200, 201):
                text = await resp.text()
                print(f"[whatsapp] Twilio error {resp.status}: {text}", file=sys.stderr)


# ══════════════════════════════════════════════════════════════════════════════
# § 13 · BOLTZHUB INTEGRATION
# ══════════════════════════════════════════════════════════════════════════════

BOLTZHUB_API = "https://boltzhub.com/bz-appstore-api"
BOLTZHUB_AUTH = "https://boltzhub.com"
# Matches what the VS Code plugin excludes — build output (dist/) is intentionally included
_PUSH_EXCLUDE = {".git", "node_modules", ".bzhub", "__pycache__", ".venv", "venv"}


def _boltzhub_token() -> Optional[str]:
    """Return an auth token. BZ_API_KEY (non-expiring) takes priority over OAuth JWT from credentials.json."""
    # API key first — never expires
    api_keys = _read_api_keys()
    api_key = api_keys.get("BZ_API_KEY") or os.environ.get("BZ_API_KEY")
    if api_key:
        return api_key
    # Fallback: OAuth JWT from credentials.json (can expire)
    try:
        import json as _json

        bz_home = Path(os.environ.get("BZ_HOME") or "/usr/local/boltzbit").expanduser()
        creds = _json.loads((bz_home / "credentials.json").read_text())
        tok = creds.get(BOLTZHUB_AUTH, {}).get("accessToken")
        if tok:
            return tok
    except Exception:
        pass
    return None


def _credentials_valid() -> Tuple[bool, str]:
    """Return (ok, reason). Accepts BZ_API_KEY in api_keys.json or a valid credentials.json."""
    import time as _time

    bz_home = Path(os.environ.get("BZ_HOME") or "/usr/local/boltzbit").expanduser()
    # BZ_API_KEY is sufficient — no OAuth credentials needed
    api_keys = _read_api_keys()
    if api_keys.get("BZ_API_KEY"):
        return True, ""
    creds_file = bz_home / "credentials.json"
    if not creds_file.exists():
        return False, f"credentials.json not found in {bz_home}"
    try:
        creds = json.loads(creds_file.read_text())
    except Exception as e:
        return False, f"credentials.json is unreadable: {e}"
    entry = creds.get(BOLTZHUB_AUTH, {})
    token = entry.get("accessToken", "")
    if not token:
        return False, "no accessToken in credentials.json"
    if not entry.get("refreshToken", ""):
        return False, "no refreshToken in credentials.json — please log in again"
    expires_at = entry.get("expiresAt")
    if expires_at:
        ms = int(expires_at)
        # expiresAt may be in seconds or milliseconds
        if ms < 10_000_000_000:
            ms *= 1000
        if ms < _time.time() * 1000:
            return False, "accessToken is expired — please log in again"
    return True, ""


def _read_app_config(cwd: str) -> Optional[dict]:
    try:
        return json.loads((Path(cwd) / ".bzhub" / "app_config.json").read_text())
    except Exception:
        return None


def _write_app_config(cwd: str, config: dict) -> None:
    bzhub = Path(cwd) / ".bzhub"
    bzhub.mkdir(parents=True, exist_ok=True)
    (bzhub / "app_config.json").write_text(json.dumps(config, indent=2))


def _sync_env_oauth_client_id(cwd: str, app_id: str) -> None:
    """Ensure VITE_OAUTH_CLIENT_ID in .env matches the app ID from app_config.json."""
    env_path = Path(cwd) / ".env"
    if not env_path.exists():
        return
    lines = env_path.read_text().splitlines(keepends=True)
    new_lines = []
    found = False
    for line in lines:
        if "=" in line and not line.lstrip().startswith("#"):
            key, _, _ = line.partition("=")
            if key.strip() == "VITE_OAUTH_CLIENT_ID":
                new_lines.append(f"VITE_OAUTH_CLIENT_ID={app_id}\n")
                found = True
                continue
        new_lines.append(line)
    if not found:
        new_lines.append(f"VITE_OAUTH_CLIENT_ID={app_id}\n")
    env_path.write_text("".join(new_lines))


def _bz_headers(token: str) -> dict:
    if token.startswith("bz_"):
        return {"X-API-Key": token, "Content-Type": "application/json"}
    return {"Authorization": f"Bearer {token}", "Content-Type": "application/json"}


def _bz_auth(token: str) -> dict:
    """Auth-only headers (no Content-Type) for multipart/form-data requests."""
    if token.startswith("bz_"):
        return {"X-API-Key": token}
    return {"Authorization": f"Bearer {token}"}


# ══════════════════════════════════════════════════════════════════════════════
# § 15 · BATCH EXECUTION
# ══════════════════════════════════════════════════════════════════════════════


_batch_store: dict = {}  # batchId -> {"items": [...], "created": float}


class _BatchItem:
    """Runs a single bzcode process in YOLO mode and collects the response."""

    def __init__(self, cwd: str, bzcode_path: str, resume_session_id: str = "") -> None:
        self.cwd = cwd
        self.dir_name = Path(cwd).name
        self.bzcode_path = bzcode_path
        self.resume_session_id = resume_session_id  # if set, resume this session
        self.status = "pending"  # pending | running | done | error
        self.output = ""
        self.error_msg = ""
        self.session_id = resume_session_id  # filled/confirmed from bzcode's session message
        self._buf: list = []
        self._done = asyncio.Event()
        self._proc = None
        self._msg_sent = False  # True once the user message has been written to stdin

    def to_dict(self) -> dict:
        return {
            "cwd": self.cwd,
            "dirName": self.dir_name,
            "status": self.status,
            "output": self.output,
            "error": self.error_msg,
            "sessionId": self.session_id,
        }

    async def run(self, message: str) -> None:
        self.status = "running"
        _running_sessions.add(self.resume_session_id or self.session_id)
        try:
            if not self.resume_session_id:
                import secrets as _sec

                self.resume_session_id = f"bz-{_sec.token_hex(6)}"
                self.session_id = self.resume_session_id
            cmd = [self.bzcode_path, "--stdio", "--resume", self.resume_session_id]
            self._proc = await asyncio.create_subprocess_exec(
                *cmd,
                stdin=asyncio.subprocess.PIPE,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.DEVNULL,
                cwd=self.cwd,
                env={**os.environ, **_read_api_keys(), "BZ_PYTHON": sys.executable},
                limit=16 * 1024 * 1024,
            )
            asyncio.create_task(self._read_loop())
            # Allow bzcode to send its session/ready message
            await asyncio.sleep(0.8)
            # YOLO mode — auto-approve all tool permissions
            self._proc.stdin.write(b'{"type":"setMode","mode":"yolo"}\n')
            await self._proc.stdin.drain()
            await asyncio.sleep(0.1)
            # Send the user message
            payload = json.dumps({"type": "user", "content": message}) + "\n"
            self._proc.stdin.write(payload.encode())
            await self._proc.stdin.drain()
            self._msg_sent = True  # now it's safe to honour status: idle
            # Wait up to 3 minutes for completion
            await asyncio.wait_for(self._done.wait(), timeout=180)
            self.output = "\n\n".join(self._buf).strip()
            self.status = "done"
            # Give bzcode a moment to finish writing the session file to disk
            await asyncio.sleep(0.5)
        except asyncio.TimeoutError:
            self.status = "error"
            self.error_msg = "Timed out after 3 minutes"
        except Exception as exc:
            self.status = "error"
            self.error_msg = str(exc)
        finally:
            _running_sessions.discard(self.resume_session_id or self.session_id)
            if self._proc:
                try:
                    # Close stdin gracefully — signals bzcode to save session and exit
                    if self._proc.stdin:
                        self._proc.stdin.close()
                    await asyncio.wait_for(self._proc.wait(), timeout=8)
                except asyncio.TimeoutError:
                    self._proc.kill()
                except Exception:
                    try:
                        self._proc.kill()
                    except Exception:
                        pass

    async def _read_loop(self) -> None:
        assert self._proc and self._proc.stdout
        while True:
            line = await self._proc.stdout.readline()
            if not line:
                break
            raw = line.decode().strip()
            if not raw:
                continue
            try:
                msg = json.loads(raw)
                t = msg.get("type")
                if t == "session":
                    self.session_id = msg.get("sessionId", "")
                elif t == "assistant":
                    for block in msg.get("content", []):
                        if block.get("type") == "text" and block.get("text"):
                            self._buf.append(block["text"])
                elif t == "result":
                    if msg.get("output"):
                        self._buf.append(msg["output"])
                    if msg.get("usage"):
                        _add_tokens(msg["usage"])
                elif t == "status" and msg.get("status") == "idle":
                    # Only complete after the user message has been sent —
                    # ignores startup auto-run idle that arrives before our message
                    if self._msg_sent:
                        self._done.set()
            except Exception:
                pass


# ══════════════════════════════════════════════════════════════════════════════
# § 19 · AGENT POOL
# ══════════════════════════════════════════════════════════════════════════════
# Decouples bzcode process lifetime from WebSocket lifetime.
# Process-facing tasks (stdout/stderr readers) are long-lived.
# WebSocket-facing tasks (send_to_client, relay_client_messages) attach/detach.

import time as _time


class AgentPoolEntry:
    """A pooled bzcode process that outlives individual WebSocket connections."""

    def __init__(self, session_id: str, cwd: str, mode: str):
        self.session_id = session_id
        self.cwd = cwd
        self.mode = mode  # agent mode from agent_modes.json (general/widget/worker/coder)
        self.proc: Optional[asyncio.subprocess.Process] = None

        # Per-agent state tracked from bzcode stdout messages
        self.agent_status = "starting"  # starting | idle | running | waiting_permission | waiting_input
        self.session_mode = "default"  # runtime mode from bzcode: default | yolo | plan
        self.model_info: dict = {}
        self._pending_request_id: Optional[str] = None

        # Agent config captured from bzcode's initial session message
        self.available_modes: list = []
        self.available_commands: list = []

        # Process-facing output queue — fed by read_bzcode_stdout and drain_bzcode_stderr
        self._out_queue: asyncio.Queue = asyncio.Queue()
        self._ready_event = asyncio.Event()

        # Per-turn replay buffer: accumulates messages during the current turn
        # so reconnecting SSE clients can catch up. Cleared on turn boundaries.
        self._turn_buffer: list = []

        # Subscriber queues — each SSE or WS connection gets one
        self._subscribers: set = set()

        # WebSocket attachment state (legacy, kept for backward compat)
        self._ws = None  # currently attached WS (aiohttp-like interface)
        self._ws_fwd_queue: Optional[asyncio.Queue] = None
        self._ws_send_task: Optional[asyncio.Task] = None
        self._ws_relay_task: Optional[asyncio.Task] = None

        # Process-facing tasks (long-lived)
        self._stdout_task: Optional[asyncio.Task] = None
        self._stderr_task: Optional[asyncio.Task] = None
        self._dispatcher_task: Optional[asyncio.Task] = None

        # Lifecycle
        self._created_at = _time.monotonic()
        self._last_ws_detach_at: Optional[float] = None
        self._attach_count = 0  # how many times a WS has been attached
        self._shutting_down = False
        self._attach_lock = asyncio.Lock()

    async def start(self, bzcode_path: str, cmd: list, env: dict) -> None:
        """Spawn the bzcode process and start process-facing tasks."""
        self.proc = await asyncio.create_subprocess_exec(
            *cmd,
            stdin=asyncio.subprocess.PIPE,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            cwd=self.cwd,
            env=env,
        )
        self._stdout_task = asyncio.create_task(
            read_bzcode_stdout(
                self.proc, self._out_queue, self._ready_event, session_id=self.session_id, mode=self.mode
            )
        )
        self._stderr_task = asyncio.create_task(drain_bzcode_stderr(self.proc, self._out_queue))
        self._dispatcher_task = asyncio.create_task(self._dispatch_stdout())

        # Send setMode once bzcode emits its session-ready message. Done as a
        # background task so start() (and POST /api/pool/connect) returns as soon
        # as the process is spawned — not after bzcode finishes restoring history.
        entry = _mode_entry(self.mode)
        session_mode = (entry.get("settings") or {}).get("mode", "")
        if session_mode:
            _sid = self.session_id
            _proc = self.proc
            _evt = self._ready_event
            _sm = session_mode

            async def _send_mode_once_ready() -> None:
                try:
                    await asyncio.wait_for(_evt.wait(), timeout=30)
                except asyncio.TimeoutError:
                    pass
                if _proc.returncode is None:
                    _proc.stdin.write(json.dumps({"type": "setMode", "mode": _sm}).encode() + b"\n")
                    await _proc.stdin.drain()
                    print(f"[pool] sent setMode={_sm} to {_sid}", file=sys.stderr)

            asyncio.create_task(_send_mode_once_ready())

    async def _dispatch_stdout(self) -> None:
        """Read from _out_queue, track agent state, auto-approve in yolo, fan out."""
        while True:
            raw = await self._out_queue.get()
            if raw is None:
                self.agent_status = "dead"
                print(f"[{self.session_id}] stdout closed", file=sys.stderr)
                for q in list(self._subscribers):
                    try:
                        q.put_nowait(None)
                    except asyncio.QueueFull:
                        pass
                break

            forward = True  # whether to fan out to subscribers
            _was_running = self.agent_status == "running"  # snapshot before any update

            if raw and raw[0] == "{":
                try:
                    _msg = json.loads(raw)
                    _t = _msg.get("type", "")

                    # ── State tracking ───────────────────────────────────
                    if _t == "session":
                        if isinstance(_msg.get("modes"), list):
                            self.available_modes = _msg["modes"]
                        if isinstance(_msg.get("commands"), list):
                            self.available_commands = _msg["commands"]

                    elif _t == "status":
                        _s = _msg.get("status", "")
                        if _s == "running":
                            # Don't reset the replay buffer here. A turn's buffer is
                            # seeded with the user prompt when the message is sent
                            # (seed_user_turn), so it must span the whole turn —
                            # including the prompt and any intra-turn "running"
                            # re-emissions (tool calls, permission/input replies).
                            # Clearing on "running" would drop the start of the turn
                            # (prompt + early output) on a mid-turn reconnect.
                            self.agent_status = "running"
                        elif _s == "idle":
                            self.agent_status = "idle"
                        if _msg.get("mode"):
                            self.session_mode = _msg["mode"]
                        if _msg.get("model"):
                            self.model_info = _msg["model"] if isinstance(_msg["model"], dict) else {}

                    elif _t == "result":
                        self.agent_status = "idle"

                    elif _t == "prompt":
                        _sub = _msg.get("subtype", "")
                        _rid = _msg.get("requestId", "")
                        if _sub == "permission":
                            self._pending_request_id = _rid
                            if self.session_mode == "yolo":
                                # Auto-approve: write directly to stdin, skip frontend
                                self.agent_status = "running"
                                _resp = (
                                    json.dumps(
                                        {
                                            "type": "user",
                                            "subtype": "permission",
                                            "requestId": _rid,
                                            "behavior": "always",
                                        }
                                    )
                                    + "\n"
                                )
                                self.proc.stdin.write(_resp.encode())
                                asyncio.ensure_future(self.proc.stdin.drain())
                                print(f"[{self.session_id}] auto-approved {_msg.get('tool', '?')}", file=sys.stderr)
                                forward = False
                            else:
                                self.agent_status = "waiting_permission"
                        elif _sub == "input":
                            self.agent_status = "waiting_input"
                            self._pending_request_id = _rid

                    # ── Logging ──────────────────────────────────────────
                    if _t == "delta":
                        pass
                    elif _t == "assistant":
                        _preview = ""
                        for _b in _msg.get("content") or []:
                            if isinstance(_b, dict) and _b.get("type") == "text":
                                _preview = str(_b.get("text", ""))[:120]
                                break
                        print(f"[{self.session_id}] assistant: {_preview}", file=sys.stderr)
                    elif forward:
                        _preview = raw[:200] if len(raw) > 200 else raw
                        print(f"[{self.session_id}] {_t}: {_preview}", file=sys.stderr)

                except Exception:
                    print(f"[{self.session_id}] raw: {raw[:200]}", file=sys.stderr)

            if forward:
                self._turn_buffer.append(raw)
                for q in list(self._subscribers):
                    try:
                        q.put_nowait(raw)
                    except asyncio.QueueFull:
                        pass
                # Clear buffer only after a complete turn (running → idle/result).
                # Spurious idle events (e.g. bzcode acking setMode between turns)
                # must not wipe a freshly-seeded user message from the buffer.
                if self.agent_status == "idle" and _was_running:
                    self._turn_buffer.clear()

    def subscribe(self, replay: bool = False) -> asyncio.Queue:
        """Add a subscriber queue. If replay=True, pre-fill with current turn's buffered messages."""
        q: asyncio.Queue = asyncio.Queue(maxsize=1000)
        if replay:
            import sys as _sys
            _types = []
            for _m in self._turn_buffer:
                try: _types.append(json.loads(_m).get("type", "?"))
                except Exception: _types.append("?")
            print(f"[{self.session_id}] subscribe replay buffer ({len(self._turn_buffer)} msgs): {_types}", file=_sys.stderr)
        if replay:
            still_waiting = self.agent_status in ("waiting_input", "waiting_permission")
            for msg in self._turn_buffer:
                # Skip prompt events that the user has already answered so they
                # don't re-appear on reconnect.
                if not still_waiting and msg and msg[0] == "{":
                    try:
                        if json.loads(msg).get("type") == "prompt":
                            continue
                    except Exception:
                        pass
                try:
                    q.put_nowait(msg)
                except asyncio.QueueFull:
                    break
        self._subscribers.add(q)
        self._last_ws_detach_at = None
        return q

    def unsubscribe(self, q: asyncio.Queue) -> None:
        """Remove a subscriber queue."""
        self._subscribers.discard(q)
        if not self._subscribers and self._ws is None:
            self._last_ws_detach_at = _time.monotonic()

    def seed_user_turn(self, raw: str) -> None:
        """Publish the user's prompt to the SSE stream at the start of a turn.

        The prompt is written to bzcode's stdin and isn't echoed on stdout, and it
        isn't in the .jsonl transcript until the turn completes. So we inject it into
        the stream ourselves: it becomes the first entry of the turn's replay buffer
        (so a later reconnect/refresh replays it) AND is fanned out to any already
        connected subscribers (so the live client renders it). The buffer is cleared
        again when the turn finishes (agent_status == "idle", see _dispatch_stdout).
        """
        self._turn_buffer = [raw]
        for q in list(self._subscribers):
            try:
                q.put_nowait(raw)
            except asyncio.QueueFull:
                pass

    def _read_session_messages(self) -> list:
        """Read .jsonl transcript and return messages list (timestamps stripped)."""
        jsonl_path = SESSIONS_DIR / f"{self.session_id}.jsonl"
        if not jsonl_path.exists():
            return []
        try:
            messages = []
            with open(jsonl_path, encoding="utf-8") as f:
                for line in f:
                    line = line.strip()
                    if not line:
                        continue
                    try:
                        entry = json.loads(line)
                        entry.pop("timestamp", None)
                        if entry.get("type") == "session":
                            continue
                        messages.append(entry)
                    except json.JSONDecodeError:
                        continue
            return messages
        except Exception:
            return []

    async def attach_ws(self, ws) -> tuple:
        """Attach a WebSocket. Returns (send_task, relay_task) for the caller to await."""
        async with self._attach_lock:
            if self._shutting_down:
                raise RuntimeError("Agent is shutting down")

            # If another WS is attached, detach it first
            if self._ws is not None:
                await self._detach_ws_internal()

            self._attach_count += 1
            is_reattach = self._attach_count > 1
            self._ws = ws
            self._ws_fwd_queue = self.subscribe()  # register as subscriber
            self._last_ws_detach_at = None

            # On reattach, send conversation history since bzcode won't re-emit it
            if is_reattach:
                messages = self._read_session_messages()
                session_msg = json.dumps(
                    {
                        "type": "session",
                        "sessionId": self.session_id,
                        "messages": messages,
                    }
                )
                try:
                    await ws.send_str(session_msg)
                except Exception:
                    pass

            # Start WS-facing tasks
            self._ws_send_task = asyncio.create_task(send_to_client(self._ws_fwd_queue, ws))
            self._ws_relay_task = asyncio.create_task(relay_client_messages(self.proc, ws, self._ready_event))

            return self._ws_send_task, self._ws_relay_task

    async def detach_ws(self) -> None:
        """Detach the current WebSocket. Process keeps running."""
        async with self._attach_lock:
            await self._detach_ws_internal()

    async def _detach_ws_internal(self) -> None:
        """Internal detach — caller must hold _attach_lock."""
        if self._ws_relay_task and not self._ws_relay_task.done():
            self._ws_relay_task.cancel()
            try:
                await self._ws_relay_task
            except (asyncio.CancelledError, Exception):
                pass
        if self._ws_send_task and not self._ws_send_task.done():
            if self._ws_fwd_queue:
                try:
                    self._ws_fwd_queue.put_nowait(None)
                except asyncio.QueueFull:
                    pass
            self._ws_send_task.cancel()
            try:
                await self._ws_send_task
            except (asyncio.CancelledError, Exception):
                pass
        if self._ws_fwd_queue:
            self.unsubscribe(self._ws_fwd_queue)
        self._ws = None
        self._ws_send_task = None
        self._ws_relay_task = None
        self._ws_fwd_queue = None

    async def shutdown(self, reason: str = "idle") -> None:
        """Gracefully stop the bzcode process."""
        if self._shutting_down:
            return
        self._shutting_down = True
        print(
            f"[pool] shutting down {self.session_id} reason={reason} pid={self.proc.pid if self.proc else '?'}",
            file=sys.stderr,
        )
        await self.detach_ws()
        # Close stdin — signals bzcode to save session and exit
        if self.proc and self.proc.stdin:
            try:
                self.proc.stdin.close()
            except Exception:
                pass
        if self.proc:
            try:
                await asyncio.wait_for(self.proc.wait(), timeout=8)
            except asyncio.TimeoutError:
                try:
                    self.proc.kill()
                except ProcessLookupError:
                    pass
            except ProcessLookupError:
                pass
        # Cancel process-facing tasks
        for task in (self._stdout_task, self._stderr_task, self._dispatcher_task):
            if task and not task.done():
                task.cancel()
                try:
                    await task
                except (asyncio.CancelledError, Exception):
                    pass
        _active_sessions.discard(self.session_id)
        _running_sessions.discard(self.session_id)

    @property
    def is_dead(self) -> bool:
        return self.proc is not None and self.proc.returncode is not None

    @property
    def has_ws(self) -> bool:
        return self._ws is not None

    @property
    def has_clients(self) -> bool:
        return self._ws is not None or len(self._subscribers) > 0


class AgentPool:
    """Manages a pool of bzcode processes keyed by sessionId."""

    def __init__(self, idle_timeout: float = 300.0):
        self._entries: dict[str, AgentPoolEntry] = {}
        self._lock: Optional[asyncio.Lock] = None  # created in start() inside the running loop
        self._idle_timeout = idle_timeout
        self._idle_check_task: Optional[asyncio.Task] = None

    async def start(self) -> None:
        """Start the idle-timeout sweeper. Call from FastAPI lifespan startup."""
        # Create the lock here so it's bound to the running event loop (Python 3.8 compat).
        self._lock = asyncio.Lock()
        self._idle_check_task = asyncio.create_task(self._idle_sweeper())
        print(f"[pool] started  idle_timeout={self._idle_timeout}s", file=sys.stderr)

    async def stop(self) -> None:
        """Shutdown all agents. Call from FastAPI lifespan shutdown."""
        if self._idle_check_task:
            self._idle_check_task.cancel()
            try:
                await self._idle_check_task
            except (asyncio.CancelledError, Exception):
                pass
        async with self._lock:
            for entry in list(self._entries.values()):
                await entry.shutdown(reason="server_shutdown")
            self._entries.clear()
        print("[pool] stopped — all agents shut down", file=sys.stderr)

    async def get_or_create(
        self,
        session_id: str,
        cwd: str,
        mode: str,
        bzcode_path: str,
        cmd: list,
        env: dict,
    ) -> AgentPoolEntry:
        """Return existing entry for session_id, or create a new one."""
        if session_id.startswith("bz-probe-"):
            raise ValueError(f"probe session {session_id!r} must not enter the pool")
        async with self._lock:
            entry = self._entries.get(session_id)

            if entry is not None:
                if entry.is_dead:
                    print(f"[pool] removing dead entry {session_id}", file=sys.stderr)
                    del self._entries[session_id]
                    entry = None
                elif entry._shutting_down:
                    await entry.shutdown()
                    self._entries.pop(session_id, None)
                    entry = None

            if entry is None:
                entry = AgentPoolEntry(session_id=session_id, cwd=cwd, mode=mode)
                await entry.start(bzcode_path, cmd, env)
                self._entries[session_id] = entry
                print(f"[pool] spawned {session_id} pid={entry.proc.pid}", file=sys.stderr)
            else:
                print(f"[pool] reusing {session_id} pid={entry.proc.pid}", file=sys.stderr)

            return entry

    async def remove(self, session_id: str) -> None:
        """Remove and shut down an agent."""
        async with self._lock:
            entry = self._entries.pop(session_id, None)
        if entry:
            await entry.shutdown(reason="explicit_remove")

    async def flush_all(self, reason: str = "api_key_reset") -> int:
        """Shut down all running agents so they restart with fresh env vars."""
        async with self._lock:
            entries = list(self._entries.items())
            self._entries.clear()
        for sid, entry in entries:
            await entry.shutdown(reason=reason)
            print(f"[pool] flushed {sid} ({reason})", file=sys.stderr)
        return len(entries)

    async def _idle_sweeper(self) -> None:
        """Periodically check for idle agents and shut them down."""
        while True:
            await asyncio.sleep(30)
            now = _time.monotonic()
            to_remove: list[str] = []

            async with self._lock:
                for sid, entry in list(self._entries.items()):
                    if entry.is_dead:
                        to_remove.append(sid)
                        continue
                    if (
                        not entry.has_clients
                        and entry._last_ws_detach_at is not None
                        and (now - entry._last_ws_detach_at) > self._idle_timeout
                        and entry.agent_status == "idle"
                    ):
                        to_remove.append(sid)

                for sid in to_remove:
                    entry = self._entries.pop(sid, None)
                    if entry:
                        asyncio.create_task(entry.shutdown(reason="idle_timeout"))

    def status(self) -> list:
        """Return pool status for monitoring."""
        now = _time.monotonic()
        return [
            {
                "sessionId": sid,
                "cwd": e.cwd,
                "mode": e.mode,
                "pid": e.proc.pid if e.proc else None,
                "alive": e.proc is not None and e.proc.returncode is None,
                "agent_status": e.agent_status,
                "session_mode": e.session_mode,
                "model": e.model_info.get("displayName") or e.model_info.get("name"),
                "ws_attached": e._ws is not None,
                "subscribers": len(e._subscribers),
                "idle_seconds": round(now - e._last_ws_detach_at, 1) if e._last_ws_detach_at else None,
            }
            for sid, e in self._entries.items()
        ]


agent_pool = AgentPool(idle_timeout=float(os.environ.get("AGENT_IDLE_TIMEOUT", "300")))


# ══════════════════════════════════════════════════════════════════════════════
# § 18 · MISC / UTILITY ROUTES
# ══════════════════════════════════════════════════════════════════════════════


# ── Settings ──────────────────────────────────────────────────────────────────


# ── Database health ───────────────────────────────────────────────────────────


# ══════════════════════════════════════════════════════════════════════════════
# § 12 · DATABASE — WIDGET DATA
# ══════════════════════════════════════════════════════════════════════════════
# Each widget placement stores its records in a single JSON file.
# Complex queries are handled by executing a Python snippet server-side —
# this is intentional: widgets are vibe-coded through bzcode, not hand-injected.
#
# File location: server_data/widget_data/{canvasId}.json
# File format:   { "_next_id": N, "records": [{id, created_at, ...}, ...] }

import re as _re
import threading as _threading

_CANVAS_ID_RE = _re.compile(r"^[a-z0-9][a-z0-9-]{3,63}$")
WIDGET_DATA_DIR = SERVER_DATA_DIR / "widget_data"
_widget_locks: dict = {}  # per-canvasId write locks
_widget_locks_meta = _threading.Lock()


def _widget_lock(canvas_id: str) -> _threading.Lock:
    with _widget_locks_meta:
        if canvas_id not in _widget_locks:
            _widget_locks[canvas_id] = _threading.Lock()
        return _widget_locks[canvas_id]


def _widget_path(canvas_id: str, session_id: str = "") -> Path:
    if not _CANVAS_ID_RE.match(canvas_id):
        raise ValueError(f"Invalid canvasId: {canvas_id!r}")
    # Prefer session-scoped widget_data when a session owns this widget
    if session_id:
        session_dir = SESSIONS_DIR / session_id / "widget_data"
        p = session_dir / f"{canvas_id}.json"
        if p.exists() or not (WIDGET_DATA_DIR / f"{canvas_id}.json").exists():
            session_dir.mkdir(parents=True, exist_ok=True)
            return p
    WIDGET_DATA_DIR.mkdir(parents=True, exist_ok=True)
    return WIDGET_DATA_DIR / f"{canvas_id}.json"


def _widget_load(canvas_id: str, session_id: str = "") -> dict:
    p = _widget_path(canvas_id, session_id)
    if not p.exists():
        return {"_next_id": 1, "records": []}
    with open(p) as f:
        return json.load(f)


def _widget_save(canvas_id: str, data: dict, session_id: str = "") -> None:
    p = _widget_path(canvas_id, session_id)
    with open(p, "w") as f:
        json.dump(data, f, indent=2, default=str)


_WIDGET_CORS = {"Access-Control-Allow-Origin": "*"}


def _wj(data: object, status: int = 200) -> web.Response:
    return web.Response(
        text=json.dumps(data, default=str),
        status=status,
        content_type="application/json",
        headers=_WIDGET_CORS,
    )


# ── Agent mode + file endpoints ───────────────────────────────────────────────


# ══════════════════════════════════════════════════════════════════════════════
# § 7 · DOCUMENT PARSING
# ══════════════════════════════════════════════════════════════════════════════

_MAX_DOC_CHARS = 80_000
_MAX_DOC_BYTES = 50 * 1024 * 1024  # 50 MB

# ── DOCX ↔ Block JSON conversion (bz-office format) ─────────────────────────


def _docx_to_blocks(data: bytes) -> list:
    """Convert DOCX binary → Block[] in bz-office JSON format."""
    import base64
    import io
    import secrets

    import docx as _docx
    from docx.oxml.ns import qn

    doc = _docx.Document(io.BytesIO(data))
    blocks = []

    # Resolve major/minor theme fonts once for the whole document
    _major_font = _minor_font = None
    try:
        import xml.etree.ElementTree as _ET
        import zipfile as _zf

        with _zf.ZipFile(io.BytesIO(data)) as _z:
            _theme_names = [n for n in _z.namelist() if n.lower().endswith("theme1.xml") and "theme" in n.lower()]
            if _theme_names:
                _theme_xml = _z.read(_theme_names[0])
                _theme_root = _ET.fromstring(_theme_xml)
                _ans = "http://schemas.openxmlformats.org/drawingml/2006/main"
                _fs = _theme_root.find(f".//{{{_ans}}}fontScheme")
                if _fs is not None:
                    _mj = _fs.find(f"{{{_ans}}}majorFont/{{{_ans}}}latin")
                    _mn = _fs.find(f"{{{_ans}}}minorFont/{{{_ans}}}latin")
                    if _mj is not None:
                        _major_font = _mj.get("typeface")
                    if _mn is not None:
                        _minor_font = _mn.get("typeface")
    except Exception:
        pass

    # Resolve document-default font (Normal style or docDefaults)
    # Note: _resolve_font is defined below; inline the theme-ref logic here
    def _resolve_theme(name):
        if not name:
            return None
        if name in ("+mn-lt", "+Body"):
            return _minor_font
        if name in ("+mj-lt", "+Heading"):
            return _major_font
        if name.startswith("+"):
            return None
        return name

    _default_font = None
    try:
        v = doc.styles["Normal"].font.name
        _default_font = _resolve_theme(v) or v or None
    except Exception:
        pass

    # Resolve the document-default body font size (from Normal style or docDefaults).
    _default_font_size_pt = None
    try:
        sz = doc.styles["Normal"].font.size
        if sz:
            _default_font_size_pt = int(sz.pt)
    except Exception:
        pass
    if not _default_font_size_pt:
        try:
            _wns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
            # docDefaults lives in the styles XML part, NOT in doc.element (body)
            styles_el = doc.part.styles._element
            defs = styles_el.find(f"{{{_wns}}}docDefaults")
            if defs is not None:
                rPr = defs.find(f"{{{_wns}}}rPrDefault/{{{_wns}}}rPr")
                if rPr is not None:
                    sz_el = rPr.find(f"{{{_wns}}}sz")
                    if sz_el is not None:
                        val = sz_el.get(qn("w:val"))
                        if val:
                            _default_font_size_pt = int(val) // 2  # half-points
        except Exception:
            pass
    if not _default_font:
        try:
            _wns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
            defs = doc.element.find(f".//{{{_wns}}}docDefaults")
            if defs is not None:
                rFonts = defs.find(f".//{{{_wns}}}rFonts")
                if rFonts is not None:
                    raw = rFonts.get(qn("w:ascii")) or rFonts.get(qn("w:hAnsi"))
                    _default_font = _resolve_theme(raw)
        except Exception:
            pass

    def _resolve_font(name):
        """Resolve a font name, including theme references."""
        if not name:
            return None
        if name in ("+mj-lt", "+Heading"):
            return _major_font
        if name in ("+mn-lt", "+Body"):
            return _minor_font
        # w:asciiTheme / w:hAnsiTheme values like "majorHAnsi", "minorHAnsi"
        nl = name.lower()
        if nl.startswith("major"):
            return _major_font
        if nl.startswith("minor"):
            return _minor_font
        if name.startswith("+"):
            return None  # unknown theme slot
        return name

    def _rFonts_font(rFonts):
        """Extract resolved font from an rFonts element, checking all relevant attrs."""
        for attr in (qn("w:ascii"), qn("w:hAnsi"), qn("w:asciiTheme"), qn("w:hAnsiTheme"), qn("w:cs")):
            v = _resolve_font(rFonts.get(attr))
            if v:
                return v
        return None

    def _get_run_font(run, para=None):
        """Return effective font name for a run, tracing style inheritance."""
        # 1. Directly set on the run's rFonts XML
        try:
            rPr = run._r.rPr
            if rPr is not None:
                rFonts = rPr.find(qn("w:rFonts"))
                if rFonts is not None:
                    v = _rFonts_font(rFonts)
                    if v:
                        return v
        except Exception:
            pass
        # 2. Paragraph style's character rPr (for runs that inherit from the para style)
        if para is not None:
            try:
                _wns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
                style_el = para.style.element if para.style else None
                if style_el is not None:
                    style_rPr = style_el.find(f".//{{{_wns}}}rPr")
                    if style_rPr is not None:
                        rFonts = style_rPr.find(qn("w:rFonts"))
                        if rFonts is not None:
                            v = _rFonts_font(rFonts)
                            if v:
                                return v
            except Exception:
                pass
        # 3. python-docx high-level accessor
        try:
            v = _resolve_font(run.font.name)
            if v:
                return v
        except Exception:
            pass
        # 4. Document default font (minor/body theme font preferred for body text)
        return _minor_font or _default_font or None

    def _effective_font_size_pt(run, para):
        """Resolve the effective font size in points by walking the style chain."""
        # 1. Explicit size on the run
        if run.font.size:
            return int(run.font.size.pt)
        # 2. Character style on the run
        try:
            if run.style and run.style.font.size:
                return int(run.style.font.size.pt)
        except Exception:
            pass
        # 3. Paragraph style
        try:
            if para.style and para.style.font.size:
                return int(para.style.font.size.pt)
        except Exception:
            pass
        # 4. Document default (Normal style / docDefaults)
        return _default_font_size_pt

    def _run_styles(para) -> list:
        styles, pos = [], 0
        for run in para.runs:
            n = len(run.text)
            if not n:
                pos += n
                continue
            sr = {"start": pos, "end": pos + n}
            if run.bold:
                sr["isBold"] = True
            if run.italic:
                sr["isItalic"] = True
            if run.underline:
                sr["isUnderlined"] = True
            if getattr(run.font, "strike", None):
                sr["isStrikethrough"] = True
            eff_size = _effective_font_size_pt(run, para)
            if eff_size:
                sr["fontSize"] = eff_size
            if run.font.color and run.font.color.type is not None:
                try:
                    sr["textColor"] = f"#{run.font.color.rgb}"
                except Exception:
                    pass
            fname = _get_run_font(run, para)
            if fname:
                sr["fontFamily"] = fname
            if len(sr) > 2:
                styles.append(sr)
            pos += n
        return styles

    # Pre-read full heading properties (size, bold, italic, color) by walking the style chain.
    _H_WNS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    _FALLBACK_HEADING_SIZES = {"1": 24, "2": 20, "3": 18, "4": 16}

    def _resolve_style_size_pt(style_name: str):
        visited: set = set()
        try:
            s = doc.styles[style_name]
        except Exception:
            return None
        while s is not None and s.name not in visited:
            visited.add(s.name)
            if s.font.size:
                return int(s.font.size.pt)
            s = s.base_style
        return _default_font_size_pt

    _heading_props: dict = {}  # level str -> {fontSize, isBold, isItalic, textColor, fontFamily}
    for _level in ("1", "2", "3", "4"):
        _props: dict = {}
        _sz = _resolve_style_size_pt(f"Heading {_level}")
        if _sz:
            _props["fontSize"] = _sz
        try:
            _hs = doc.styles[f"Heading {_level}"]
            if _hs.font.bold:
                _props["isBold"] = True
            if _hs.font.italic:
                _props["isItalic"] = True
            _rPr = _hs._element.find(f"{{{_H_WNS}}}rPr")
            _col = _rPr.find(f"{{{_H_WNS}}}color") if _rPr is not None else None
            if _col is not None:
                _cv = _col.get(qn("w:val"))
                if _cv and _cv.lower() != "auto":
                    _props["textColor"] = f"#{_cv.upper()}"
        except Exception:
            pass
        _heading_props[_level] = _props

    def _heading_size(style_name: str):
        for level in ("1", "2", "3", "4"):
            if style_name == f"Heading {level}":
                return _heading_props[level].get("fontSize") or _FALLBACK_HEADING_SIZES[level]
        return None

    _ALIGN_MAP = {}
    try:
        from docx.enum.text import WD_ALIGN_PARAGRAPH as _WAP

        _ALIGN_MAP = {
            _WAP.CENTER: "center",
            _WAP.RIGHT: "right",
            _WAP.JUSTIFY: "justify",
        }
    except Exception:
        pass

    def _extract_drawing_style(drawing, para):
        """Extract image data and dimensions from a <w:drawing> element."""
        try:
            _wp = "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
            _a = "http://schemas.openxmlformats.org/drawingml/2006/main"
            _r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
            container = drawing.find(f"{{{_wp}}}inline") or drawing.find(f"{{{_wp}}}anchor")
            if container is None:
                return None
            extent = container.find(f"{{{_wp}}}extent")
            cx = int(extent.get("cx", 0)) if extent is not None else 0
            cy = int(extent.get("cy", 0)) if extent is not None else 0
            width_px = max(1, round(cx / 9525))
            height_px = max(1, round(cy / 9525))
            blip = container.find(f".//{{{_a}}}blip")
            if blip is None:
                return None
            r_id = blip.get(f"{{{_r}}}embed")
            if not r_id:
                return None
            image_part = para.part.related_parts[r_id]
            b64 = base64.b64encode(image_part.blob).decode("ascii")
            data_url = f"data:{image_part.content_type};base64,{b64}"
            return {"imageUrl": data_url, "imageWidth": width_px, "imageHeight": height_px}
        except Exception:
            return None

    def _para_to_block(para) -> dict:
        text = para.text
        styles = _run_styles(para)
        block = {"text": text, "styles": styles}

        # Scan ALL <w:r> elements in the paragraph XML, not just para.runs —
        # python-docx omits drawing-only runs (no <w:t>) from .runs, so image runs
        # would be silently skipped if we used that property.
        pos = 0
        insertions = []
        for run_el in para._p.findall(qn("w:r")):
            drawing = run_el.find(qn("w:drawing"))
            if drawing is not None:
                img = _extract_drawing_style(drawing, para)
                if img:
                    insertions.append((pos, img))
            t_el = run_el.find(qn("w:t"))
            pos += len(t_el.text if t_el is not None and t_el.text else "")

        if insertions:
            chars = list(block["text"])
            cur_styles = list(block["styles"])
            for ins_pos, img_style in sorted(insertions, key=lambda x: -x[0]):
                chars.insert(ins_pos, " ")
                shifted = []
                for sr in cur_styles:
                    nr = dict(sr)
                    if nr["start"] >= ins_pos:
                        nr["start"] += 1
                        nr["end"] += 1
                    elif nr["end"] > ins_pos:
                        nr["end"] += 1
                    shifted.append(nr)
                shifted.append({"start": ins_pos, "end": ins_pos + 1, **img_style})
                cur_styles = sorted(shifted, key=lambda s: s["start"])
            block["text"] = "".join(chars)
            block["styles"] = cur_styles

        # Heading → override styles with properties from the document's heading style
        sname = para.style.name if para.style else ""
        size = _heading_size(sname)
        if size:
            level_key = sname.split()[-1] if sname.startswith("Heading ") else None
            props = _heading_props.get(level_key, {}) if level_key else {}
            heading_font = _major_font or _minor_font or None
            heading_style: dict = {"start": 0, "end": len(block["text"]), "fontSize": size}
            if props.get("isBold", True):
                heading_style["isBold"] = True
            if props.get("isItalic"):
                heading_style["isItalic"] = True
            if props.get("textColor"):
                heading_style["textColor"] = props["textColor"]
            if heading_font:
                heading_style["fontFamily"] = heading_font
            block["styles"] = [heading_style]
            block["headingLevel"] = int(level_key) if level_key else None

        # Paragraph alignment
        if para.alignment in _ALIGN_MAP:
            block["alignment"] = _ALIGN_MAP[para.alignment]

        # Bullet / numbered list
        try:
            numPr = para._p.pPr.numPr if para._p.pPr is not None else None
            if numPr is not None:
                block["prefix"] = "•"
                ilvl = numPr.ilvl
                block["indent"] = int(ilvl.val) + 1 if ilvl is not None else 1
        except Exception:
            pass

        return block

    for child in doc.element.body.iterchildren():
        tag = child.tag.split("}")[-1]

        if tag == "p":
            try:
                from docx.text.paragraph import Paragraph

                para = Paragraph(child, doc)
                block = _para_to_block(para)
                if block["text"].strip() or not blocks:
                    blocks.append(block)
            except Exception:
                pass

        elif tag == "tbl":
            try:
                from docx.table import Table

                table = Table(child, doc)
                tid = secrets.token_hex(8)
                n_rows = len(table.rows)
                n_cols = max((len(r.cells) for r in table.rows), default=0)
                for r_idx, row in enumerate(table.rows):
                    for c_idx, cell in enumerate(row.cells):
                        # Build per-cell styles by walking each paragraph's runs
                        cell_text = ""
                        cell_styles = []
                        for p_idx, para in enumerate(cell.paragraphs):
                            if p_idx > 0:
                                cell_text += "\n"
                            offset = len(cell_text)
                            for sr in _run_styles(para):
                                cell_styles.append(
                                    {
                                        **sr,
                                        "start": sr["start"] + offset,
                                        "end": sr["end"] + offset,
                                    }
                                )
                            cell_text += para.text
                        blocks.append(
                            {
                                "text": cell_text,
                                "styles": cell_styles,
                                "isTableCell": True,
                                "tableId": tid,
                                "rowIndex": r_idx,
                                "columnIndex": c_idx,
                                "numberOfRows": n_rows,
                                "numberOfColumns": n_cols,
                            }
                        )
            except Exception:
                pass

    effective_default = _minor_font or _default_font or None
    return {"blocks": blocks, "defaultFont": effective_default}


def _blocks_to_docx(blocks: list) -> bytes:
    """Convert Block[] (bz-office format) → DOCX binary."""
    import base64
    import io

    import docx as _docx
    from docx.oxml.ns import qn
    from docx.shared import Emu, Pt, RGBColor
    from lxml import etree

    PX_TO_EMU = 9525  # 1 CSS px at 96 DPI
    WP = "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"

    def _make_float(run, wrap, placed_x_css):
        """Convert the inline picture in `run` to a floating wp:anchor element."""
        drawing = run._r.find(qn("w:drawing"))
        if drawing is None:
            return
        inline = drawing.find(f"{{{WP}}}inline")
        if inline is None:
            return

        extent = inline.find(f"{{{WP}}}extent")
        cx = extent.get("cx", "0") if extent is not None else "0"
        cy = extent.get("cy", "0") if extent is not None else "0"

        behind = "1" if wrap == "behind" else "0"
        # imagePlacedX is from the canvas/page left edge; DOCX column-relative
        # coordinates start at the text column left (~100 CSS px into the canvas).
        CANVAS_LEFT_CSS = 100  # matches bz-agent START_X / SF
        COLUMN_W_CSS = 621  # matches bz-agent (END_X - START_X) / SF
        col_x = max(0, (placed_x_css or 0) - CANVAS_LEFT_CSS)
        pos_x = int(col_x * PX_TO_EMU)
        # Float direction: image in left half → text wraps right; right half → text wraps left
        wrap_side = "right" if col_x < COLUMN_W_CSS / 2 else "left"
        wrap_el = "<wp:wrapBehindDoc/>" if wrap == "behind" else f'<wp:wrapSquare wrapText="{wrap_side}"/>'

        anchor = etree.fromstring(
            f'<wp:anchor xmlns:wp="{WP}"'
            f' distT="0" distB="114300" distL="114300" distR="114300"'
            f' simplePos="0" relativeHeight="251658240" behindDoc="{behind}"'
            f' locked="0" layoutInCell="1" allowOverlap="0">'
            f'<wp:simplePos x="0" y="0"/>'
            f'<wp:positionH relativeFrom="column"><wp:posOffset>{pos_x}</wp:posOffset></wp:positionH>'
            f'<wp:positionV relativeFrom="paragraph"><wp:posOffset>0</wp:posOffset></wp:positionV>'
            f'<wp:extent cx="{cx}" cy="{cy}"/>'
            f'<wp:effectExtent l="0" t="0" r="0" b="0"/>'
            f"{wrap_el}"
            f"</wp:anchor>"
        )
        # Move docPr, cNvGraphicFramePr, graphic from inline → anchor
        for tag in (f"{{{WP}}}docPr", f"{{{WP}}}cNvGraphicFramePr", f"{{{A}}}graphic"):
            child = inline.find(tag)
            if child is not None:
                anchor.append(child)

        drawing.remove(inline)
        drawing.append(anchor)

    doc = _docx.Document()

    # Group consecutive table cells by tableId
    i = 0
    while i < len(blocks):
        b = blocks[i]

        if b.get("isTableCell"):
            tid = b.get("tableId")
            cells = [c for c in blocks if c.get("tableId") == tid]
            n_rows = b.get("numberOfRows", 1)
            n_cols = b.get("numberOfColumns", 1)
            tbl = doc.add_table(rows=n_rows, cols=n_cols)
            tbl.style = "Table Grid"
            for cell in cells:
                r, c = cell.get("rowIndex", 0), cell.get("columnIndex", 0)
                try:
                    tbl.rows[r].cells[c].text = cell.get("text", "")
                    if r == 0:
                        for run in tbl.rows[r].cells[c].paragraphs[0].runs:
                            run.bold = True
                except Exception:
                    pass
            # Skip all cells belonging to this table
            while i < len(blocks) and blocks[i].get("tableId") == tid:
                i += 1
            continue

        # Regular paragraph
        text = b.get("text", "")
        styles = b.get("styles", [])
        prefix = b.get("prefix", "")
        indent = b.get("indent", 0)

        # Use stored headingLevel (set during parse) or fall back to none
        heading_size = b.get("headingLevel") or None

        if heading_size:
            para = doc.add_heading(text, level=heading_size)
            heading_font = next((sr.get("fontFamily") for sr in styles if sr.get("fontFamily")), None)
            if heading_font:
                for run in para.runs:
                    run.font.name = heading_font
        elif prefix == "•":
            para = doc.add_paragraph(style="List Bullet")
            para.add_run(text)
        else:
            para = doc.add_paragraph()
            if not styles:
                para.add_run(text)
            else:
                # Apply style ranges
                cursor = 0
                for sr in sorted(styles, key=lambda s: s.get("start", 0)):
                    s, e = sr.get("start", 0), sr.get("end", len(text))
                    if cursor < s:
                        para.add_run(text[cursor:s])
                    if sr.get("imageUrl"):
                        try:
                            data_url = sr["imageUrl"]
                            if data_url.startswith("data:"):
                                _, b64_data = data_url.split(",", 1)
                                img_bytes = base64.b64decode(b64_data)
                            else:
                                import urllib.request as _req

                                with _req.urlopen(data_url, timeout=10) as resp:
                                    img_bytes = resp.read()
                            wrap = sr.get("imageWrap", "inline")
                            img_run = para.add_run()
                            w_emu = sr.get("imageWidth", 64) * PX_TO_EMU
                            h_emu = sr.get("imageHeight", 64) * PX_TO_EMU
                            img_run.add_picture(io.BytesIO(img_bytes), width=Emu(w_emu), height=Emu(h_emu))
                            if wrap in ("square", "behind"):
                                _make_float(img_run, wrap, sr.get("imagePlacedX"))
                        except Exception:
                            para.add_run(text[s:e])
                    else:
                        run = para.add_run(text[s:e])
                        run.bold = sr.get("isBold", False)
                        run.italic = sr.get("isItalic", False)
                        run.underline = sr.get("isUnderlined", False)
                        run.font.strike = sr.get("isStrikethrough", False) or None
                        if sr.get("fontFamily"):
                            run.font.name = sr["fontFamily"]
                        if sr.get("fontSize"):
                            run.font.size = Pt(sr["fontSize"])
                        if sr.get("textColor"):
                            try:
                                hex_c = sr["textColor"].lstrip("#")
                                run.font.color.rgb = RGBColor(
                                    int(hex_c[0:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16)
                                )
                            except Exception:
                                pass
                    cursor = e
                if cursor < len(text):
                    para.add_run(text[cursor:])

        if indent and not heading_size:
            from docx.shared import Inches

            para.paragraph_format.left_indent = Inches(indent * 0.25)

        i += 1

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _parse_pdf(data: bytes):
    import io

    import pypdf

    reader = pypdf.PdfReader(io.BytesIO(data))
    pages = len(reader.pages)
    parts = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        if text.strip():
            parts.append(f"# Page {i}\n\n{text.strip()}")
    return pages, "\n\n".join(parts)


def _parse_docx(data: bytes):
    import io

    import docx

    doc = docx.Document(io.BytesIO(data))
    parts = []
    heading_map = {1: "#", 2: "##", 3: "###", 4: "####"}
    for para in doc.paragraphs:
        style = para.style.name if para.style else ""
        text = para.text.strip()
        if not text:
            continue
        level = next((int(s) for s in ("1", "2", "3", "4") if style == f"Heading {s}"), None)
        if level:
            parts.append(f"{heading_map[level]} {text}")
        else:
            parts.append(text)
    for table in doc.tables:
        rows = []
        for i, row in enumerate(table.rows):
            cells = [c.text.strip() for c in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
            if i == 0:
                rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
        parts.append("\n".join(rows))
    page_count = max(1, len(parts) // 10)  # approximate
    return page_count, "\n\n".join(parts)


def _parse_xlsx(data: bytes):
    import io

    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        # Drop fully-empty rows
        rows = [r for r in rows if any(c is not None for c in r)]
        if not rows:
            continue
        rows = rows[:1000]  # cap at 1000 rows
        parts.append(f"## Sheet: {sheet_name}")
        header = rows[0]
        parts.append("| " + " | ".join(str(c) if c is not None else "" for c in header) + " |")
        parts.append("| " + " | ".join(["---"] * len(header)) + " |")
        for row in rows[1:]:
            parts.append("| " + " | ".join(str(c) if c is not None else "" for c in row) + " |")
    return len(wb.sheetnames), "\n\n".join(parts)


def _parse_pptx(data: bytes):
    import io

    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        title_text = ""
        body_lines = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if shape.shape_type == 13:  # picture
                continue
            if not title_text and hasattr(shape, "placeholder_format") and shape.placeholder_format:
                title_text = text
            else:
                for para in shape.text_frame.paragraphs:
                    pt = para.text.strip()
                    if pt:
                        body_lines.append(f"- {pt}")
        header = f"## Slide {i}: {title_text}" if title_text else f"## Slide {i}"
        parts.append(header + ("\n" + "\n".join(body_lines) if body_lines else ""))
    return len(prs.slides), "\n\n".join(parts)


_DOCX_EXTS = {".docx", ".doc"}


def _detect_and_parse(filename: str, data: bytes) -> dict:
    ext = Path(filename).suffix.lower()
    fmt = ext.lstrip(".")
    if fmt in ("doc", "xls", "ppt"):
        fmt = {"doc": "docx", "xls": "xlsx", "ppt": "pptx"}[fmt]

    # DOCX/DOC → return Block[] (bz-office format); other formats → markdown text
    if ext in _DOCX_EXTS:
        result = _docx_to_blocks(data)
        blocks = result["blocks"]
        word_count = sum(len(b.get("text", "").split()) for b in blocks)
        return {
            "filename": filename,
            "type": fmt,
            "pages": max(1, len([b for b in blocks if not b.get("isTableCell")]) // 30),
            "wordCount": word_count,
            "truncated": False,
            "blocks": blocks,
            "defaultFont": result.get("defaultFont"),
        }

    parsers = {
        ".pdf": _parse_pdf,
        ".xlsx": _parse_xlsx,
        ".xls": _parse_xlsx,
        ".pptx": _parse_pptx,
        ".ppt": _parse_pptx,
    }
    if ext not in parsers:
        raise ValueError(f"unsupported format: {ext or '(no extension)'}")
    pages, content = parsers[ext](data)
    truncated = len(content) > _MAX_DOC_CHARS
    if truncated:
        content = content[:_MAX_DOC_CHARS]
    return {
        "filename": filename,
        "type": fmt,
        "pages": pages,
        "wordCount": len(content.split()),
        "truncated": truncated,
        "content": content,
    }


# ══════════════════════════════════════════════════════════════════════════════
# § 8 · EXCEL
# ══════════════════════════════════════════════════════════════════════════════


def _eval_excel_formula(formula: str, cells: dict) -> object:
    """Evaluate common Excel formulas against a cell dict {cell_id: {value: ...}}."""
    import re as _re

    if not formula.startswith("="):
        return None
    expr = formula[1:].strip()

    def cell_val(cid):
        cid = cid.upper()
        cd = cells.get(cid, {})
        v = cd.get("value")
        if v is None:
            return 0
        try:
            return float(v)
        except:
            return 0

    def expand_range(r):
        """Expand A1:A10 to list of cell ids."""
        m = _re.match(r"^([A-Z]+)(\d+):([A-Z]+)(\d+)$", r.upper())
        if not m:
            return [r.upper()]
        import openpyxl.utils as _ou

        c1 = _ou.column_index_from_string(m.group(1))
        r1 = int(m.group(2))
        c2 = _ou.column_index_from_string(m.group(3))
        r2 = int(m.group(4))
        return [f"{_ou.get_column_letter(c)}{r}" for r in range(r1, r2 + 1) for c in range(c1, c2 + 1)]

    try:
        # Handle SUM(range)
        m = _re.fullmatch(r"SUM\(([^)]+)\)", expr, _re.I)
        if m:
            vals = [cell_val(cid) for cid in expand_range(m.group(1).strip())]
            return sum(vals)

        # Handle AVERAGE(range)
        m = _re.fullmatch(r"AVERAGE\(([^)]+)\)", expr, _re.I)
        if m:
            vals = [cell_val(cid) for cid in expand_range(m.group(1).strip())]
            return sum(vals) / len(vals) if vals else 0

        # Handle COUNT(range)
        m = _re.fullmatch(r"COUNT\(([^)]+)\)", expr, _re.I)
        if m:
            vals = [
                1
                for cid in expand_range(m.group(1).strip())
                if cells.get(cid.upper(), {}).get("value") not in (None, "")
            ]
            return sum(vals)

        # Handle MIN/MAX(range)
        m = _re.fullmatch(r"(MIN|MAX)\(([^)]+)\)", expr, _re.I)
        if m:
            vals = [cell_val(cid) for cid in expand_range(m.group(2).strip())]
            return min(vals) if m.group(1).upper() == "MIN" else max(vals)

        # Replace cell references in arithmetic expression (e.g. A1+B2*C3)
        def repl_cell(m2):
            return str(cell_val(m2.group(0)))

        arith = _re.sub(r"[A-Z]+\d+", repl_cell, expr.upper())
        # Only evaluate if it's a simple arithmetic expression
        if _re.fullmatch(r"[\d\s\.\+\-\*\/\(\)]+", arith):
            result = eval(arith, {"__builtins__": {}})  # nosec — restricted input
            return result
    except Exception:
        pass
    return None


# ══════════════════════════════════════════════════════════════════════════════
# § 9 · POWERPOINT
# ══════════════════════════════════════════════════════════════════════════════


# ══════════════════════════════════════════════════════════════════════════════
# § 16 · DEV SERVER
# ══════════════════════════════════════════════════════════════════════════════
# Runs `pnpm dev` (or npm/yarn) in the project cwd and exposes the URL.
# Only one dev server per cwd is tracked; a second start call returns the URL.

_dev_servers: dict = {}  # cwd → { proc, url }


async def _find_free_port() -> int:
    import socket as _socket

    with _socket.socket() as s:
        s.bind(("", 0))
        return s.getsockname()[1]


# ══════════════════════════════════════════════════════════════════════════════
# § 17 · AGENT WEBSOCKET BRIDGE
# ══════════════════════════════════════════════════════════════════════════════


async def read_bzcode_stdout(
    proc: asyncio.subprocess.Process,
    out_queue: asyncio.Queue,
    ready_event: asyncio.Event,
    session_id: str = "",
    mode: str = "general",
) -> None:
    try:
        while True:
            try:
                line = await proc.stdout.readline()
            except ValueError:
                # Single output line exceeded the 16 MB StreamReader limit.
                # Drain until the next newline so the stream stays healthy.
                print("[ws] bzcode emitted an oversized output line (>16 MB) — skipping", file=sys.stderr)
                while True:
                    chunk = await proc.stdout.read(1 << 20)
                    if not chunk or b"\n" in chunk:
                        break
                continue
            if not line:
                break
            raw = line.decode().rstrip("\n")
            if not raw:
                continue
            await out_queue.put(raw)
            if raw[0] == "{":
                try:
                    msg = json.loads(raw)
                    mtype = msg.get("type")
                    mstatus = msg.get("status")
                    if mtype == "status" and mstatus == "running":
                        ready_event.clear()
                        if session_id:
                            _running_sessions.add(session_id)
                    elif mtype == "status" and mstatus == "idle":
                        ready_event.set()
                        if session_id:
                            _running_sessions.discard(session_id)
                    elif mtype == "result":
                        ready_event.set()
                        if session_id:
                            _running_sessions.discard(session_id)
                        if msg.get("usage"):
                            _add_tokens(msg["usage"])
                except Exception:
                    pass
    finally:
        await out_queue.put(None)
        ready_event.set()
        if session_id:
            _running_sessions.discard(session_id)


async def send_to_client(queue: asyncio.Queue, ws: "web.WebSocketResponse") -> None:
    while True:
        raw = await queue.get()
        if raw is None:
            break
        if raw and raw[0] == "{":
            try:
                await ws.send_str(_redact(raw))
            except Exception:
                # Socket already closing — stop sending
                break


async def drain_bzcode_stderr(
    proc: asyncio.subprocess.Process,
    out_queue: "asyncio.Queue | None" = None,
) -> None:
    """Read bzcode stderr, log it, and forward auth errors to the client."""
    _AUTH_KEYWORDS = (
        "token is expired",
        "Token refresh failed",
        "invalid authentication token",
        "invalid_token",
        "unauthorized",
        "401",
    )
    while True:
        line = await proc.stderr.readline()
        if not line:
            break
        text = line.decode().rstrip()
        print(f"[bzcode] {text}", file=sys.stderr)
        # Forward authentication errors so the frontend can prompt re-login
        if out_queue and any(k.lower() in text.lower() for k in _AUTH_KEYWORDS):
            msg = json.dumps(
                {
                    "type": "system",
                    "event": "auth-error",
                    "message": "Your authentication token has expired. Please sign in again.",
                }
            )
            try:
                out_queue.put_nowait(msg)
            except Exception:
                pass


async def relay_client_messages(
    proc: asyncio.subprocess.Process,
    ws: "web.WebSocketResponse",
    ready_event: asyncio.Event,
) -> None:
    async for msg in ws:
        if msg.type == aiohttp.WSMsgType.TEXT:
            raw: str = msg.data
            # Keepalive ping — reply with pong so the reverse proxy sees bidirectional activity.
            try:
                if json.loads(raw).get("type") == "ping":
                    await ws.send_str('{"type":"pong"}')
                    continue
            except Exception:
                pass
            if not raw.endswith("\n"):
                raw += "\n"
            proc.stdin.write(raw.encode())
            await proc.stdin.drain()
        elif msg.type in (aiohttp.WSMsgType.ERROR, aiohttp.WSMsgType.CLOSE):
            break


async def handle_ws_client(request: web.Request, bzcode_path: str, default_cwd: str) -> web.WebSocketResponse:
    """aiohttp WebSocket handler — mounted at GET /ws."""
    ws = web.WebSocketResponse()
    await ws.prepare(request)

    params = request.rel_url.query

    req_session_id = params.get("sessionId") or None
    req_cwd = params.get("cwd") or default_cwd

    # Validate cwd; fall back to default if the path doesn't exist
    effective_cwd = req_cwd if os.path.isdir(req_cwd) else default_cwd

    # Determine mode.
    req_mode = params.get("mode") or _load_mode_config().get("default", "general")

    # Generate a session ID if none was provided — bzcode accepts any ID.
    if not req_session_id:
        import secrets as _secrets

        req_session_id = f"bz-{_secrets.token_hex(6)}"
        print(f"[ws] generated new sessionId={req_session_id}", file=sys.stderr)

    # Write IDENTITY.md, SOUL.md, settings.json, meta.json into the session config dir.
    # Done before spawning so bzcode picks them up on first startup.
    _write_session_config(req_session_id, req_mode, working_dir=effective_cwd)

    # Always use --resume so bzcode loads the session config directory.
    # If no .jsonl exists yet, bzcode starts a fresh conversation with this ID.
    cmd = [bzcode_path, "--stdio", "--resume", req_session_id]

    print(f"[ws] connect  cwd={effective_cwd}  sessionId={req_session_id}  mode={req_mode}", file=sys.stderr)
    _active_sessions.add(req_session_id)

    try:
        proc = await asyncio.create_subprocess_exec(
            *cmd,
            stdin=asyncio.subprocess.PIPE,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            cwd=effective_cwd,
            env={**os.environ, **_read_api_keys(), "BZ_PYTHON": sys.executable},
            limit=16 * 1024 * 1024,  # 16 MB — large sessions can emit long lines
        )
    except FileNotFoundError:
        await ws.send_str(
            json.dumps(
                {
                    "type": "result",
                    "status": "error",
                    "error": f"bzcode not found: {bzcode_path}",
                }
            )
        )
        await ws.close()
        return ws

    out_queue = asyncio.Queue()
    ready_event = asyncio.Event()

    _entry_cfg = _mode_entry(req_mode)
    _session_mode = (_entry_cfg.get("settings") or {}).get("mode", "")
    if _session_mode:
        proc.stdin.write(json.dumps({"type": "setMode", "mode": _session_mode}).encode() + b"\n")
        await proc.stdin.drain()

    try:
        await asyncio.gather(
            read_bzcode_stdout(proc, out_queue, ready_event, session_id=req_session_id, mode=req_mode),
            send_to_client(out_queue, ws),
            drain_bzcode_stderr(proc, out_queue),
            relay_client_messages(proc, ws, ready_event),
        )
    except (BrokenPipeError, ConnectionResetError, asyncio.CancelledError):
        pass
    finally:
        _active_sessions.discard(req_session_id)
        # Check exit code — if bzcode crashed (non-zero), notify the client before
        # the handler returns so the frontend shows a meaningful error instead of a
        # silent disconnect.
        exit_code = proc.returncode
        if exit_code not in (None, 0):
            print(f"[ws] bzcode exited with code {exit_code}  pid={proc.pid}", file=sys.stderr)
            try:
                await ws.send_str(
                    json.dumps(
                        {
                            "type": "system",
                            "message": f"⚠ bzcode process exited unexpectedly (code {exit_code}). Reconnecting…",
                        }
                    )
                )
            except Exception:
                pass
        else:
            print(f"[ws] disconnect  pid={proc.pid}", file=sys.stderr)
        # Process may have already exited (e.g. crashed) — ignore lookup errors
        try:
            proc.terminate()
        except ProcessLookupError:
            pass
        try:
            await asyncio.wait_for(proc.wait(), timeout=5)
        except (ProcessLookupError, asyncio.TimeoutError):
            try:
                proc.kill()
            except ProcessLookupError:
                pass
    return ws


# ── Entry point ───────────────────────────────────────────────────────────────


if __name__ == "__main__":
    main()
