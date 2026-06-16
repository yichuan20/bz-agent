#!/usr/bin/env python3
"""
read-doc.py — parse PDF, DOCX, XLSX or PPTX and print extracted text.

Used by the worker agent via the Bash tool to read documents on disk.

Usage:
  python3 bzcode/scripts/read-doc.py --path /abs/path/to/report.pdf
  python3 bzcode/scripts/read-doc.py --path budget.xlsx --sheet "Q2 2026"
  python3 bzcode/scripts/read-doc.py --path big.pdf --max-chars 200000
"""
import argparse
import json
import sys
from pathlib import Path

_MAX_CHARS = 80_000
_MAX_BYTES = 50 * 1024 * 1024  # 50 MB


def parse_pdf(data: bytes) -> tuple[int, str]:
    import pypdf, io
    reader = pypdf.PdfReader(io.BytesIO(data))
    pages = len(reader.pages)
    parts = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        if text.strip():
            parts.append(f"# Page {i}\n\n{text.strip()}")
    return pages, "\n\n".join(parts)


def parse_docx(data: bytes) -> tuple[int, str]:
    import docx, io
    doc = docx.Document(io.BytesIO(data))
    heading_map = {1: "#", 2: "##", 3: "###", 4: "####"}
    parts = []
    for para in doc.paragraphs:
        style = para.style.name if para.style else ""
        text  = para.text.strip()
        if not text:
            continue
        level = next((int(s) for s in ("1", "2", "3", "4") if style == f"Heading {s}"), None)
        parts.append(f"{heading_map[level]} {text}" if level else text)
    for table in doc.tables:
        rows = []
        for i, row in enumerate(table.rows):
            cells = [c.text.strip() for c in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
            if i == 0:
                rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
        parts.append("\n".join(rows))
    return max(1, len(parts) // 10), "\n\n".join(parts)


def parse_xlsx(data: bytes, sheet_filter: str | None = None) -> tuple[int, str]:
    import openpyxl, io
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    names = [s for s in wb.sheetnames if sheet_filter is None or s == sheet_filter]
    parts = []
    for name in names:
        ws = wb[name]
        rows = [r for r in ws.iter_rows(values_only=True) if any(c is not None for c in r)][:1000]
        if not rows:
            continue
        parts.append(f"## Sheet: {name}")
        header = rows[0]
        parts.append("| " + " | ".join(str(c) if c is not None else "" for c in header) + " |")
        parts.append("| " + " | ".join(["---"] * len(header)) + " |")
        for row in rows[1:]:
            parts.append("| " + " | ".join(str(c) if c is not None else "" for c in row) + " |")
    return len(names), "\n\n".join(parts)


def parse_pptx(data: bytes) -> tuple[int, str]:
    from pptx import Presentation
    import io
    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        title_text, body_lines = "", []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if not title_text and hasattr(shape, "placeholder_format") and shape.placeholder_format:
                title_text = text
            else:
                for para in shape.text_frame.paragraphs:
                    pt = para.text.strip()
                    if pt:
                        body_lines.append(f"- {pt}")
        header = f"## Slide {i}: {title_text}" if title_text else f"## Slide {i}"
        parts.append(header + ("\n" + "\n".join(body_lines) if body_lines else ""))
    return len(prs.slides), "\n\n".join(parts)


PARSERS = {
    ".pdf":  lambda d, _: parse_pdf(d),
    ".docx": lambda d, _: parse_docx(d),
    ".doc":  lambda d, _: parse_docx(d),
    ".xlsx": lambda d, s: parse_xlsx(d, s),
    ".xls":  lambda d, s: parse_xlsx(d, s),
    ".pptx": lambda d, _: parse_pptx(d),
    ".ppt":  lambda d, _: parse_pptx(d),
}


def main() -> None:
    parser = argparse.ArgumentParser(description="Parse a document and print extracted text as JSON.")
    parser.add_argument("--path",      required=True, help="Absolute path to the document.")
    parser.add_argument("--sheet",     default=None,  help="Excel sheet name (XLSX only).")
    parser.add_argument("--max-chars", type=int, default=_MAX_CHARS, help="Character budget (default 80 000).")
    args = parser.parse_args()

    p = Path(args.path)
    if not p.exists():
        print(json.dumps({"error": f"file not found: {p}"}), flush=True)
        sys.exit(1)

    if p.stat().st_size > _MAX_BYTES:
        print(json.dumps({"error": "file too large (max 50 MB)"}), flush=True)
        sys.exit(1)

    ext = p.suffix.lower()
    if ext not in PARSERS:
        print(json.dumps({"error": f"unsupported format: {ext or '(no extension)'}"}), flush=True)
        sys.exit(1)

    try:
        data = p.read_bytes()
        pages, content = PARSERS[ext](data, args.sheet)
    except Exception as exc:
        print(json.dumps({"error": f"could not parse: {exc}"}), flush=True)
        sys.exit(1)

    truncated = len(content) > args.max_chars
    if truncated:
        content = content[:args.max_chars]

    result = {
        "filename":  p.name,
        "type":      ext.lstrip(".").replace("doc", "docx").replace("xls", "xlsx").replace("ppt", "pptx"),
        "pages":     pages,
        "wordCount": len(content.split()),
        "truncated": truncated,
        "content":   content,
    }
    print(json.dumps(result, ensure_ascii=False), flush=True)


if __name__ == "__main__":
    main()
