#!/usr/bin/env python3
from __future__ import annotations
"""
create-pptx.py — generate a PowerPoint presentation from a simple JSON outline.

Two modes:
  1. Template mode (--template-pptx):  inherits slide masters, fonts, colors and
     layouts from an existing .pptx file.  The agent should first call read-pptx.py
     to discover the available layout names, then reference those names in the
     outline JSON ("layout" field on each slide).

  2. Scratch mode (default):  builds slides using hardcoded light/dark themes.
     Use only when no template is available.

Usage:
  # Template mode — user's own PPTX as visual base
  python3 create-pptx.py --template-pptx company.pptx \\
      --outline-file /tmp/outline.json --out /tmp/new_deck.pptx

  # Scratch mode
  python3 create-pptx.py --outline-file /tmp/outline.json --out /tmp/deck.pptx

  # Inline JSON (both modes)
  python3 create-pptx.py --outline '{"title":"Q2","slides":[...]}' --out deck.pptx

Outline JSON schema:
  {
    "title":  "Deck Title",
    "author": "Name",            # optional
    "theme":  "dark"|"light",   # optional, scratch mode only; default "light"
    "slides": [
      {
        "layout": "Title Slide",     # optional; template mode: layout name from read-pptx.py
                                     # scratch mode: ignored
        "title":   "Slide Title",
        "bullets": ["Point 1", "Point 2"],  # optional
        "body":    "Free text",             # optional, used when no bullets
        "notes":   "Speaker notes"          # optional
      }
    ]
  }

Output:
  Prints JSON: { "ok": true, "path": "/abs/path/to/deck.pptx", "slides": N,
                 "mode": "template"|"scratch", "template": "/path/used" }
"""

import argparse
import json
import sys
from pathlib import Path

# ── Scratch-mode themes ───────────────────────────────────────────────────────

THEMES = {
    "light": {
        "bg":       (255, 255, 255),
        "title_fg": (28,  25,  23),
        "body_fg":  (107, 99,  88),
        "accent":   (20, 115, 223),
    },
    "dark": {
        "bg":       (10,  10,  10),
        "title_fg": (255, 255, 255),
        "body_fg":  (163, 163, 163),
        "accent":   (46, 136, 255),
    },
}

# ── Helpers ───────────────────────────────────────────────────────────────────

def _rgb(triple):
    from pptx.dml.color import RGBColor
    return RGBColor(*triple)


def _find_layout(prs, name: str | None):
    """Return the slide layout matching `name` (case-insensitive substring).
    Falls back to layout index 1 ('Title and Content') then index 0."""
    if name:
        lower = name.lower()
        for layout in prs.slide_layouts:
            if lower in layout.name.lower() or layout.name.lower() in lower:
                return layout
    # Prefer a 'Title and Content'-style layout as default
    for layout in prs.slide_layouts:
        if "content" in layout.name.lower():
            return layout
    return prs.slide_layouts[min(1, len(prs.slide_layouts) - 1)]


def _fill_placeholder(slide, idx: int, text: str) -> bool:
    """Write text into placeholder by idx. Returns True on success."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            try:
                ph.text = text
                return True
            except Exception:
                pass
    return False


def _fill_body_placeholder(slide, bullets: list[str], body_text: str) -> bool:
    """Write bullets or body into the first non-title placeholder. Returns True on success."""
    from pptx.util import Pt
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            continue  # skip title
        if not ph.has_text_frame:
            continue
        tf = ph.text_frame
        tf.clear()
        if bullets:
            for j, bullet in enumerate(bullets):
                para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                para.text = bullet
                para.level = 0
        elif body_text:
            tf.paragraphs[0].text = body_text
        return True
    return False


# ── Template mode ─────────────────────────────────────────────────────────────

def make_from_template(outline: dict, template_path: Path, out_path: Path) -> int:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation(str(template_path))
    slides_data = outline.get("slides", [])
    slide_count = 0

    for sdata in slides_data:
        layout = _find_layout(prs, sdata.get("layout"))
        slide  = prs.slides.add_slide(layout)
        slide_count += 1

        title_text = sdata.get("title", "")
        bullets    = sdata.get("bullets", [])
        body_text  = sdata.get("body", "")
        notes_text = sdata.get("notes", "")

        # Fill title placeholder (idx 0)
        if title_text:
            if not _fill_placeholder(slide, 0, title_text):
                # Fallback: write into the first CENTER_TITLE or TITLE shape
                for ph in slide.placeholders:
                    try:
                        ph.text = title_text
                        break
                    except Exception:
                        pass

        # Fill body placeholder
        if bullets or body_text:
            _fill_body_placeholder(slide, bullets, body_text)

        # Speaker notes
        if notes_text:
            try:
                slide.notes_slide.notes_text_frame.text = notes_text
            except Exception:
                pass

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return slide_count


# ── Scratch mode ──────────────────────────────────────────────────────────────

def make_from_scratch(outline: dict, out_path: Path) -> int:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    theme = THEMES.get(outline.get("theme", "light"), THEMES["light"])

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    slides_data = outline.get("slides", [])
    slide_count = 0

    for i, sdata in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        slide_count += 1

        # Background
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(*theme["bg"])

        # Accent bar at top
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*theme["accent"])
        bar.line.fill.background()

        # Title
        title_text = sdata.get("title", f"Slide {i + 1}")
        tx_title = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(1.2))
        tf = tx_title.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        run = p.runs[0]
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*theme["title_fg"])

        # Slide number
        num_box = slide.shapes.add_textbox(Inches(11.8), Inches(6.9), Inches(1.4), Inches(0.4))
        ntf = num_box.text_frame.paragraphs[0]
        ntf.text = str(i + 1)
        ntf.alignment = PP_ALIGN.RIGHT
        ntf.runs[0].font.size = Pt(11)
        ntf.runs[0].font.color.rgb = RGBColor(*theme["body_fg"])

        # Body
        bullets   = sdata.get("bullets", [])
        body_text = sdata.get("body", "")
        body_box  = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(12), Inches(5.0))
        btf = body_box.text_frame
        btf.word_wrap = True

        if bullets:
            for j, bullet in enumerate(bullets):
                para = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
                para.text = f"• {bullet}"
                para.space_before = Pt(4)
                run = para.runs[0]
                run.font.size = Pt(20)
                run.font.color.rgb = RGBColor(*theme["body_fg"])
        elif body_text:
            btf.paragraphs[0].text = body_text
            run = btf.paragraphs[0].runs[0]
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(*theme["body_fg"])

        # Speaker notes
        notes_text = sdata.get("notes", "")
        if notes_text:
            slide.notes_slide.notes_text_frame.text = notes_text

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return slide_count


# ── Entry point ───────────────────────────────────────────────────────────────

def main() -> None:
    parser = argparse.ArgumentParser(description="Generate a PPTX from a JSON outline.")
    group = parser.add_mutually_exclusive_group(required=True)
    group.add_argument("--outline",      help="Inline JSON outline string.")
    group.add_argument("--outline-file", help="Path to a JSON outline file.")
    parser.add_argument("--out",           required=True, help="Output .pptx path.")
    parser.add_argument("--template-pptx", default="",   help="Existing .pptx to use as visual template.")
    args = parser.parse_args()

    # Load outline
    if args.outline_file:
        f = Path(args.outline_file)
        if not f.exists():
            print(json.dumps({"error": f"outline file not found: {f}"}))
            sys.exit(1)
        try:
            outline = json.loads(f.read_text(encoding="utf-8"))
        except json.JSONDecodeError as e:
            print(json.dumps({"error": f"invalid JSON: {e}"}))
            sys.exit(1)
    else:
        try:
            outline = json.loads(args.outline)
        except json.JSONDecodeError as e:
            print(json.dumps({"error": f"invalid JSON: {e}"}))
            sys.exit(1)

    out_path = Path(args.out)
    if out_path.suffix.lower() != ".pptx":
        out_path = out_path.with_suffix(".pptx")

    # Resolve template
    template_path = Path(args.template_pptx) if args.template_pptx else None
    if template_path and not template_path.exists():
        print(json.dumps({"error": f"template not found: {template_path}"}))
        sys.exit(1)
    if template_path and template_path.suffix.lower() != ".pptx":
        print(json.dumps({"error": "template must be a .pptx file"}))
        sys.exit(1)

    try:
        if template_path:
            n = make_from_template(outline, template_path, out_path)
            mode = "template"
        else:
            n = make_from_scratch(outline, out_path)
            mode = "scratch"
    except Exception as exc:
        print(json.dumps({"error": f"could not create presentation: {exc}"}))
        sys.exit(1)

    result = {"ok": True, "path": str(out_path.resolve()), "slides": n, "mode": mode}
    if template_path:
        result["template"] = str(template_path.resolve())
    print(json.dumps(result))


if __name__ == "__main__":
    main()
