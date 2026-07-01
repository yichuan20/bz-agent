#!/usr/bin/env python3
from __future__ import annotations
"""
extract-doc.py — extract only the relevant sections from a document, saving tokens.

Instead of sending an entire 50-page PDF to the LLM, this script returns only
the paragraphs/sections that are semantically relevant to a query.

Usage:
  python3 bzcode/scripts/extract-doc.py --path report.pdf --query "revenue figures"
  python3 bzcode/scripts/extract-doc.py --path budget.xlsx --query "Q2 expenses" --sheet "Q2"
  python3 bzcode/scripts/extract-doc.py --path deck.pptx --query "market size"
  python3 bzcode/scripts/extract-doc.py --path report.pdf --query "risks" --top 5 --max-chars 8000

Output:
  Prints JSON:
  {
    "filename": "report.pdf",
    "query": "revenue figures",
    "matched_sections": N,
    "total_sections": M,
    "content": "<extracted text>",
    "truncated": false
  }
"""
import argparse
import json
import re
import sys
from pathlib import Path

_MAX_CHARS = 12_000
_MAX_BYTES = 50 * 1024 * 1024


# ── Text extraction (reuses read-doc logic) ────────────────────────────────────

def _extract_text(path: Path, sheet: str | None) -> tuple[str, int]:
    """Return (full_text, page_count)."""
    ext = path.suffix.lower()
    data = path.read_bytes()

    if ext == ".pdf":
        import pypdf, io
        reader = pypdf.PdfReader(io.BytesIO(data))
        parts = []
        for i, page in enumerate(reader.pages, 1):
            t = (page.extract_text() or "").strip()
            if t:
                parts.append(f"[Page {i}]\n{t}")
        return "\n\n".join(parts), len(reader.pages)

    if ext in (".docx", ".doc"):
        import docx, io
        doc = docx.Document(io.BytesIO(data))
        parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(parts), max(1, len(parts) // 8)

    if ext in (".xlsx", ".xls"):
        import openpyxl, io
        wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
        names = [s for s in wb.sheetnames if sheet is None or s == sheet]
        parts = []
        for name in names:
            ws = wb[name]
            rows = [r for r in ws.iter_rows(values_only=True) if any(c is not None for c in r)][:500]
            if rows:
                lines = ["  ".join(str(c) if c is not None else "" for c in row) for row in rows]
                parts.append(f"[Sheet: {name}]\n" + "\n".join(lines))
        return "\n\n".join(parts), len(names)

    if ext in (".pptx", ".ppt"):
        from pptx import Presentation
        import io
        prs = Presentation(io.BytesIO(data))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [s.text_frame.text.strip() for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
            if texts:
                parts.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(parts), len(prs.slides)

    raise ValueError(f"Unsupported format: {ext}")


# ── Scoring ────────────────────────────────────────────────────────────────────

def _score_section(section: str, query_terms: list[str]) -> float:
    """Simple TF-like score: count query term occurrences (case-insensitive)."""
    text_lower = section.lower()
    score = 0.0
    for term in query_terms:
        count = text_lower.count(term.lower())
        if count:
            score += count * (1 + len(term) / 10)  # longer terms weighted more
    # Boost short sections slightly (more focused)
    if len(section) < 400:
        score *= 1.2
    return score


def _tokenise_query(query: str) -> list[str]:
    """Split query into meaningful terms (drop stopwords)."""
    STOP = {"a","an","the","is","are","was","were","be","been","being",
            "in","on","at","to","for","of","and","or","but","with","by",
            "from","what","how","when","who","which","this","that"}
    words = re.findall(r"[a-z0-9']+", query.lower())
    # Include multi-word phrases (bigrams)
    terms = [w for w in words if w not in STOP and len(w) > 2]
    bigrams = [f"{words[i]} {words[i+1]}" for i in range(len(words) - 1)
               if words[i] not in STOP and words[i+1] not in STOP]
    return terms + bigrams


# ── Main ───────────────────────────────────────────────────────────────────────

def main() -> None:
    parser = argparse.ArgumentParser(description="Extract relevant sections from a document.")
    parser.add_argument("--path",      required=True, help="Absolute path to the document.")
    parser.add_argument("--query",     required=True, help="What to look for.")
    parser.add_argument("--sheet",     default=None,  help="XLSX sheet name filter.")
    parser.add_argument("--top",       type=int, default=8, help="Max sections to return (default 8).")
    parser.add_argument("--max-chars", type=int, default=_MAX_CHARS, help="Character budget.")
    args = parser.parse_args()

    p = Path(args.path)
    if not p.exists():
        print(json.dumps({"error": f"file not found: {p}"}), flush=True)
        sys.exit(1)
    if p.stat().st_size > _MAX_BYTES:
        print(json.dumps({"error": "file too large (max 50 MB)"}), flush=True)
        sys.exit(1)

    try:
        full_text, page_count = _extract_text(p, args.sheet)
    except Exception as exc:
        print(json.dumps({"error": f"could not parse: {exc}"}), flush=True)
        sys.exit(1)

    # Split into sections (paragraphs / page blocks / slide blocks)
    sections = [s.strip() for s in re.split(r'\n{2,}|\[(?:Page|Slide|Sheet)', full_text) if s.strip()]
    # Re-prepend markers that were consumed by the split
    restored = []
    for i, s in enumerate(sections):
        restored.append(s)
    sections = restored

    query_terms = _tokenise_query(args.query)

    # Score and rank
    scored = [(s, _score_section(s, query_terms)) for s in sections]
    scored.sort(key=lambda x: x[1], reverse=True)

    # Take top N with score > 0; fall back to first sections if nothing matches
    top = [s for s, score in scored[:args.top] if score > 0]
    if not top:
        top = sections[:min(args.top, len(sections))]

    combined = "\n\n---\n\n".join(top)
    truncated = len(combined) > args.max_chars
    if truncated:
        combined = combined[:args.max_chars]

    result = {
        "filename":        p.name,
        "query":           args.query,
        "matched_sections": len(top),
        "total_sections":  len(sections),
        "pages":           page_count,
        "content":         combined,
        "truncated":       truncated,
    }
    print(json.dumps(result, ensure_ascii=False), flush=True)


if __name__ == "__main__":
    main()
